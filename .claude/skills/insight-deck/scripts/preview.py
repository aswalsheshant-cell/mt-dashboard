#!/usr/bin/env python3
"""Render any .pptx to an HTML preview (and optionally PNGs) — no PowerPoint.

    python preview.py deck.pptx -o preview.html          # HTML, open in a browser
    python preview.py deck.pptx -o preview.html --png    # + one PNG per slide

Use it to eyeball a generated deck for overflow, overlap and empty tiles before
sending it. It reads the real .pptx, so it checks the file you will actually
send. Approximate by design: fonts and wrapping are the browser's, not
PowerPoint's — treat it as a layout check, not a pixel proof.
PNG export shells out to `npx playwright` and is skipped if that is missing.
"""

import argparse
import html
import os
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

PX = 96.0          # px per inch
EMU = 914400.0
ALIGN = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}
ANCHOR = {MSO_ANCHOR.MIDDLE: "center", MSO_ANCHOR.BOTTOM: "flex-end"}

PAGE_CSS = """
body{margin:0;background:#3a3f45;font-family:Aptos,'Segoe UI',Calibri,system-ui,
Arial,sans-serif}
.slide{position:relative;margin:18px auto;background:#fff;overflow:hidden;
box-shadow:0 6px 24px rgba(0,0,0,.35)}
.sh{position:absolute;box-sizing:border-box;display:flex;flex-direction:column}
.sh p{margin:0}
.n{position:absolute;left:0;top:-15px;color:#cfd4d9;font-size:11px}
"""


def color_of(fmt):
    try:
        if fmt.type is not None and fmt.fore_color.type is not None:
            return "#%s" % str(fmt.fore_color.rgb)
    except Exception:                                   # noqa: BLE001
        pass
    return None


def line_of(shape):
    try:
        if shape.line.fill.type is not None and shape.line.color.rgb is not None:
            w = shape.line.width.pt if shape.line.width else 0.75
            return f"{max(0.5, w)}px solid #{shape.line.color.rgb}"
    except Exception:                                   # noqa: BLE001
        pass
    return None


def para_html(p, default_size=11.0):
    align = ALIGN.get(p.alignment, "left")
    ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
    out = []
    for r in p.runs:
        st = [f"font-size:{(r.font.size.pt if r.font.size else default_size) * 1.333:.1f}px"]
        if r.font.bold:
            st.append("font-weight:700")
        try:
            if r.font.color and r.font.color.type is not None:
                st.append(f"color:#{r.font.color.rgb}")
        except Exception:                               # noqa: BLE001
            pass
        out.append(f"<span style=\"{';'.join(st)}\">{html.escape(r.text)}</span>")
    if not out:
        return ""
    return (f'<p style="text-align:{align};line-height:{ls * 1.15:.2f}">'
            f'{"".join(out)}</p>')


def shape_html(sh, media_dir):
    x, y = (sh.left or 0) / EMU * PX, (sh.top or 0) / EMU * PX
    w, h = (sh.width or 0) / EMU * PX, (sh.height or 0) / EMU * PX
    st = [f"left:{x:.1f}px", f"top:{y:.1f}px", f"width:{w:.1f}px",
          f"height:{h:.1f}px"]
    if sh.shape_type == 13 and media_dir:                # PICTURE
        try:
            img = sh.image
            path = os.path.join(media_dir, f"img{id(sh)}.{img.ext}")
            with open(path, "wb") as fh:
                fh.write(img.blob)
            st.append(f"background-image:url('{os.path.basename(path)}')")
            st += ["background-size:cover", "background-position:center"]
        except Exception:                               # noqa: BLE001
            st.append("background:#dcdcdc")
        return f'<div class="sh" style="{";".join(st)}"></div>'
    fill = color_of(sh.fill) if hasattr(sh, "fill") else None
    if fill:
        st.append(f"background:{fill}")
    ln = line_of(sh)
    if ln:
        st.append(f"border:{ln}")
    body = ""
    if sh.has_text_frame:
        tf = sh.text_frame
        st.append(f"justify-content:{ANCHOR.get(tf.vertical_anchor, 'flex-start')}")
        st.append(f"padding:{(tf.margin_top or 0) / EMU * PX:.0f}px "
                  f"{(tf.margin_right or 0) / EMU * PX:.0f}px "
                  f"{(tf.margin_bottom or 0) / EMU * PX:.0f}px "
                  f"{(tf.margin_left or 0) / EMU * PX:.0f}px")
        body = "".join(para_html(p) for p in tf.paragraphs)
        if body:
            st.append("overflow:hidden")
    if not fill and not ln and not body:
        return ""
    # rounded rectangles read better with a radius; harmless on plain rects
    st.append("border-radius:6px" if fill and min(w, h) > 12 else "border-radius:0")
    return f'<div class="sh" style="{";".join(st)}">{body}</div>'


def render(path, out):
    prs = Presentation(path)
    w = prs.slide_width / EMU * PX
    h = prs.slide_height / EMU * PX
    media = os.path.dirname(os.path.abspath(out)) or "."
    parts = [f"<style>{PAGE_CSS}.slide{{width:{w:.0f}px;height:{h:.0f}px}}</style>"]
    for i, s in enumerate(prs.slides, 1):
        shapes = "".join(shape_html(sh, media) for sh in s.shapes)
        parts.append(f'<div class="slide" id="s{i}"><div class="n">slide {i}'
                     f'</div>{shapes}</div>')
    with open(out, "w", encoding="utf-8") as fh:
        fh.write("\n".join(parts))
    return len(prs.slides), int(w), int(h)


PNG_JS = """
const {chromium} = require('playwright');
(async () => {
  const [file, n, w, h, prefix] = process.argv.slice(2);
  const b = await chromium.launch();
  const p = await b.newPage({viewport:{width:+w, height:+h}});
  await p.goto('file://' + file);
  for (let i = 1; i <= +n; i++) {
    await p.locator('#s' + i).screenshot({path: `${prefix}-${i}.png`});
  }
  await b.close();
})();
"""


def to_png(html_path, n, w, h, prefix):
    node = shutil.which("npx") or "/opt/node22/bin/npx"
    if not os.path.exists(node) and not shutil.which("npx"):
        print("npx not found — skipping PNG export", file=sys.stderr)
        return
    with tempfile.NamedTemporaryFile("w", suffix=".js", delete=False) as fh:
        fh.write(PNG_JS)
        js = fh.name
    try:
        subprocess.run([node, "playwright", "screenshot", "--help"],
                       capture_output=True, timeout=120)
        subprocess.run(["node", js, os.path.abspath(html_path), str(n),
                        str(w), str(h), prefix], check=True,
                       env={**os.environ,
                            "NODE_PATH": os.environ.get("NODE_PATH", "")
                            or "/opt/node22/lib/node_modules"})
        print(f"wrote {prefix}-1..{n}.png")
    except Exception as exc:                            # noqa: BLE001
        print(f"PNG export failed ({exc}) — open the HTML instead",
              file=sys.stderr)
    finally:
        os.unlink(js)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="preview.html")
    ap.add_argument("--png", action="store_true", help="also export PNGs")
    a = ap.parse_args()
    n, w, h = render(a.pptx, a.out)
    print(f"wrote {a.out} — {n} slide(s)")
    if a.png:
        to_png(a.out, n, w, h, os.path.splitext(a.out)[0])


if __name__ == "__main__":
    main()
