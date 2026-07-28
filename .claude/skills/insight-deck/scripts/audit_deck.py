#!/usr/bin/env python3
"""Audit an existing .pptx and print a page-compression plan.

    python audit_deck.py deck.pptx [--json plan.json]

Finds the three things that inflate a deck to 20-25 pages:
  1. CLONES     - slides with the same structure and title pattern, one per
                  zone / chain / month. Collapse into one small-multiple page.
  2. PHOTO PAGES- slides that are mostly pictures. Collapse into a contact
                  sheet, or move to appendix.
  3. THIN PAGES - slides carrying very little information for a whole page.
                  Merge into a neighbour.
It reports, it never edits the source deck.
"""

import argparse
import json
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Words that vary between otherwise identical slides (the small-multiple axis).
VARY = re.compile(
    r"\b(east|west|north|south[- ]?\d?|central|zone \w+|q[1-4]|fy\d{2}|"
    r"jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\b['’\-\s]*\d*",
    re.I)


EMU_IN = 914400


def slide_title(slide):
    """Biggest text near the top wins — that is what reads as the title."""
    best, best_score = "", -1e9
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt or len(txt.split()) > 22:
            continue
        sz = 0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sz = max(sz, r.font.size.pt)
        score = sz - 8.0 * (sh.top or 0) / EMU_IN
        if score > best_score:
            best_score, best = score, txt.split("\n")[0]
    return best[:110]


def profile(slide):
    pics = charts = tables = texts = 0
    words = 0
    for sh in slide.shapes:
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        elif getattr(sh, "has_chart", False) and sh.has_chart:
            charts += 1
        elif getattr(sh, "has_table", False) and sh.has_table:
            tables += 1
        if sh.has_text_frame and sh.text_frame.text.strip():
            texts += 1
            words += len(sh.text_frame.text.split())
    return {"pics": pics, "charts": charts, "tables": tables,
            "texts": texts, "words": words, "shapes": len(slide.shapes)}


STOP = {"the", "a", "of", "and", "for", "in", "by", "vs", "wise", "view",
        "summary", "snapshot", "performance", "overview"}


def stem(title):
    """Title reduced to its topic tokens, with the varying axis removed."""
    t = VARY.sub(" ", title.lower())
    toks = [w for w in re.sub(r"[^a-z ]+", " ", t).split()
            if len(w) > 1 and w not in STOP]
    return set(toks)


def jaccard(a, b):
    return len(a & b) / len(a | b) if a and b else 0.0


def shape_sig(s):
    return (s["pics"], s["charts"], s["tables"], s["shapes"] // 6)


def same(a, b, thresh):
    """Same page repeated? Similar title, or identical build when titles are weak."""
    if jaccard(a["stem"], b["stem"]) >= thresh:
        return True
    weak = len(a["stem"]) < 2 or len(b["stem"]) < 2
    return weak and shape_sig(a) == shape_sig(b) and a["shapes"] > 6


def cluster(slides, thresh=0.5):
    """Greedy clustering — same topic, different filter value."""
    out = []
    for s in slides:
        for grp in out:
            if any(same(s, m, thresh) for m in grp["members"]):
                grp["members"].append(s)
                break
        else:
            out.append({"members": [s]})
    return [g["members"] for g in out if len(g["members"]) > 1]


def audit(path):
    prs = Presentation(path)
    slides = []
    for i, s in enumerate(prs.slides, 1):
        t = slide_title(s)
        slides.append({"page": i, "title": t, **profile(s), "stem": stem(t)})

    clones = cluster(slides)
    cloned = {s["page"] for g in clones for s in g}
    photos = [s for s in slides if s["page"] not in cloned and s["pics"] >= 3
              and s["words"] / max(1, s["pics"]) < 40]
    thin = [s for s in slides if s["page"] not in cloned and s["words"] < 25
            and s["pics"] < 3 and s["charts"] + s["tables"] <= 1]

    accounted, plan = set(), []
    for g in sorted(clones, key=lambda g: -len(g)):
        pages = [s["page"] for s in g]
        accounted.update(pages)
        plan.append({
            "action": "COLLAPSE", "pages": pages, "to": 1,
            "why": f"{len(g)} structural clones of '{g[0]['title'][:60]}'",
            "how": "one page, `metrics` chip grid (one chip per variant) + a "
                   "`bars` tile ranking them + 2-3 exception bullets",
        })
    ph = [s["page"] for s in photos if s["page"] not in accounted]
    if len(ph) > 1:
        accounted.update(ph)
        plan.append({
            "action": "CONTACT-SHEET", "pages": ph,
            "to": max(1, -(-len(ph) // 8)),
            "why": f"{len(ph)} photo-led pages",
            "how": "one `image` tile, 4 cols x 2 rows with captions; the rest "
                   "to appendix",
        })
    th = [s["page"] for s in thin if s["page"] not in accounted]
    if th:
        plan.append({
            "action": "MERGE", "pages": th, "to": max(1, len(th) // 3),
            "why": "thin pages (<25 words, <=1 visual)",
            "how": "fold into the nearest themed page as a tile",
        })
    keep = [s["page"] for s in slides
            if s["page"] not in accounted and s["page"] not in th]
    target = sum(p["to"] for p in plan) + len(keep)
    return {"file": path, "slides": len(slides), "target": target,
            "plan": plan, "keep": keep, "detail": [
                {k: v for k, v in s.items() if k != "stem"} for s in slides]}


def report(a):
    print(f"\n{a['file']}\n{'=' * len(a['file'])}")
    print(f"{a['slides']} slides  ->  ~{a['target']} after compression "
          f"({100 - round(100 * a['target'] / max(1, a['slides']))}% fewer pages)\n")
    print(f"{'PAGE':>4}  {'WORDS':>5} {'PIC':>3} {'CHT':>3} {'TBL':>3}  TITLE")
    for s in a["detail"]:
        print(f"{s['page']:>4}  {s['words']:>5} {s['pics']:>3} {s['charts']:>3} "
              f"{s['tables']:>3}  {s['title'][:70]}")
    print("\nCOMPRESSION PLAN")
    if not a["plan"]:
        print("  no structural redundancy found — compress by content, not by "
              "structure (see SKILL.md step 2)")
    for p in a["plan"]:
        pages = ",".join(map(str, p["pages"]))
        print(f"  [{p['action']}] p{pages}  ->  {p['to']} page(s)")
        print(f"      why : {p['why']}")
        print(f"      how : {p['how']}")
    if a["keep"]:
        print(f"  [KEEP]   p{','.join(map(str, a['keep']))} — distinct content, "
              "still check each against the one-insight-per-page rule")


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx")
    ap.add_argument("--json", help="also write the plan as JSON")
    a = ap.parse_args()
    res = audit(a.pptx)
    report(res)
    if a.json:
        with open(a.json, "w", encoding="utf-8") as fh:
            json.dump(res, fh, indent=2)
        print(f"\nwrote {a.json}")


if __name__ == "__main__":
    main()
