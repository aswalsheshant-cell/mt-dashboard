#!/usr/bin/env python3
"""Render a dense one-pager insight deck from a JSON spec.

    python build_deck.py spec.json -o deck.pptx

One spec "slide" = one high-density page: header band, KPI strip, a grid of
tiles, and a takeaway bar. The point is to fit what usually sprawls over 5-6
slides onto a single page that still reads cleanly.

See reference/spec_schema.md for the full spec. Nothing here is MT-specific
except the default palette, which comes from dashboard/index.html.
"""

import argparse
import json
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- theme ----

THEMES = {
    # Honasa / MT dashboard palette (dashboard/index.html :13-16)
    "honasa": {
        "accent": "2D9B7F", "accent_d": "1F7A63", "bg": "F4F1EC",
        "card": "FFFFFF", "ink": "1F2933", "muted": "6B7682",
        "line": "E4E0D8", "good": "1E8E3E", "risk": "C0392B",
        "warn": "C77D17", "info": "2B6CB0", "band_text": "FFFFFF",
        "font": "Aptos", "font_head": "Aptos",
    },
    "slate": {
        "accent": "2B6CB0", "accent_d": "1A3F6B", "bg": "F2F4F7",
        "card": "FFFFFF", "ink": "16202B", "muted": "667085",
        "line": "E2E6EC", "good": "12805C", "risk": "B42318",
        "warn": "B54708", "info": "175CD3", "band_text": "FFFFFF",
        "font": "Aptos", "font_head": "Aptos",
    },
}

TONE_KEY = {"good": "good", "risk": "risk", "warn": "warn",
            "info": "info", "neutral": "accent", "accent": "accent"}

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.30
GUTTER = 0.14


def C(theme, key):
    """Resolve a palette key (or a literal RRGGBB) to an RGBColor."""
    val = theme.get(key, key)
    return RGBColor.from_string(str(val).lstrip("#").upper())


def tone_color(theme, tone):
    return C(theme, TONE_KEY.get(tone or "neutral", "accent"))


# --------------------------------------------------------------- shapes ----

def box(slide, x, y, w, h, fill=None, line=None, lw=0.75, radius=None,
        shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        # adjustment is a fraction of the shape's shorter side
        try:
            shp.adjustments[0] = max(0.0, min(0.5, radius / max(0.01, min(w, h))))
        except (IndexError, ValueError):
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    if not shadow:
        shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    return shp


def _runs(para, s, theme, size, color, bold, font):
    """Split **bold** spans into separate runs."""
    for i, chunk in enumerate(re.split(r"\*\*(.+?)\*\*", str(s))):
        if not chunk:
            continue
        r = para.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.bold = bold or bool(i % 2)
        r.font.color.rgb = color
        r.font.name = font


def text(shape, lines, theme, size=11, color=None, bold=False, align="l",
         anchor="t", space=2, line_spacing=0.95, font=None, inset=0.07):
    """Write `lines` (str or list) into a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(inset)
    tf.margin_top = tf.margin_bottom = Inches(max(0.02, inset * 0.5))
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    if isinstance(lines, str):
        lines = [lines]
    color = color if color is not None else C(theme, "ink")
    font = font or theme["font"]
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                          "r": PP_ALIGN.RIGHT}[align]
        para.line_spacing = line_spacing
        para.space_after = Pt(space if i < len(lines) - 1 else 0)
        _runs(para, ln, theme, size, color, bold, font)
    return tf


def label(slide, x, y, w, h, s, theme, **kw):
    shp = box(slide, x, y, w, h)
    text(shp, s, theme, **kw)
    return shp


def _plain(s):
    return re.sub(r"\*\*", "", str(s))


def n_lines(s, w_in, size_pt, cw=0.53):
    """Estimated wrapped line count for `s` in a box `w_in` inches wide."""
    per_line = max(1, int(w_in * 96 / (size_pt * 1.333 * cw)))
    return max(1, -(-len(_plain(s)) // per_line))


def text_h(lines, w_in, size_pt, lead=1.06):
    """Estimated height in inches for a list of strings at `size_pt`."""
    if isinstance(lines, str):
        lines = [lines]
    n = sum(n_lines(x, w_in, size_pt) for x in lines)
    return n * size_pt * 1.333 * lead / 96


def fit_size(lines, w_in, h_in, base, floor=7.5, step=0.5):
    """Largest size <= base at which `lines` still fit in the box."""
    size = base
    while size > floor and text_h(lines, w_in, size) * 1.06 > h_in:
        size -= step
    return max(floor, size)


def delta_color(theme, s, direction=None):
    d = (direction or "").lower()
    s = str(s or "")
    if d in ("up", "good") or s.startswith("▲") or (d == "" and s.startswith("+")):
        return C(theme, "good")
    if d in ("down", "bad") or s.startswith("▼") or (d == "" and s.startswith("-")):
        return C(theme, "risk")
    return C(theme, "muted")


# ---------------------------------------------------------------- blocks ----

def draw_header(slide, sp, theme, y):
    """Dark band: eyebrow + headline (the conclusion) + right-hand stamp.

    The band grows to fit the headline, so a two-line conclusion never clips.
    """
    stamp = sp.get("stamp")
    stamp_w = max(1.5, 0.105 * len(str(stamp)) + 0.5) if stamp else 0.0
    tw = SLIDE_W - 2 * MARGIN - (stamp_w + 0.25 if stamp else 0)

    hl = sp.get("headline", "")
    size = 21.0
    while size > 14 and n_lines(hl, tw, size, cw=0.53) > 2:
        size -= 0.5
    hl_h = text_h(hl, tw, size, lead=1.0) * (1.0 if size > 14 else 1.15)

    top_pad, eyebrow_h = 0.11, 0.21 if sp.get("eyebrow") else 0.0
    sub_h = 0.0
    if sp.get("subhead"):
        sub_h = text_h(sp["subhead"], tw, 10) + 0.04
    h = sp.get("header_h") or max(0.86, top_pad + eyebrow_h + hl_h + sub_h + 0.12)

    box(slide, 0, y, SLIDE_W, h, fill=C(theme, "accent_d"))
    box(slide, 0, y + h, SLIDE_W, 0.045, fill=C(theme, "accent"))
    if stamp:
        s = box(slide, SLIDE_W - MARGIN - stamp_w, y + (h - 0.42) / 2,
                stamp_w, 0.42, fill=C(theme, "accent"), radius=0.09)
        text(s, stamp, theme, size=11.5, bold=True,
             color=C(theme, "band_text"), align="c", anchor="m")

    ty = y + top_pad
    if eyebrow_h:
        label(slide, MARGIN, ty, tw, eyebrow_h, sp["eyebrow"].upper(), theme,
              size=9.5, bold=True, color=C(theme, "accent"), inset=0)
        ty += eyebrow_h
    label(slide, MARGIN, ty, tw, hl_h, hl, theme, size=size, bold=True,
          color=C(theme, "band_text"), inset=0, line_spacing=0.95,
          font=theme["font_head"])
    if sub_h:
        label(slide, MARGIN, ty + hl_h + 0.02, tw, sub_h, sp["subhead"], theme,
              size=10, color=RGBColor.from_string("CFE3DB"), inset=0)
    return y + h + 0.045 + 0.10


def draw_kpis(slide, kpis, theme, y):
    """Evenly spread KPI cards with a coloured left rail."""
    h = 0.80
    n = len(kpis)
    w = (SLIDE_W - 2 * MARGIN - GUTTER * (n - 1)) / n
    for i, k in enumerate(kpis):
        x = MARGIN + i * (w + GUTTER)
        box(slide, x, y, w, h, fill=C(theme, "card"), line=C(theme, "line"),
            radius=0.07)
        box(slide, x, y + 0.08, 0.055, h - 0.16,
            fill=tone_color(theme, k.get("tone")))
        tx, tw = x + 0.20, w - 0.30
        lab = str(k.get("label", "")).upper()
        label(slide, tx, y + 0.06, tw, 0.20, lab, theme,
              size=fit_size(lab, tw, 0.20, 8.5, floor=6.5), bold=True,
              color=C(theme, "muted"), inset=0)
        val = str(k.get("value", ""))
        label(slide, tx, y + 0.25, tw, 0.32, val, theme,
              size=fit_size(val, tw, 0.30, 19, floor=11), bold=True, inset=0,
              font=theme["font_head"])
        if k.get("delta"):
            d = str(k["delta"])
            label(slide, tx, y + 0.57, tw, 0.20, d, theme,
                  size=fit_size(d, tw, 0.19, 9.5, floor=7), bold=True,
                  color=delta_color(theme, d, k.get("dir")), inset=0)
    return y + h + GUTTER


def draw_footer(slide, items, theme, y, h=0.58, title=None):
    """Takeaway / action strip: the 'so what' in one line each."""
    box(slide, 0, y, SLIDE_W, h, fill=C(theme, "ink"))
    x = MARGIN
    if title:
        tw = 0.085 * len(title) + 0.34
        b = box(slide, x, y + (h - 0.34) / 2, tw, 0.34, fill=C(theme, "accent"),
                radius=0.07)
        text(b, title.upper(), theme, size=9.5, bold=True,
             color=C(theme, "band_text"), align="c", anchor="m")
        x += tw + 0.18
    n = max(1, len(items))
    w = (SLIDE_W - x - MARGIN - 0.55 - GUTTER * (n - 1)) / n   # 0.55 = page no.
    flat = [(i.get("lead", "") + "  " + i.get("text", "")) if isinstance(i, dict)
            else str(i) for i in items]
    fsize = min(9.5, fit_size(max(flat, key=len), w - 0.06, h - 0.16, 9.5,
                              floor=7))
    for i, it in enumerate(items):
        cx = x + i * (w + GUTTER)
        if i:
            box(slide, cx - GUTTER / 2, y + 0.13, 0.02, h - 0.26,
                fill=RGBColor.from_string("4A5560"))
        if isinstance(it, dict):
            lead, body = it.get("lead"), it.get("text", "")
        else:
            lead, body = None, it
        shp = box(slide, cx, y + 0.04, w, h - 0.08)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.03)
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 0.92
        if lead:
            r = p.add_run()
            r.text = lead.upper() + "  "
            r.font.size, r.font.bold = Pt(fsize - 0.5), True
            r.font.color.rgb = C(theme, "accent")
            r.font.name = theme["font"]
        _runs(p, body, theme, fsize, RGBColor.from_string("EDEFF1"), False,
              theme["font"])
    return y + h


# ----------------------------------------------------------------- tiles ----

def tile_frame(slide, x, y, w, h, t, theme):
    """Card + coloured title bar. Returns the y where body content starts."""
    tone = tone_color(theme, t.get("tone"))
    box(slide, x, y, w, h, fill=C(theme, "card"), line=C(theme, "line"),
        radius=0.08)
    if t.get("title"):
        bh = 0.28
        box(slide, x, y, w, bh, fill=tone, radius=0.08)
        box(slide, x, y + bh - 0.09, w, 0.09, fill=tone)
        title = str(t["title"]).upper()
        tsize = 10 if len(title) <= int(w * 20) else 8.8
        label(slide, x + 0.10, y, w - 0.2, bh, title, theme, size=tsize,
              bold=True, color=C(theme, "band_text"), anchor="m", inset=0)
        return y + bh + 0.07
    return y + 0.09


def t_bullets(slide, x, y, w, h, t, theme):
    by = tile_frame(slide, x, y, w, h, t, theme)
    items = t.get("items", t.get("bullets", []))
    avail = (y + h) - by - 0.06
    n = max(1, len(items))
    flat = [(it.get("lead", "") + "  " + it.get("text", "")) if isinstance(it, dict)
            else str(it) for it in items]
    size = t.get("size") or fit_size(flat, w - 0.34, avail - 0.04, 10, floor=7)
    rowh = avail / n
    for i, it in enumerate(items):
        ry = by + i * rowh
        if isinstance(it, dict):
            lead, body, tone = it.get("lead"), it.get("text", ""), it.get("tone")
        else:
            lead, body, tone = None, it, None
        # marker sits on the FIRST line of the bullet, not the middle of it
        body_w = w - 0.36
        lines = n_lines((lead + "  " if lead else "") + str(body), body_w, size)
        lh = size * 1.333 * 0.92 / 96
        my = ry + rowh / 2 - (lines - 1) * lh / 2 - 0.035
        box(slide, x + 0.11, max(ry + 0.02, my), 0.07, 0.07,
            fill=tone_color(theme, tone or t.get("tone")))
        shp = box(slide, x + 0.26, ry, body_w, rowh)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 0.92
        if lead:
            r = p.add_run()
            r.text = lead + "  "
            r.font.size, r.font.bold = Pt(size), True
            r.font.color.rgb = tone_color(theme, tone or t.get("tone"))
            r.font.name = theme["font"]
        _runs(p, body, theme, size, C(theme, "ink"), False, theme["font"])


def t_bars(slide, x, y, w, h, t, theme):
    """Horizontal bars — label, bar, value, optional delta badge."""
    by = tile_frame(slide, x, y, w, h, t, theme)
    rows = t.get("items", [])
    if not rows:
        return
    vals = [float(r.get("value", 0) or 0) for r in rows]
    top = max(vals + [0.0001])
    avail = (y + h) - by - 0.06
    rowh = min(0.42, avail / len(rows))
    by += max(0.0, (avail - rowh * len(rows)) / 2)      # centre short blocks
    size = 9 if rowh > 0.3 else 8.2
    lab_w = t.get("label_w", max(0.62, min(1.5, w * 0.26)))
    note_w = 0.72 if any(r.get("note") for r in rows) else 0.0
    val_w = 0.62
    bar_x = x + 0.12 + lab_w + 0.06
    bar_w = w - 0.24 - lab_w - 0.06 - val_w - note_w
    for i, r in enumerate(rows):
        ry = by + i * rowh
        label(slide, x + 0.12, ry, lab_w, rowh, str(r.get("label", "")), theme,
              size=size, bold=True, anchor="m", inset=0)
        bh = min(0.19, rowh * 0.55)
        cy = ry + (rowh - bh) / 2
        box(slide, bar_x, cy, bar_w, bh, fill=C(theme, "bg"), radius=0.04)
        frac = max(0.02, float(r.get("value", 0) or 0) / top)
        box(slide, bar_x, cy, max(0.06, bar_w * frac), bh,
            fill=tone_color(theme, r.get("tone") or t.get("tone")), radius=0.04)
        label(slide, bar_x + bar_w + 0.03, ry, val_w, rowh,
              str(r.get("display", r.get("value", ""))), theme, size=size,
              bold=True, anchor="m", align="r", inset=0)
        if r.get("note"):
            label(slide, bar_x + bar_w + 0.03 + val_w, ry, note_w, rowh,
                  str(r["note"]), theme, size=size - 0.5, bold=True, anchor="m",
                  align="r", color=delta_color(theme, r["note"], r.get("dir")),
                  inset=0)


def t_metrics(slide, x, y, w, h, t, theme):
    """Chip grid — the small-multiple killer: one chip per zone/chain/SKU."""
    by = tile_frame(slide, x, y, w, h, t, theme)
    items = t.get("items", [])
    if not items:
        return
    cols = t.get("cols") or (3 if len(items) > 4 else max(1, len(items)))
    rows = -(-len(items) // cols)
    avail = (y + h) - by - 0.08
    cw = (w - 0.2 - GUTTER * (cols - 1)) / cols
    ch = (avail - 0.08 * (rows - 1)) / rows
    for i, it in enumerate(items):
        cx = x + 0.10 + (i % cols) * (cw + GUTTER)
        cy = by + (i // cols) * (ch + 0.08)
        box(slide, cx, cy, cw, ch, fill=C(theme, "bg"), radius=0.06)
        box(slide, cx, cy, cw, 0.045, fill=tone_color(theme, it.get("tone")))
        tail = it.get("delta") or it.get("note")
        block = 0.17 + 0.30 + (0.19 if tail else 0)
        ty = cy + max(0.07, (ch - block) / 2)          # centre the stack
        tw = cw - 0.12
        lab = str(it.get("label", "")).upper()
        label(slide, cx + 0.06, ty, tw, 0.17, lab, theme,
              size=fit_size(lab, tw, 0.17, 8, floor=6), bold=True,
              color=C(theme, "muted"), inset=0)
        val = str(it.get("value", ""))
        label(slide, cx + 0.06, ty + 0.17, tw, 0.30, val, theme,
              size=fit_size(val, tw, 0.28, 15, floor=9.5), bold=True, inset=0,
              font=theme["font_head"])
        if tail:
            label(slide, cx + 0.06, ty + 0.47, tw, 0.19, str(tail), theme,
                  size=fit_size(str(tail), tw, 0.19, 8.5, floor=6.5),
                  bold=bool(it.get("delta")), inset=0,
                  color=delta_color(theme, tail, it.get("dir"))
                  if it.get("delta") else C(theme, "muted"))


def t_callout(slide, x, y, w, h, t, theme):
    tone = tone_color(theme, t.get("tone"))
    box(slide, x, y, w, h, fill=C(theme, "bg"), line=C(theme, "line"),
        radius=0.08)
    box(slide, x, y, 0.075, h, fill=tone)
    ty = y + 0.10
    if t.get("title"):
        label(slide, x + 0.20, ty, w - 0.32, 0.20, str(t["title"]).upper(),
              theme, size=8.8, bold=True, color=tone, inset=0)
        ty += 0.21
    tw = w - 0.34
    big = str(t.get("value", ""))
    body = t.get("text")
    if big:
        # give the body whatever it needs first, then let the number use the rest
        body_h = min(text_h(body, tw, 9.5) + 0.06, ((y + h) - ty) * 0.62) if body else 0
        bh = max(0.24, (y + h) - ty - body_h - 0.06)
        bsize = t.get("size") or fit_size(big, tw, bh, 20, floor=12)
        label(slide, x + 0.20, ty, tw, bh, big, theme, size=bsize, bold=True,
              inset=0, font=theme["font_head"])
        ty += bh + 0.02
    if body:
        bh = (y + h) - ty - 0.07
        label(slide, x + 0.20, ty, tw, bh, body, theme,
              size=t.get("text_size") or fit_size(body, tw, bh, 9.5, floor=7),
              inset=0, line_spacing=1.02)


def t_table(slide, x, y, w, h, t, theme):
    by = tile_frame(slide, x, y, w, h, t, theme)
    head, rows = t.get("head", []), t.get("items", [])
    ncol = len(head) or (len(rows[0]) if rows else 1)
    weights = t.get("weights") or [1.6] + [1.0] * (ncol - 1)
    tot = sum(weights)
    widths = [(w - 0.2) * wt / tot for wt in weights]
    avail = (y + h) - by - 0.06
    rowh = min(0.34, avail / max(1, len(rows) + (1 if head else 0)))
    size = 8.8 if rowh > 0.24 else 8
    ry = by
    if head:
        box(slide, x + 0.10, ry, w - 0.2, rowh, fill=C(theme, "bg"), radius=0.04)
        cx = x + 0.10
        for j, hd in enumerate(head):
            label(slide, cx + 0.05, ry, widths[j] - 0.1, rowh, str(hd).upper(),
                  theme, size=fit_size(str(hd).upper(), widths[j] - 0.1, rowh,
                                       size - 0.6, floor=5.5),
                  bold=True, color=C(theme, "muted"),
                  anchor="m", align="l" if j == 0 else "r", inset=0)
            cx += widths[j]
        ry += rowh
    for i, row in enumerate(rows):
        if i % 2:
            box(slide, x + 0.10, ry, w - 0.2, rowh, fill=C(theme, "bg"))
        cx = x + 0.10
        for j, cell in enumerate(row[:ncol]):
            col = C(theme, "ink")
            if j and re.match(r"^\s*[▲▼+-]", str(cell)):
                col = delta_color(theme, str(cell).strip())
            label(slide, cx + 0.05, ry, widths[j] - 0.1, rowh, str(cell), theme,
                  size=min(size, fit_size(str(cell), widths[j] - 0.1, rowh,
                                          size, floor=6)),
                  bold=(j == 0), color=col, anchor="m",
                  align="l" if j == 0 else "r", inset=0)
            cx += widths[j]
        ry += rowh


def t_image(slide, x, y, w, h, t, theme):
    """Proof photos, contact-sheet style — many per page, never one per page."""
    by = tile_frame(slide, x, y, w, h, t, theme)
    paths = t.get("items", [])
    cap_h = 0.22 if any(isinstance(p, dict) and p.get("caption") for p in paths) else 0
    cols = t.get("cols") or min(4, max(1, len(paths)))
    rows = -(-max(1, len(paths)) // cols)
    cw = (w - 0.2 - 0.08 * (cols - 1)) / cols
    ch = ((y + h) - by - 0.06 - 0.08 * (rows - 1)) / rows
    for i, p in enumerate(paths):
        path, cap = (p.get("path"), p.get("caption")) if isinstance(p, dict) else (p, None)
        cx = x + 0.10 + (i % cols) * (cw + 0.08)
        cy = by + (i // cols) * (ch + 0.08)
        cell_h = ch - cap_h
        try:
            pic = slide.shapes.add_picture(path, Inches(cx), Inches(cy))
            # scale to fit the cell, keeping the photo's aspect ratio
            k = min(cw / (pic.width / 914400), cell_h / (pic.height / 914400))
            pic.width, pic.height = int(pic.width * k), int(pic.height * k)
            pic.left = Inches(cx + (cw - pic.width / 914400) / 2)
            pic.top = Inches(cy + (cell_h - pic.height / 914400) / 2)
        except Exception as exc:                       # noqa: BLE001
            print(f"  ! image {path}: {exc}", file=sys.stderr)
            box(slide, cx, cy, cw, cell_h, fill=C(theme, "bg"), radius=0.05)
        if cap:
            label(slide, cx, cy + ch - cap_h, cw, cap_h, cap, theme, size=7.5,
                  color=C(theme, "muted"), align="c", inset=0.02)


def t_text(slide, x, y, w, h, t, theme):
    by = tile_frame(slide, x, y, w, h, t, theme)
    label(slide, x + 0.10, by, w - 0.2, (y + h) - by - 0.06,
          t.get("items", t.get("text", "")), theme,
          size=t.get("size", 9.5), inset=0, line_spacing=1.02, space=4)


TILES = {"bullets": t_bullets, "bars": t_bars, "metrics": t_metrics,
         "callout": t_callout, "table": t_table, "image": t_image,
         "text": t_text}


# ----------------------------------------------------------------- slide ----

def draw_rows(slide, rows, theme, y, bottom):
    """Rows of tiles. Row height is proportional to `weight` (default 1)."""
    if not rows:
        return
    avail = bottom - y - GUTTER * (len(rows) - 1)
    fixed = sum(float(r["h"]) for r in rows if r.get("h"))
    flex = max(0.4, avail - fixed)
    total_w = sum(float(r.get("weight", 1)) for r in rows if not r.get("h")) or 1
    for row in rows:
        rh = float(row["h"]) if row.get("h") else \
            flex * float(row.get("weight", 1)) / total_w
        tiles = row.get("tiles", [])
        if not tiles:
            y += rh + GUTTER
            continue
        span = sum(float(t.get("span", 1)) for t in tiles)
        x = MARGIN
        usable = SLIDE_W - 2 * MARGIN - GUTTER * (len(tiles) - 1)
        for t in tiles:
            tw = usable * float(t.get("span", 1)) / span
            TILES.get(t.get("kind", "bullets"), t_bullets)(
                slide, x, y, tw, rh, t, theme)
            x += tw + GUTTER
        y += rh + GUTTER


def build_slide(prs, sp, theme, page, pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C(theme, "bg"))
    y = draw_header(slide, sp, theme, 0)
    if sp.get("kpis"):
        y = draw_kpis(slide, sp["kpis"], theme, y)
    foot = sp.get("footer") or sp.get("takeaways")
    fh = sp.get("footer_h", 0.58) if foot else 0.0
    bottom = SLIDE_H - (fh + 0.10 if foot else 0.30)
    draw_rows(slide, sp.get("rows", []), theme, y, bottom)
    if foot:
        draw_footer(slide, foot, theme, SLIDE_H - fh, fh,
                    sp.get("footer_title", "SO WHAT"))
    if sp.get("page_numbers", True):
        # sits in the footer bar when there is one, otherwise on the canvas
        py = SLIDE_H - (fh - 0.02 if foot else 0.26)
        label(slide, SLIDE_W - MARGIN - 0.55, py, 0.55, 0.20,
              f"{page}/{pages}", theme, size=8, align="r", inset=0,
              color=RGBColor.from_string("8A939C" if foot else "6B7682"))
    return slide


def build(spec, out):
    theme = dict(THEMES.get(spec.get("theme", "honasa"), THEMES["honasa"]))
    theme.update(spec.get("palette", {}))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
    slides = spec.get("slides", [])
    for i, sp in enumerate(slides, 1):
        build_slide(prs, sp, theme, i, len(slides))
    prs.save(out)
    return len(slides)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("spec")
    ap.add_argument("-o", "--out", default="deck.pptx")
    a = ap.parse_args()
    with open(a.spec, encoding="utf-8") as fh:
        spec = json.load(fh)
    n = build(spec, a.out)
    print(f"Wrote {a.out} — {n} slide(s)")


if __name__ == "__main__":
    main()
