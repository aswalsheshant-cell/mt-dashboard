#!/usr/bin/env python3
"""
Generate 1-pager executive PowerPoint presentation from MT Primary vs Offtake Analysis template.
Reads metrics from Excel, applies RAG status logic, and builds a 16:9 slide.
Idempotent: generates the same output for the same input.
"""

import os
import sys
from pathlib import Path
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

EXCEL_PATH = "MT_Primary_vs_Offtake_Analysis_Template.xlsx"
OUTPUT_PPTX = "MT_Primary_vs_Offtake_1Pager.pptx"

# RAG Status Colors
COLOR_RED = RGBColor(190, 40, 40)
COLOR_AMBER = RGBColor(200, 120, 0)
COLOR_GREEN = RGBColor(40, 140, 40)
COLOR_DARK = RGBColor(24, 43, 73)
COLOR_GRAY = RGBColor(110, 110, 110)
COLOR_BODY = RGBColor(16, 24, 40)

def safe_float(val, default=0.0):
    """Safely convert value to float, handling None, strings, and errors."""
    if val is None:
        return default
    try:
        return float(val)
    except (ValueError, TypeError):
        return default

def safe_str(val, default="–"):
    """Safely convert value to string, handling None."""
    if val is None or val == "":
        return default
    return str(val).strip()

def get_rag_status(gap_pct, primary_val=None):
    """
    Determine RAG status based on alignment gap percentage.
    gap_pct: absolute gap percentage or value
    Returns: ("Green"|"Amber"|"Red", color_rgb)
    """
    if gap_pct is None:
        return ("Gray", RGBColor(150, 150, 150))

    gap = abs(safe_float(gap_pct, 0))

    # Thresholds: <2% = Green, 2-5% = Amber, >5% = Red
    if gap <= 2.0:
        return ("Green", COLOR_GREEN)
    elif gap <= 5.0:
        return ("Amber", COLOR_AMBER)
    else:
        return ("Red", COLOR_RED)

def load_excel_metrics(filepath):
    """
    Extracts summary metrics and zone data from Excel template.
    Returns: (kpis, zone_data, alerts)
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"Excel template not found: {filepath}")

    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        ws = wb.active
    except Exception as e:
        raise RuntimeError(f"Failed to load Excel: {e}")

    # --- 1. Extract Executive Summary Metrics ---
    primary_nsv = safe_float(ws["B7"].value, 0.0)
    primary_mom = safe_float(ws["B8"].value, 0.0)
    primary_yoy = safe_float(ws["B9"].value, 0.0)

    offtake_nsv = safe_float(ws["B11"].value, 0.0)
    offtake_mom = safe_float(ws["B12"].value, 0.0)
    offtake_yoy = safe_float(ws["B13"].value, 0.0)

    gap_pct = safe_float(ws["B15"].value, 0.0)
    status_text = safe_str(ws["B16"].value, "Unknown")

    # Calculate gap if not explicitly provided
    if gap_pct == 0.0 and primary_nsv > 0:
        gap_pct = abs((primary_nsv - offtake_nsv) / primary_nsv * 100)

    gap_status, gap_color = get_rag_status(gap_pct)

    # Build KPI cards
    # Safe formatting for KPI display
    primary_str = f"₹{primary_nsv:.1f} Cr" if primary_nsv > 0 else "No data"
    offtake_str = f"₹{offtake_nsv:.1f} Cr" if offtake_nsv > 0 else "No data"
    gap_str = f"{gap_pct:.1f}%" if gap_pct > 0 else "TBD"

    primary_mom_str = (f"{primary_mom:+.1f}% MoM" if primary_mom != 0 else "MoM Data")
    offtake_mom_str = (f"{offtake_mom:+.1f}% MoM" if offtake_mom != 0 else "MoM Data")

    kpis = [
        ("Primary Sales", primary_str,
         primary_mom_str,
         COLOR_GREEN if primary_mom > 0 else (COLOR_RED if primary_mom < 0 else COLOR_GRAY)),

        ("Offtake Sales", offtake_str,
         offtake_mom_str,
         COLOR_GREEN if offtake_mom > 0 else (COLOR_RED if offtake_mom < 0 else COLOR_GRAY)),

        ("Alignment Gap", gap_str,
         f"Status: {gap_status}",
         gap_color),
    ]

    # --- 2. Extract Zone Breakdown Table (rows 27-32) ---
    zone_rows = []
    zone_names = ["North", "South-1", "South-2", "East", "West", "Central"]

    for idx, row_num in enumerate(range(27, 33)):
        zone_name = safe_str(ws.cell(row=row_num, column=1).value, zone_names[idx] if idx < len(zone_names) else f"Zone-{idx+1}")
        p_val = safe_float(ws.cell(row=row_num, column=2).value, 0.0)
        o_val = safe_float(ws.cell(row=row_num, column=3).value, 0.0)
        gap_val = safe_float(ws.cell(row=row_num, column=4).value, 0.0)

        # Calculate gap if not provided
        if gap_val == 0.0 and p_val > 0:
            gap_val = ((p_val - o_val) / p_val * 100)

        z_status, z_color = get_rag_status(gap_val)

        zone_rows.append({
            'name': zone_name,
            'primary': p_val,
            'offtake': o_val,
            'gap': gap_val,
            'status': z_status,
            'color': z_color
        })

    # --- 3. Extract Alert Bullets (from column F or use defaults) ---
    alerts = []
    for row_num in range(8, 12):
        alert_text = safe_str(ws.cell(row=row_num, column=6).value, None)
        if alert_text and alert_text != "–" and len(alert_text) > 5:
            alerts.append(alert_text)

    # If no alerts extracted, generate from data
    if not alerts:
        if gap_pct > 5.0:
            alerts.append(f"🔴 Critical: Alignment gap at {gap_pct:.1f}% exceeds 5% threshold.")
        elif gap_pct > 2.0:
            alerts.append(f"🟡 Alert: Alignment gap at {gap_pct:.1f}% requires attention.")
        else:
            alerts.append(f"🟢 On Track: Alignment gap at {gap_pct:.1f}% within target.")

        red_zones = [z['name'] for z in zone_rows if z['status'] == 'Red']
        if red_zones:
            alerts.append(f"🔴 {', '.join(red_zones)}: Zones show elevated gaps.")

        if primary_nsv > 0:
            conversion = (offtake_nsv / primary_nsv * 100)
            alerts.append(f"Secondary conversion: {conversion:.0f}% of primary.")
        alerts.append("Review with zone heads and category managers by week-end.")

    wb.close()
    return kpis, zone_rows, alerts

def build_presentation(kpis, zone_data, alerts):
    """Build and return PowerPoint presentation object."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # --- Title ---
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Monthly Modern Trade Primary vs. Offtake Snapshot"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    # --- KPI Cards (Top Row) ---
    card_w, card_h = Inches(2.7), Inches(1.15)
    start_x, start_y = Inches(0.8), Inches(1.25)

    for i, (title, val, sub, color) in enumerate(kpis):
        card_x = start_x + i * Inches(3.0)
        box = slide.shapes.add_textbox(card_x, start_y, card_w, card_h)
        tf = box.text_frame
        tf.word_wrap = True

        # Title
        p_title = tf.paragraphs[0]
        p_title.text = title.upper()
        p_title.font.size = Pt(9)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_GRAY
        p_title.space_after = Pt(2)

        # Value
        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.size = Pt(18)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_BODY
        p_val.space_after = Pt(2)

        # Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(9)
        p_sub.font.color.rgb = color

    # --- Zone Performance Table ---
    table_rows = len(zone_data) + 1
    table_shape = slide.shapes.add_table(table_rows, 4, Inches(0.8), Inches(2.65), Inches(6.8), Inches(4.0))
    table = table_shape.table

    # Header row
    headers = ["Zone", "Primary (₹ Cr)", "Offtake (₹ Cr)", "Gap Status"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cp = cell.text_frame.paragraphs[0]
        cp.font.bold = True
        cp.font.size = Pt(11)
        cp.font.color.rgb = RGBColor(255, 255, 255)
        # Set cell background to dark blue
        from pptx.dml.color import RGBColor as DmlColor
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK

    # Data rows
    for row_idx, zone in enumerate(zone_data, start=1):
        # Zone name
        cell = table.cell(row_idx, 0)
        cell.text = zone['name']
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Primary
        cell = table.cell(row_idx, 1)
        cell.text = f"{zone['primary']:.1f}"
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Offtake
        cell = table.cell(row_idx, 2)
        cell.text = f"{zone['offtake']:.1f}"
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Gap Status
        cell = table.cell(row_idx, 3)
        cell.text = f"{zone['status']}"
        cp = cell.text_frame.paragraphs[0]
        cp.font.size = Pt(10)
        cp.font.bold = True
        cp.font.color.rgb = zone['color']

    # --- Executive Action & Alerts Callout (Right side) ---
    callout_x, callout_y = Inches(8.0), Inches(2.65)
    callout_w, callout_h = Inches(4.5), Inches(4.0)

    callout = slide.shapes.add_textbox(callout_x, callout_y, callout_w, callout_h)
    cf = callout.text_frame
    cf.word_wrap = True
    cf.margin_bottom = Inches(0.1)
    cf.margin_left = Inches(0.1)
    cf.margin_right = Inches(0.1)

    # Callout title
    cp = cf.paragraphs[0]
    cp.text = "Executive Action & Alerts"
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = COLOR_RED
    cp.space_after = Pt(8)

    # Alert bullets
    for alert in alerts:
        p_bullet = cf.add_paragraph()
        p_bullet.text = alert
        p_bullet.font.size = Pt(10)
        p_bullet.space_after = Pt(6)
        p_bullet.level = 0

    # --- Footer / Metadata ---
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(12.0), Inches(0.4))
    ff = footer.text_frame
    p_footer = ff.paragraphs[0]
    p_footer.text = "Source: MT Dashboard Primary vs Offtake Analysis | Generated automatically from Excel template"
    p_footer.font.size = Pt(8)
    p_footer.font.italic = True
    p_footer.font.color.rgb = RGBColor(150, 150, 150)

    return prs

def main():
    try:
        print(f"📊 Loading Excel template: {EXCEL_PATH}")
        kpis, zone_data, alerts = load_excel_metrics(EXCEL_PATH)

        print(f"✓ Extracted {len(zone_data)} zones and {len(alerts)} alerts")

        print(f"🎨 Building presentation...")
        prs = build_presentation(kpis, zone_data, alerts)

        print(f"💾 Saving to {OUTPUT_PPTX}")
        prs.save(OUTPUT_PPTX)

        file_size = os.path.getsize(OUTPUT_PPTX)
        print(f"✅ Success! Generated {OUTPUT_PPTX} ({file_size:,} bytes)")
        return 0

    except Exception as e:
        print(f"❌ Error: {e}", file=sys.stderr)
        return 1

if __name__ == "__main__":
    sys.exit(main())
