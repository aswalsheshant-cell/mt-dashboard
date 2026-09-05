"""
BUILD_MT_MONTHLY_PPT.PY — Parameterized 18-Slide MT Leadership Deck Engine
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Purpose:
  Generate dynamic, reusable Modern Trade monthly performance decks with:
  - 18 modular slides (diagnostic + execution architecture)
  - Parameterized inputs (zone data, promo scenarios, color customization)
  - Clean Python-PPTX generation (no overlaps, validated shapes)
  - Dual export (PowerPoint .pptx + Google Slides JSON)

Usage:
  python build_mt_monthly_ppt.py --month september --year 2026 --theme teal

Architecture:
  Phase 1: Modular Engine (18 slide functions)
  Phase 2: Dynamic Math (calculations for waterfall, scenario, matrix)
  Phase 3: Dual Export (.pptx + .json for Google Slides)

Author: Claude Haiku 4.5 | Session: MT Intelligence Framework
Date: 09-Sep-2026
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""

import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
import os

# Import analytics engine for dynamic calculations (Phase 2)
from mt_analytics_engine import (
    calculate_waterfall_bridge,
    calculate_scenario_roi,
    calculate_matrix_coordinates
)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 1: COLOR PALETTE & THEME TOKENS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class MTTheme:
    """Reusable color palette for MT decks. Swap RGB values for brand customization."""

    def __init__(self, theme_name="default"):
        self.name = theme_name

        # Core backgrounds
        self.NAVY = RGBColor(13, 27, 42)           # #0D1B2A (dark)
        self.WHITE = RGBColor(255, 255, 255)
        self.LIGHT_GREY = RGBColor(240, 240, 240)
        self.MEDIUM_GREY = RGBColor(150, 150, 150)

        # Semantic colors (data viz)
        self.TEAL = RGBColor(42, 157, 176)         # #2A9DB0 (growth/hero)
        self.ACCENT_RED = RGBColor(230, 57, 70)    # #E63946 (alert/urgent)
        self.ACCENT_GREEN = RGBColor(42, 157, 126) # #2A9D7E (success)
        self.ACCENT_ORANGE = RGBColor(247, 162, 97)# #F7A261 (caution/watch)
        self.ACCENT_PURPLE = RGBColor(150, 100, 200) # Velocity/future

    def get_status_color(self, status):
        """Return color for status indicator."""
        status_map = {
            "URGENT": self.ACCENT_RED,
            "FIX": self.ACCENT_RED,
            "WATCH": self.ACCENT_ORANGE,
            "OK": self.ACCENT_GREEN,
            "BENCHMARK": self.ACCENT_GREEN,
        }
        return status_map.get(status, self.NAVY)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 2: CONFIGURATION & DATA STRUCTURES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

DEFAULT_CONFIG = {
    "month": "September",
    "year": 2026,
    "title": "Modern Trade Leadership Review",

    # Executive metrics
    "period": "Q1-Sep",
    "q1_offtake": "₹114.39 Cr",
    "4m_offtake": "₹185.81 Cr",
    "month_offtake": "₹71.42 Cr",
    "q1_growth_seq": "+27%",
    "q1_growth_yoy": "+64%",

    # Zone data (name → {offtake, conversion%, gap, status, yoy_growth})
    "zones": {
        "Central": {"offtake": "₹2.12 Cr", "conv": 78.8, "gap": "₹0.57 Cr", "status": "WATCH", "yoy_growth": 18},
        "North": {"offtake": "₹6.99 Cr", "conv": 58.5, "gap": "₹4.97 Cr", "status": "FIX", "yoy_growth": 22},
        "East": {"offtake": "₹3.55 Cr", "conv": 45.3, "gap": "₹4.28 Cr", "status": "URGENT", "yoy_growth": 15},
        "South-1": {"offtake": "₹8.19 Cr", "conv": 83.6, "gap": "₹1.61 Cr", "status": "WATCH", "yoy_growth": 26},
        "South-2": {"offtake": "₹4.91 Cr", "conv": 71.3, "gap": "₹1.98 Cr", "status": "FIX", "yoy_growth": 20},
        "West": {"offtake": "₹8.28 Cr", "conv": 82.3, "gap": "₹1.78 Cr", "status": "WATCH", "yoy_growth": 24},
    },

    # Diagnostic chain data (for Slide 5c waterfall auto-balancing)
    "diagnostic_chain": {
        "chain_name": "Reliance",
        "primary": 2.40,  # ₹ Cr
        "offtake": 1.25,  # ₹ Cr
    },

    # Zones detail for Risk-Opportunity Matrix (Slide 7)
    "zones_detail": [
        {"name": "East", "conversion": 45.3, "nsv": 3.55},
        {"name": "South-2", "conversion": 71.3, "nsv": 4.91},
        {"name": "North", "conversion": 58.5, "nsv": 6.99},
        {"name": "West", "conversion": 82.3, "nsv": 8.28},
        {"name": "Central", "conversion": 78.8, "nsv": 2.12},
        {"name": "South-1", "conversion": 83.6, "nsv": 8.19},
    ],

    # Chain data (name → {offtake, conversion%})
    "chains": [
        {"name": "DMart", "offtake": "₹52.1 Cr", "conv": 45.6},
        {"name": "Reliance", "offtake": "₹28.7 Cr", "conv": 25.1},
        {"name": "Apollo", "offtake": "₹18.2 Cr", "conv": 15.9},
    ],

    # Brand data
    "brands": [
        {"name": "MAMAEARTH", "offtake": "₹82.5 Cr", "growth": "+33%"},
        {"name": "TDC", "offtake": "₹29.5 Cr", "growth": "+365%"},
        {"name": "AQUALOGICA", "offtake": "₹1.96 Cr", "growth": "+16%"},
    ],

    # Scenario analysis defaults (Slide 12)
    "scenario": {
        "zone": "East",
        "current_conv": 0.453,
        "promo_spend": 30.0,  # ₹ Lakhs
        "promo_discount": 0.10,  # 10%
        "target_conv": 0.70,
        "days": 21,
        "lift_pp": 25,  # percentage points
    },

    # Scenario ROI params for dynamic calculation (Phase 2)
    "scenario_params": {
        "current_offtake_weekly": 7.0,  # ₹ Cr (East zone weekly)
        "current_conv": 45.3,  # % (East zone current conversion)
        "target_conv": 70.0,  # % (target)
        "promo_spend": 30.0,  # ₹ Lakhs
        "promo_days": 21,
        "gross_margin_pct": 0.45,  # 45%
        "discount_pct": 10.0,  # 10% discount
    },

    # Theme
    "theme_name": "default",
}

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 3: HELPER FUNCTIONS — SHAPE & TEXT UTILITIES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def add_background_color(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def clear_slide(slide):
    """Remove all default shapes from a slide."""
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

def add_heading(slide, text, top_inches=0.4, size_pt=40, color=None, bold=True):
    """Add a headline text box (top of slide)."""
    if color is None:
        color = RGBColor(255, 255, 255)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(top_inches), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_subtitle(slide, text, top_inches=1.8, size_pt=18, color=None):
    """Add a subtitle/context line."""
    if color is None:
        color = RGBColor(200, 200, 200)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(top_inches), Inches(9), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.color.rgb = color
    return box

def add_footer(slide, text, size_pt=10, color=None):
    """Add footer text (bottom of slide)."""
    if color is None:
        color = RGBColor(150, 150, 150)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.color.rgb = color
    return box

def add_kpi_card(slide, label, value, top_inches, left_inches, status_color, theme):
    """Add a KPI metric card (number + label + status indicator)."""
    # Background box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left_inches), Inches(top_inches),
        Inches(2.0), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.NAVY
    shape.line.color.rgb = status_color
    shape.line.width = Pt(2)

    # Value text
    value_box = slide.shapes.add_textbox(
        Inches(left_inches + 0.1), Inches(top_inches + 0.05),
        Inches(1.8), Inches(0.45)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = theme.WHITE

    # Label text
    label_box = slide.shapes.add_textbox(
        Inches(left_inches + 0.1), Inches(top_inches + 0.5),
        Inches(1.8), Inches(0.35)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = theme.MEDIUM_GREY

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 4: 18 SLIDE GENERATION FUNCTIONS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def slide_01_title(prs, config, theme):
    """Slide 1: Title Slide with month + theme."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    # Title
    add_heading(slide, f"Modern Trade Leadership Review — {config['month']} {config['year']}",
                top_inches=2.0, size_pt=44)

    # Subtitle
    add_subtitle(slide,
                 f"Q1–{config['month'][:3]} FY27 | Zone Accountability & Market Share Analysis",
                 top_inches=3.5, size_pt=18, color=theme.LIGHT_GREY)

    # Date
    add_footer(slide, f"{datetime.now().strftime('%d-%B-%Y')} | MAMAEARTH MT INTELLIGENCE",
               size_pt=11, color=theme.MEDIUM_GREY)

def slide_02_toc(prs, config, theme):
    """Slide 2: Table of Contents (18-slide structure)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "18-Slide Diagnostic Architecture", top_inches=0.3, size_pt=36)

    # Two-column TOC
    toc_items = [
        "1. Title Slide",
        "2. Table of Contents",
        "3. Executive Summary",
        "4. Market Context",
        "5. Primary Trend (3M)",
        "5b. Offtake Trend (3M)",
        "5c. Multi-Step Waterfall",
        "6. Zone Primary Performance",
        "7. Zone Risk-Opportunity Matrix",
        "8. Zone Conversion % & Status",
        "9. Chain Concentration Risk",
        "9b. 4-Pillar Strategy (Focus Zone)",
        "10. Category Mix & Brands",
        "11. Multi-Period Comparison (Sparklines)",
        "12. Scenario Analysis & ROI",
        "13. 30-Day Execution Roadmap",
        "14. Action Register (Live Tracking)",
        "15–18. Appendix & Closing",
    ]

    # Left column (1-9)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(toc_items[:9]):
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = theme.WHITE
        p.space_before = Pt(3)

    # Right column (10-18)
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(toc_items[9:]):
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = theme.WHITE
        p.space_before = Pt(3)

    add_footer(slide, "Navigate: Diagnostic context → Zone prioritization → Execution roadmap → Accountability",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_03_exec_summary(prs, config, theme):
    """Slide 3: Executive Summary (Q1-Month metrics + 3 bullets)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Executive Summary: Q1–Sep Performance", top_inches=0.3, size_pt=36)

    # KPI cards
    add_kpi_card(slide, "Q1 Offtake", config["q1_offtake"], 1.2, 0.5,
                 theme.TEAL, theme)
    add_kpi_card(slide, "4-Month Total", config["4m_offtake"], 1.2, 2.8,
                 theme.TEAL, theme)
    add_kpi_card(slide, f"{config['month']} Offtake", config["month_offtake"], 1.2, 5.1,
                 theme.TEAL, theme)
    add_kpi_card(slide, "YoY Growth", config["q1_growth_yoy"], 1.2, 7.4,
                 theme.ACCENT_GREEN, theme)

    # Three bullets
    bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True

    bullets = [
        "Three zones below target conversion (75-80%): East 45.3%, North 58.5%, South-2 71.3% = ₹10.2 Cr gap",
        "DMart concentrated at 45.6% conversion — Reliance only 25.1% — working capital stranded in retail",
        "Priority: 30-day recovery sprint (East + North), hero SKU focus, weekly velocity tracking, promo mechanics",
    ]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(13)
        p.font.color.rgb = theme.WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)

    add_footer(slide, f"Baseline: 75-80% conversion target | Source: FY27 Q1 offtake + Aug sell-out audit",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_04_market_context(prs, config, theme):
    """Slide 4: Market Context (Competitive positioning + category shifts)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Market Context: Mamaearth position vs. competition", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "MT Personal Care ₹4,200L NSV | Mamaearth 6.4% share | Nielsen MAT validation",
                 top_inches=1.0, size_pt=13)

    # Left: Market share pie
    context_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf = context_box.text_frame
    tf.word_wrap = True

    context_lines = [
        "Market Share (MT Category)",
        "",
        "HUL:  28.0% (₹1,176 Cr) — Leader",
        "P&G:  15.0% (₹630 Cr) — Strong #2",
        "ITC:   9.0% (₹378 Cr) — Regional",
        "Mamamearth: 6.4% (₹270 Cr) — Challenger",
        "Others: 41.6% (₹1,746 Cr)",
        "",
        "Category Growth: +11% YoY",
        "Value Tier (<₹400): +18% YoY ⚠️",
    ]

    for i, line in enumerate(context_lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
        else:
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.font.size = Pt(11)

        p.text = line
        p.font.color.rgb = theme.WHITE if line else theme.NAVY

    # Right: Key insight
    insight_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.5))
    tf = insight_box.text_frame
    tf.word_wrap = True

    # Insight box background
    insight_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.0), Inches(1.8),
        Inches(4.5), Inches(3.0)
    )
    insight_shape.fill.solid()
    insight_shape.fill.fore_color.rgb = RGBColor(30, 50, 70)
    insight_shape.line.color.rgb = theme.ACCENT_ORANGE
    insight_shape.line.width = Pt(2)

    insight_text = """KEY INSIGHT

Value tier gaining 18% YoY while Mamaearth premium lines flat.

Action:
Launch sub-₹500 line by Oct 1 to defend share in value migration."""

    insight_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.1), Inches(2.6))
    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = insight_text
    p.font.size = Pt(11)
    p.font.color.rgb = theme.WHITE
    p.line_spacing = 1.3

    add_footer(slide, "Source: Nielsen MAT FY27 Q1 | Category: Face Wash + Shampoo + Sun Care",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_05_primary_trend(prs, config, theme):
    """Slide 5: Primary Trend (3-month line chart)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Primary Revenue Trend: 3-Month MoM Growth", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "Billed purchases by MT chains (FY27 Q1-Sep)",
                 top_inches=1.0, size_pt=12)

    # Placeholder: Chart would be inserted here (Phase 2)
    chart_placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
    tf = chart_placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[LINE CHART: July ₹71.42 Cr → Aug ₹72.15 Cr → Sep ₹73.24 Cr (+0.9% → +1.5% MoM)]"
    p.font.size = Pt(14)
    p.font.color.rgb = theme.LIGHT_GREY
    p.alignment = PP_ALIGN.CENTER

    # Data table below chart
    table_data = [
        ["Month", "Primary (₹ Cr)", "MoM %", "YoY %"],
        ["July-26", "₹71.42", "+4.2%", "+78%"],
        ["Aug-26", "₹72.15", "+0.9%", "+82%"],
        ["Sep-26", "₹73.24", "+1.5%", "+85%"],
    ]

    table_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.5))
    tf = table_box.text_frame
    tf.word_wrap = True

    for row_idx, row in enumerate(table_data):
        for col_idx, cell in enumerate(row):
            p = tf.add_paragraph() if row_idx > 0 or col_idx > 0 else tf.paragraphs[0]
            p.text = f"{cell:>15}" if col_idx > 0 else cell
            p.font.size = Pt(11)
            p.font.bold = (row_idx == 0)
            p.font.color.rgb = theme.ACCENT_GREEN if row_idx == 0 else theme.WHITE

    add_footer(slide, "Consistent +1-2% MoM growth; YoY acceleration driven by Reliance + DMart base expansion",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_05b_offtake_trend(prs, config, theme):
    """Slide 5b: Offtake Trend (inventory pulled by MT)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Offtake Trend: Inventory Pulled by MT (3-Month View)", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "Store-level sell-through validation (FY27 Q1-Sep)",
                 top_inches=1.0, size_pt=12)

    # Placeholder for chart
    chart_placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
    tf = chart_placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[LINE CHART: July ₹52.68 Cr → Aug ₹53.94 Cr → Sep ₹55.18 Cr (trend up, but below primary)]"
    p.font.size = Pt(14)
    p.font.color.rgb = theme.LIGHT_GREY
    p.alignment = PP_ALIGN.CENTER

    # Data table
    table_data = [
        ["Month", "Offtake (₹ Cr)", "Conversion %", "Gap (₹ Cr)"],
        ["July-26", "₹52.68", "73.8%", "₹18.74"],
        ["Aug-26", "₹53.94", "74.8%", "₹18.21"],
        ["Sep-26", "₹55.18", "75.3%", "₹18.06"],
    ]

    table_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(8), Inches(1.5))
    tf = table_box.text_frame

    for row_idx, row in enumerate(table_data):
        for col_idx, cell in enumerate(row):
            p = tf.add_paragraph() if row_idx > 0 or col_idx > 0 else tf.paragraphs[0]
            p.text = f"{cell:>15}" if col_idx > 0 else cell
            p.font.size = Pt(11)
            p.font.bold = (row_idx == 0)
            p.font.color.rgb = theme.ACCENT_GREEN if row_idx == 0 else theme.WHITE

    add_footer(slide, "Conversion slowly improving (73.8% → 75.3%) but gap still ₹18+ Cr — East/North drag",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_05c_waterfall(prs, config, theme):
    """Slide 5c: Multi-Step Conversion Waterfall (diagnostic root cause)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    # Extract diagnostic inputs and calculate waterfall dynamically
    chain_data = config.get("diagnostic_chain", {})
    primary = chain_data.get("primary", 2.40)
    realized_offtake = chain_data.get("offtake", 1.25)
    chain_name = chain_data.get("chain_name", "Reliance")

    # Run calculation engine
    bridge = calculate_waterfall_bridge(primary, realized_offtake)

    add_heading(slide, f"Multi-Step Waterfall: Where Primary Gets Lost ({chain_name} Case Study)",
                top_inches=0.3, size_pt=32)
    add_subtitle(slide, f"Root cause diagnostic: {chain_name} primary ₹{bridge['primary']:.2f} Cr → current offtake ₹{bridge['realized_offtake']:.2f} Cr ({bridge['conversion_rate']}% realized)",
                 top_inches=1.0, size_pt=12)

    # Waterfall cards (5 boxes) — dynamically calculated
    waterfall_cards = [
        ("DISPATCHED", f"₹{bridge['primary']:.2f} Cr", theme.TEAL, "Primary Inflow"),
        ("LOSS 1: Shelf Creep", f"−₹{bridge['shelf_loss']:.2f} Cr", theme.ACCENT_RED, "Competitor creep"),
        ("LOSS 2: Price Cliff", f"−₹{bridge['price_loss']:.2f} Cr", theme.ACCENT_RED, "Elasticity loss"),
        ("LOSS 3: Trapped NPI", f"−₹{bridge['stuck_inventory']:.2f} Cr", theme.ACCENT_RED, "Slow SKUs"),
        ("REALIZED OFFTAKE", f"₹{bridge['realized_offtake']:.2f} Cr", theme.ACCENT_GREEN, f"{bridge['conversion_rate']}% Conv"),
    ]

    left_offset = 0.4
    for idx, (label, value, color, detail) in enumerate(waterfall_cards):
        x = left_offset + (idx * 1.75)

        # Card box
        card_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(2.2),
            Inches(1.6), Inches(1.8)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(25, 45, 70)
        card_shape.line.color.rgb = color
        card_shape.line.width = Pt(3)

        # Value
        value_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(2.35), Inches(1.5), Inches(0.6))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(3.0), Inches(1.5), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Detail
        detail_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(3.6), Inches(1.5), Inches(0.4))
        tf = detail_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(8)
        p.font.color.rgb = theme.LIGHT_GREY
        p.alignment = PP_ALIGN.CENTER

    # Action mandate
    action_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    action_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.4), Inches(4.4),
        Inches(9.2), Inches(1.6)
    )
    action_shape.fill.solid()
    action_shape.fill.fore_color.rgb = RGBColor(40, 60, 90)
    action_shape.line.color.rgb = theme.ACCENT_ORANGE
    action_shape.line.width = Pt(2)

    tf = action_box.text_frame
    tf.word_wrap = True

    action_text = "ACTION MANDATE: Immediately freeze non-hero NPI dispatches (Loss 3). Reallocate shelf space from competitor (Loss 1). Monitor price elasticity testing for sub-₹400 bundles (Loss 2). Target: Reliance conversion 52% → 70% by Oct 31 (release ₹1.15 Cr working capital)."

    p = tf.paragraphs[0]
    p.text = action_text
    p.font.size = Pt(11)
    p.font.color.rgb = theme.WHITE
    p.font.bold = False

    add_footer(slide, "Diagnostic: Reliance case exemplifies East/North pattern (primary growth ≠ sell-out growth)",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_06_zone_primary(prs, config, theme):
    """Slide 6: Zone-Wise Primary Performance."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Zone-Wise Primary: NSV & Growth Ranking", top_inches=0.3, size_pt=34)

    # Sort zones by offtake (descending)
    sorted_zones = sorted(config["zones"].items(),
                         key=lambda x: float(x[1]["offtake"].replace("₹", "").replace(" Cr", "")),
                         reverse=True)

    zones_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    tf = zones_box.text_frame
    tf.word_wrap = True

    header_line = "Zone           Primary (₹ Cr)   Growth YoY   Conversion %   Status"
    p = tf.paragraphs[0]
    p.text = header_line
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = theme.ACCENT_GREEN

    for zone_name, zone_data in sorted_zones:
        p = tf.add_paragraph()
        status_color = theme.get_status_color(zone_data["status"])
        yoy = zone_data.get("yoy_growth", 20)
        line = f"{zone_name:12}  {zone_data['offtake']:>10}   +{yoy:>2}%        {zone_data['conv']:>5.1f}%      {zone_data['status']}"
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = theme.WHITE

    # Insight callout
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1.2))
    tf = insight_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "INSIGHT: East lowest absolute offtake (₹3.55 Cr) + worst conversion (45.3%) = highest recovery upside. Central & South-1 are stable benchmarks (>78% conversion)."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.TEAL
    p.font.italic = False

    add_footer(slide, "Primary growth driven by chain expansion (Reliance +50 stores, DMart +80 stores); conversion varies by zone execution",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_07_risk_matrix(prs, config, theme):
    """Slide 7: Zone Risk-Opportunity Matrix (2x2) — dynamically plotted."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Territory Prioritization: Zone Risk vs. Opportunity Matrix", top_inches=0.3, size_pt=32)

    # Get zones detail and compute dynamic coordinates
    zones_data = config.get("zones_detail", [
        {"name": "East", "conversion": 45.3, "nsv": 3.55},
        {"name": "South-2", "conversion": 71.3, "nsv": 4.91},
        {"name": "North", "conversion": 58.5, "nsv": 6.99},
        {"name": "West", "conversion": 82.3, "nsv": 8.28},
        {"name": "Central", "conversion": 78.8, "nsv": 2.12},
        {"name": "South-1", "conversion": 83.6, "nsv": 8.19},
    ])

    # Canvas dimensions for the matrix area
    plot_box_left = Inches(1.5)
    plot_box_top = Inches(1.8)
    plot_box_width = Inches(7.0)
    plot_box_height = Inches(4.0)

    # Compute coordinates dynamically
    mapped_zones = calculate_matrix_coordinates(
        zones_data,
        1.5,  # box_left in inches
        1.8,  # box_top in inches
        7.0,  # box_width in inches
        4.0,  # box_height in inches
        target_conv=75.0
    )

    # Draw 2x2 grid with dividers
    mid_x_inches = plot_box_left + (plot_box_width / 2)
    mid_y_inches = plot_box_top + (plot_box_height / 2)

    # Vertical divider
    slide.shapes.add_shape(
        1,  # Rectangle
        mid_x_inches, plot_box_top,
        Inches(0.02), plot_box_height
    ).line.color.rgb = theme.MEDIUM_GREY

    # Horizontal divider
    slide.shapes.add_shape(
        1,  # Rectangle
        plot_box_left, mid_y_inches,
        plot_box_width, Inches(0.02)
    ).line.color.rgb = theme.MEDIUM_GREY

    # Quadrant labels
    q_labels = [
        ("WATCH / MEDIUM SCALE", Inches(1.8), Inches(2.0), theme.ACCENT_ORANGE),
        ("CRITICAL INTERVENTION", Inches(5.2), Inches(2.0), theme.ACCENT_RED),
        ("STABLE / SMALL SCALE", Inches(1.8), Inches(5.0), theme.MEDIUM_GREY),
        ("BENCHMARK / CORE", Inches(5.8), Inches(5.0), theme.ACCENT_GREEN),
    ]

    for label, left, top, color in q_labels:
        lbl_box = slide.shapes.add_textbox(left, top, Inches(2.6), Inches(0.35))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color

    # Plot dynamic zone bubbles
    for zone in mapped_zones:
        color_map = {
            "RED": theme.ACCENT_RED,
            "ORANGE": theme.ACCENT_ORANGE,
            "YELLOW": RGBColor(200, 160, 50),
            "GREEN": theme.ACCENT_GREEN,
        }
        bubble_color = color_map.get(zone["color_theme"], theme.MEDIUM_GREY)

        # Zone bubble box
        zone_left_inches = zone["x_coord"]
        zone_top_inches = zone["y_coord"]

        zone_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(zone_left_inches), Inches(zone_top_inches),
            Inches(1.8), Inches(0.65)
        )
        zone_shape.fill.solid()
        zone_shape.fill.fore_color.rgb = RGBColor(40, 60, 90)
        zone_shape.line.color.rgb = bubble_color
        zone_shape.line.width = Pt(2)

        # Zone text
        zone_text_box = slide.shapes.add_textbox(
            Inches(zone_left_inches + 0.05),
            Inches(zone_top_inches + 0.05),
            Inches(1.7), Inches(0.55)
        )
        tf = zone_text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{zone['name']}\n₹{zone['nsv']:.2f}Cr | {zone['conversion']:.1f}%"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = bubble_color
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "X-axis: Conversion gap vs 75% target (right = worse); Y-axis: NSV scale (top = larger). Dynamically computed quadrant placement.",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_08_zone_conversion(prs, config, theme):
    """Slide 8: Zone Conversion % with Status Indicators."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Zone Conversion %: Current Status vs. 75% Target", top_inches=0.3, size_pt=34)

    # Bar chart representation using boxes
    zones_sorted = sorted(config["zones"].items(),
                         key=lambda x: x[1]["conv"], reverse=True)

    bar_top = 1.5
    for zone_name, zone_data in zones_sorted:
        conv_pct = zone_data["conv"]
        bar_width = (conv_pct / 100) * 6  # Scale to 6 inches max
        status_color = theme.get_status_color(zone_data["status"])

        # Background bar
        bar_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2.5), Inches(bar_top),
            Inches(bar_width), Inches(0.35)
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = status_color
        bar_shape.line.width = Pt(0)

        # Zone label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(bar_top), Inches(2.0), Inches(0.35))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = zone_name
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE

        # Percentage text
        pct_box = slide.shapes.add_textbox(Inches(2.5 + bar_width + 0.1), Inches(bar_top), Inches(1.0), Inches(0.35))
        tf = pct_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{conv_pct:.1f}%"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = status_color

        bar_top += 0.5

    # Target line annotation
    target_box = slide.shapes.add_textbox(Inches(2.5 + 4.5), Inches(1.3), Inches(2), Inches(0.4))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "← 75% Target Benchmark"
    p.font.size = Pt(10)
    p.font.color.rgb = theme.ACCENT_GREEN

    # Insights
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.0))
    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Below-Target Zones: East 45.3% (−29.7pp), North 58.5% (−16.5pp), South-2 71.3% (−3.7pp) = Combined gap ₹10.23 Cr. Action intensity by zone: East (urgent 30-day sprint), North (20-store promo), South-2 (DMart shelf reset)."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.WHITE

    add_footer(slide, "Benchmark 75-80% conversion is standard for mature MT categories. Mamaearth blended: 66.8% (below target) — zone variance is execution issue, not demand",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_09_chain_concentration(prs, config, theme):
    """Slide 9: Chain Concentration Risk."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Chain Concentration: Top 3 Chains Drive 87% of NSV", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "Conversion variance by chain: DMart 45.6% vs. Apollo 68% — execution inconsistency",
                 top_inches=1.0, size_pt=12)

    # Chain cards
    chain_top = 1.8
    for chain_data in config["chains"]:
        chain_name = chain_data["name"]
        offtake = chain_data["offtake"]
        conv = chain_data["conv"]

        # Determine color by conversion
        if conv < 50:
            color = theme.ACCENT_RED
        elif conv < 70:
            color = theme.ACCENT_ORANGE
        else:
            color = theme.ACCENT_GREEN

        # Card
        card_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(chain_top),
            Inches(9.0), Inches(0.8)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(30, 50, 70)
        card_shape.line.color.rgb = color
        card_shape.line.width = Pt(2)

        # Text
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(chain_top + 0.1), Inches(8.6), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{chain_name:15} | NSV {offtake:>10} | Conversion {conv:>5.1f}% | Status: {'🟢 OK' if conv >= 70 else '🟡 WATCH' if conv >= 50 else '🔴 CRITICAL'}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE

        chain_top += 1.0

    # Summary
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
    tf = summary_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "CONCENTRATION RISK: DMart (45.6% of offtake) at lowest conversion — ₹7.2 Cr stuck. Reliance (43 stores post-expansion) needs NPI halt + hero SKU focus. Apollo healthy but smaller (15.9% mix)."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.TEAL
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "mitigation: Weekly chain-wise velocity reporting; promo intensity by chain (DMart EDLP sensitivity vs. Reliance promotional elasticity)."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.WHITE

    add_footer(slide, "Concentration creates planning complexity: 3 chains require 3 different mechanics (DMart=efficiency, Reliance=promo response, Apollo=steady fill)",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_09b_strategy_pillars(prs, config, theme):
    """Slide 9b: 4-Pillar Strategy Framework (for focus zone)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "4-Pillar East Turnaround Strategy: 45% → 70% Conversion (30 Days)",
                top_inches=0.3, size_pt=32)

    pillars = [
        {
            "num": "1",
            "name": "Hero SKU Focus",
            "icon": "⭐",
            "desc": "Restructure to Rice & Onion only for 14 days",
            "metric": "Compress SKUs 18→2; +35% velocity",
            "color": theme.TEAL,
        },
        {
            "num": "2",
            "name": "Price Elasticity",
            "icon": "💰",
            "desc": "Pilot sub-₹400 bundles at 10 Reliance stores",
            "metric": "Sep 15 rollout if margin ≥58%",
            "color": theme.ACCENT_GREEN,
        },
        {
            "num": "3",
            "name": "Shelf Excellence",
            "icon": "📦",
            "desc": "Planogram resets; eye-level reclaim from HUL",
            "metric": "Maintain ≥40% beauty bay facings",
            "color": theme.ACCENT_ORANGE,
        },
        {
            "num": "4",
            "name": "Velocity Pulse",
            "icon": "📈",
            "desc": "Weekly sell-out tracking + daily OSA flags",
            "metric": "≥5 pp weekly conversion step-up",
            "color": theme.ACCENT_PURPLE,
        },
    ]

    pillar_top = 1.3
    pillar_left_start = 0.5

    for pillar in pillars:
        left = pillar_left_start + (pillars.index(pillar) * 2.2)

        # Pillar box
        pillar_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(left), Inches(pillar_top),
            Inches(2.0), Inches(4.0)
        )
        pillar_shape.fill.solid()
        pillar_shape.fill.fore_color.rgb = RGBColor(25, 45, 70)
        pillar_shape.line.color.rgb = pillar["color"]
        pillar_shape.line.width = Pt(3)

        # Number + icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(pillar_top + 0.2), Inches(1.8), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pillar['icon']} Pillar {pillar['num']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = pillar["color"]
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(pillar_top + 0.8), Inches(1.8), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pillar["name"]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(pillar_top + 1.5), Inches(1.8), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pillar["desc"]
        p.font.size = Pt(9)
        p.font.color.rgb = theme.LIGHT_GREY

        # Metric
        metric_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(pillar_top + 2.9), Inches(1.8), Inches(1.0))
        tf = metric_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Metric:\n{pillar['metric']}"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = pillar["color"]
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "All 4 pillars must execute in parallel (not sequentially) to hit 70% conversion target by Sep 30",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_10_category_mix(prs, config, theme):
    """Slide 10: Category Mix & Brand Performance."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Brand Performance: Growth Momentum & Mix Shift", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "FY27 Q1 NSV | Mamaearth dominates but TDC explosive; Aqualogica niche",
                 top_inches=1.0, size_pt=12)

    # Brand cards
    brand_top = 1.8
    for brand in config["brands"]:
        growth_pct = int(brand["growth"].replace("%", "").replace("+", ""))
        growth_color = theme.ACCENT_GREEN if growth_pct > 0 else theme.ACCENT_RED

        # Card
        card_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.5), Inches(brand_top),
            Inches(7.0), Inches(0.75)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(30, 50, 70)
        card_shape.line.color.rgb = growth_color
        card_shape.line.width = Pt(2)

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.7), Inches(brand_top + 0.1), Inches(6.6), Inches(0.55))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{brand['name']:15} | NSV {brand['offtake']:>10} | YoY Growth {brand['growth']:>6}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = growth_color

        brand_top += 1.0

    # Insight
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
    tf = insight_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "KEY INSIGHT: TDC growing 365% YoY (from low base ₹0.81 Cr → ₹3.73 Cr); Mamaearth core +33% shows maturity. Mix evolution: Mamaearth premium positioning under pressure from value tier (see Slide 4 market context). Action: Allocate 40% of heroSKU loading to TDC in East/North to ride momentum."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.WHITE
    p.line_spacing = 1.3

    add_footer(slide, "Source: FY27 Q1 brand-level offtake | Aqualogica premium niche (₹1.96 Cr, +16%) shows sun care gaining traction",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_11_comparison_table(prs, config, theme):
    """Slide 11: Multi-Period Comparison with Sparklines (placeholder)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "3-Month Comparison: Primary vs. Offtake Trend", top_inches=0.3, size_pt=34)
    add_subtitle(slide, "July → Aug → Sep | Conversion gap narrowing but still ₹18+ Cr",
                 top_inches=1.0, size_pt=12)

    # Table header
    table_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = table_box.text_frame
    tf.word_wrap = True

    table_data = [
        ["Month", "Primary (₹ Cr)", "Offtake (₹ Cr)", "Conv %", "Gap (₹ Cr)", "Trend"],
        ["July-26", "₹71.42", "₹52.68", "73.8%", "₹18.74", "📊"],
        ["Aug-26", "₹72.15", "₹53.94", "74.8%", "₹18.21", "↗"],
        ["Sep-26", "₹73.24", "₹55.18", "75.3%", "₹18.06", "↗↗"],
    ]

    for row_idx, row in enumerate(table_data):
        for col_idx, cell in enumerate(row):
            p = tf.add_paragraph() if row_idx > 0 or col_idx > 0 else tf.paragraphs[0]

            # Format cell
            if col_idx > 0:
                cell_text = f"{cell:>12}"
            else:
                cell_text = cell

            p.text = cell_text
            p.font.size = Pt(11) if row_idx == 0 else Pt(10)
            p.font.bold = (row_idx == 0)
            p.font.color.rgb = theme.ACCENT_GREEN if row_idx == 0 else theme.WHITE

    # Summary insight
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
    tf = summary_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "INSIGHT: Conversion improving (+1.5 pp MoM) but gap remains large (₹18 Cr = ₹90+ Cr annualized loss). East + North account for ₹10.2 Cr of this gap. At current trend (+1.5 pp/month), full recovery (target 80%) will take 4+ months without intervention."
    p.font.size = Pt(11)
    p.font.color.rgb = theme.TEAL

    add_footer(slide, "Trend: Positive but sub-linear. Promo blitz + NPI halt required to accelerate conversion from +1.5% to +5% MoM",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_12_scenario_analysis(prs, config, theme):
    """Slide 12: Parameterized Scenario Analysis (ROI calculation)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    # Extract scenario params and run dynamic ROI calculation
    scenario_cfg = config.get("scenario_params", {
        "current_offtake_weekly": 7.0,
        "current_conv": 45.3,
        "target_conv": 70.0,
        "promo_spend": 30.0,
        "promo_days": 21,
        "gross_margin_pct": 0.45,
        "discount_pct": 10.0,
    })

    # Run analytics engine
    roi = calculate_scenario_roi(
        current_offtake_weekly=scenario_cfg["current_offtake_weekly"],
        current_conv=scenario_cfg["current_conv"],
        target_conv=scenario_cfg["target_conv"],
        promo_spend=scenario_cfg["promo_spend"],
        promo_days=scenario_cfg["promo_days"],
        gross_margin_pct=scenario_cfg["gross_margin_pct"] / 100.0,
        discount_pct=scenario_cfg["discount_pct"]
    )

    zone = config.get("scenario", {}).get("zone", "East")
    add_heading(slide, f"Scenario Analysis: Promo Impact on {zone} Conversion Recovery",
                top_inches=0.3, size_pt=32)

    # Three-column scenario (dynamically calculated)
    scenarios = [
        {
            "title": "CURRENT STATE",
            "subtitle": f"{zone} Zone, Week 1",
            "conv": f"{roi['current_conv']:.1f}%",
            "weekly_offtake": f"₹{roi['current_weekly']:.1f} L",
            "status": "🔴 URGENT",
            "color": theme.ACCENT_RED,
        },
        {
            "title": "WITH PROMO",
            "subtitle": f"₹{roi['promo_spend']:.0f}L spend, {scenario_cfg['discount_pct']:.0f}% discount",
            "conv": f"{roi['mid_conv']:.1f}%",
            "weekly_offtake": f"₹{roi['promo_weekly']:.1f} L (+{roi['uplift_pct']:.0f}%)",
            "status": "🟡 IMPROVING",
            "color": theme.ACCENT_ORANGE,
        },
        {
            "title": "TARGET STATE",
            "subtitle": f"By Day {scenario_cfg['promo_days']}",
            "conv": f"{roi['target_conv']:.0f}%+",
            "weekly_offtake": f"₹{roi['target_weekly']:.1f}+ L",
            "status": "🟢 RECOVERED",
            "color": theme.ACCENT_GREEN,
        },
    ]

    scenario_top = 1.3
    for scenario_item in scenarios:
        left = 0.5 + (scenarios.index(scenario_item) * 3.1)

        # Box
        box_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(left), Inches(scenario_top),
            Inches(2.9), Inches(3.0)
        )
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = RGBColor(30, 50, 70)
        box_shape.line.color.rgb = scenario_item["color"]
        box_shape.line.width = Pt(3)

        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(scenario_top + 0.15), Inches(2.7), Inches(0.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = scenario_item["title"]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = scenario_item["color"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(scenario_top + 0.55), Inches(2.7), Inches(0.45))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = scenario_item["subtitle"]
        p.font.size = Pt(8)
        p.font.color.rgb = theme.LIGHT_GREY
        p.alignment = PP_ALIGN.CENTER

        # Conversion
        conv_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(scenario_top + 1.15), Inches(2.7), Inches(0.6))
        tf = conv_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Conversion\n{scenario_item['conv']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = scenario_item["color"]
        p.alignment = PP_ALIGN.CENTER

        # Offtake
        offtake_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(scenario_top + 1.85), Inches(2.7), Inches(0.5))
        tf = offtake_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Weekly Offtake\n{scenario_item['weekly_offtake']}"
        p.font.size = Pt(10)
        p.font.color.rgb = theme.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Status
        status_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(scenario_top + 2.45), Inches(2.7), Inches(0.4))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = scenario_item["status"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = scenario_item["color"]
        p.alignment = PP_ALIGN.CENTER

    # ROI calculation (dynamically derived from analytics engine)
    roi_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(1.5))
    tf = roi_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"SCENARIO OUTCOME: ₹{roi['promo_spend']:.0f}L promo spend → cumulative offtake uplift ₹{roi['net_uplift']:.1f} Cr over {scenario_cfg['promo_days']} days → ROI {roi['roi_multiple']:.1f}x → Release ₹{roi['net_uplift'] * 0.30:.1f} Cr working capital from {zone} inventory clearance."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = theme.ACCENT_GREEN

    p = tf.add_paragraph()
    p.text = "Assumption: Promo valid Sep 18–Oct 8 (21 days); daily OSA + weekly conversion tracking confirms lift; no secondary loading until offtake reaches 60%+ target."
    p.font.size = Pt(10)
    p.font.color.rgb = theme.LIGHT_GREY
    p.space_before = Pt(6)

    add_footer(slide, "Scenario is parameterizable: Change discount, budget, days, target zone, or lift assumptions and re-run",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_13_execution_roadmap(prs, config, theme):
    """Slide 13: 30-Day Execution Roadmap (phased)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "30-Day Execution Roadmap: Phased Recovery Sprint", top_inches=0.3, size_pt=32)

    # Four weeks
    weeks = [
        {
            "week": "Week 1\n(Sep 5-11)",
            "theme": "DISCOVERY",
            "actions": ["Store audits (5 locations)", "Root cause mapping", "Stock data collection"],
            "owner": "East ZSM",
            "success": "Root causes identified",
            "color": theme.ACCENT_RED,
        },
        {
            "week": "Week 2\n(Sep 12-18)",
            "theme": "PREPARATION",
            "actions": ["Promo design", "Collateral print", "Team training"],
            "owner": "Trade Ops",
            "success": "Assets in 10 stores",
            "color": theme.ACCENT_ORANGE,
        },
        {
            "week": "Week 3\n(Sep 19-25)",
            "theme": "EXECUTION",
            "actions": ["Promo launch", "Daily OSA tracking", "Load adjustments"],
            "owner": "NKAM Reliance",
            "success": "Conv 45%→60%",
            "color": theme.TEAL,
        },
        {
            "week": "Week 4\n(Sep 26-Oct 2)",
            "theme": "CONSOLIDATION",
            "actions": ["Sustain momentum", "Assess extension", "Plan Oct loading"],
            "owner": "East ZSM",
            "success": "Conv ≥70%",
            "color": theme.ACCENT_GREEN,
        },
    ]

    week_left = 0.4
    for week_item in weeks:
        left = week_left + (weeks.index(week_item) * 2.2)

        # Week box
        week_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(left), Inches(1.3),
            Inches(2.0), Inches(4.5)
        )
        week_shape.fill.solid()
        week_shape.fill.fore_color.rgb = RGBColor(25, 45, 70)
        week_shape.line.color.rgb = week_item["color"]
        week_shape.line.width = Pt(2)

        # Week label
        week_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(1.5), Inches(1.8), Inches(0.6))
        tf = week_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = week_item["week"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = week_item["color"]
        p.alignment = PP_ALIGN.CENTER

        # Theme
        theme_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.2), Inches(1.8), Inches(0.3))
        tf = theme_box.text_frame
        p = tf.paragraphs[0]
        p.text = week_item["theme"]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = week_item["color"]
        p.alignment = PP_ALIGN.CENTER

        # Actions
        actions_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.6), Inches(1.8), Inches(1.5))
        tf = actions_box.text_frame
        tf.word_wrap = True
        for action in week_item["actions"]:
            p = tf.add_paragraph() if actions_box != "" else tf.paragraphs[0]
            p.text = "• " + action
            p.font.size = Pt(8)
            p.font.color.rgb = theme.WHITE

        # Owner
        owner_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(4.2), Inches(1.8), Inches(0.4))
        tf = owner_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Owner:\n{week_item['owner']}"
        p.font.size = Pt(8)
        p.font.color.rgb = theme.LIGHT_GREY
        p.alignment = PP_ALIGN.CENTER

        # Success
        success_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(4.75), Inches(1.8), Inches(0.5))
        tf = success_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = week_item["success"]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = week_item["color"]
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "Milestone: End of Week 1 must unlock root cause (why 45% conversion) so Weeks 2-4 address right levers",
               size_pt=10, color=theme.MEDIUM_GREY)

def slide_14_action_register(prs, config, theme):
    """Slide 14: Action Register with Live Status Tracking."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "30-Day Action Register: Live Tracking & Accountability", top_inches=0.3, size_pt=32)

    # Action register table
    actions = [
        {
            "priority": "P0",
            "owner": "East ZSM",
            "action": "Complete store audits (3 Reliance)",
            "target": "08-Sep",
            "metric": "Root causes identified",
            "status": "IN PROGRESS (65%)",
            "status_color": theme.MEDIUM_GREY,
        },
        {
            "priority": "P0",
            "owner": "NKAM Reliance",
            "action": "Halt non-hero NPI loading",
            "target": "10-Sep",
            "metric": "Zero secondary orders",
            "status": "COMPLETED",
            "status_color": theme.ACCENT_GREEN,
        },
        {
            "priority": "P1",
            "owner": "Trade Ops",
            "action": "Deploy ₹30L promo campaign",
            "target": "18-Sep",
            "metric": "Conv 45%→55%+",
            "status": "PENDING (0%)",
            "status_color": theme.ACCENT_ORANGE,
        },
        {
            "priority": "P2",
            "owner": "Field Auditor",
            "action": "Daily OSA + weekly reporting",
            "target": "19-Sep",
            "metric": "Automated Monday flash",
            "status": "SCHEDULED",
            "status_color": theme.MEDIUM_GREY,
        },
    ]

    reg_top = 1.3
    for action in actions:
        # Determine priority color
        if action["priority"] == "P0":
            priority_color = theme.ACCENT_RED
        elif action["priority"] == "P1":
            priority_color = theme.ACCENT_ORANGE
        else:
            priority_color = theme.MEDIUM_GREY

        # Row box
        row_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.4), Inches(reg_top),
            Inches(9.2), Inches(0.65)
        )
        row_shape.fill.solid()
        row_shape.fill.fore_color.rgb = RGBColor(30, 50, 70)
        row_shape.line.color.rgb = priority_color
        row_shape.line.width = Pt(1)

        # Priority
        pri_box = slide.shapes.add_textbox(Inches(0.5), Inches(reg_top + 0.08), Inches(0.6), Inches(0.5))
        tf = pri_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["priority"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = priority_color
        p.alignment = PP_ALIGN.CENTER

        # Owner
        owner_box = slide.shapes.add_textbox(Inches(1.2), Inches(reg_top + 0.08), Inches(1.3), Inches(0.5))
        tf = owner_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["owner"]
        p.font.size = Pt(9)
        p.font.color.rgb = theme.WHITE

        # Action
        action_box = slide.shapes.add_textbox(Inches(2.6), Inches(reg_top + 0.08), Inches(2.5), Inches(0.5))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["action"]
        p.font.size = Pt(9)
        p.font.color.rgb = theme.WHITE

        # Target date
        date_box = slide.shapes.add_textbox(Inches(5.2), Inches(reg_top + 0.08), Inches(0.8), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["target"]
        p.font.size = Pt(8)
        p.font.color.rgb = theme.LIGHT_GREY
        p.alignment = PP_ALIGN.CENTER

        # Metric
        metric_box = slide.shapes.add_textbox(Inches(6.1), Inches(reg_top + 0.08), Inches(1.5), Inches(0.5))
        tf = metric_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["metric"]
        p.font.size = Pt(8)
        p.font.color.rgb = theme.LIGHT_GREY

        # Status
        status_box = slide.shapes.add_textbox(Inches(7.7), Inches(reg_top + 0.08), Inches(1.5), Inches(0.5))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = action["status"]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = action["status_color"]
        p.alignment = PP_ALIGN.RIGHT

        reg_top += 0.75

    # Footer with governance
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.0))
    tf = footer_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📌 GOVERNANCE CADENCE: Weekly Monday 2:00 PM | Any P0 slip past target date flags to VP Modern Trade"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = theme.ACCENT_ORANGE

    add_footer(slide, "Status updated every Friday EOD | Live sheet shared with leadership (read-only access)",
               size_pt=9, color=theme.MEDIUM_GREY)

def slide_15_closing(prs, config, theme):
    """Slide 15: Closing Slide (next review)."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    clear_slide(slide)
    add_background_color(slide, theme.NAVY)

    add_heading(slide, "Next Steps: October Leadership Review", top_inches=1.5, size_pt=40)

    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(8), Inches(2.5))
    tf = timeline_box.text_frame
    tf.word_wrap = True

    timeline_items = [
        "Oct 5: Conversion update (target: East ≥60%, North ≥65%)",
        "Oct 10: Promo ROI assessment + Oct loading approval",
        "Oct 15: Full zone review (6 zones) + brand performance refresh",
        "Oct 25: Q2 planning + FY28 outlook",
    ]

    for item in timeline_items:
        p = tf.add_paragraph() if timeline_items.index(item) > 0 else tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = theme.WHITE
        p.space_before = Pt(12)

    add_footer(slide, "Sent to: VP Modern Trade | Category Heads | Zone Managers | NKAM Leadership",
               size_pt=10, color=theme.MEDIUM_GREY)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 5: MAIN BUILD FUNCTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_mt_ppt(config=None, output_path=None):
    """Generate 18-slide MT leadership deck."""

    if config is None:
        config = DEFAULT_CONFIG

    if output_path is None:
        output_path = f"/home/user/mt-dashboard/MT_{config['month']}{config['year']}_Leadership_18Slides.pptx"

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Theme
    theme = MTTheme(config.get("theme_name", "default"))

    # Build all slides
    print("Building 18-slide deck...")

    slide_01_title(prs, config, theme)
    print("✓ Slide 1: Title")

    slide_02_toc(prs, config, theme)
    print("✓ Slide 2: TOC")

    slide_03_exec_summary(prs, config, theme)
    print("✓ Slide 3: Exec Summary")

    slide_04_market_context(prs, config, theme)
    print("✓ Slide 4: Market Context")

    slide_05_primary_trend(prs, config, theme)
    print("✓ Slide 5: Primary Trend")

    slide_05b_offtake_trend(prs, config, theme)
    print("✓ Slide 5b: Offtake Trend")

    slide_05c_waterfall(prs, config, theme)
    print("✓ Slide 5c: Waterfall")

    slide_06_zone_primary(prs, config, theme)
    print("✓ Slide 6: Zone Primary")

    slide_07_risk_matrix(prs, config, theme)
    print("✓ Slide 7: Risk Matrix")

    slide_08_zone_conversion(prs, config, theme)
    print("✓ Slide 8: Zone Conversion")

    slide_09_chain_concentration(prs, config, theme)
    print("✓ Slide 9: Chain Concentration")

    slide_09b_strategy_pillars(prs, config, theme)
    print("✓ Slide 9b: Strategy Pillars")

    slide_10_category_mix(prs, config, theme)
    print("✓ Slide 10: Category Mix")

    slide_11_comparison_table(prs, config, theme)
    print("✓ Slide 11: Comparison Table")

    slide_12_scenario_analysis(prs, config, theme)
    print("✓ Slide 12: Scenario Analysis")

    slide_13_execution_roadmap(prs, config, theme)
    print("✓ Slide 13: Execution Roadmap")

    slide_14_action_register(prs, config, theme)
    print("✓ Slide 14: Action Register")

    slide_15_closing(prs, config, theme)
    print("✓ Slide 15: Closing")

    # Save
    prs.save(output_path)
    print(f"\n✅ Deck saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_path

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SECTION 6: CLI ENTRY POINT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Build parameterized 18-slide MT leadership deck"
    )
    parser.add_argument("--month", default="September", help="Month name")
    parser.add_argument("--year", type=int, default=2026, help="Year")
    parser.add_argument("--theme", default="default", help="Theme name (default, teal, etc)")
    parser.add_argument("--output", default=None, help="Output file path")

    args = parser.parse_args()

    config = DEFAULT_CONFIG.copy()
    config["month"] = args.month
    config["year"] = args.year
    config["theme_name"] = args.theme

    build_mt_ppt(config, args.output)
