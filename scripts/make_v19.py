"""Build V19 from V18 — two categories of fixes:

1. DATA BUGS: Insight rows 03 and 04 on every zone slide still carry pre-V18
   values (old conversion % and old gap ₹Cr). Slide 1 shape[13] (Sales Flow
   sub-label) also still shows 76.8% / ₹10.92 Cr from V13.

2. HUMANIZATION: Every zone slide has the same Management Priority text, the
   same formulaic numbered-insight structure ("Zone delivers ₹X.XXCr offtake
   (XXXK units), X.X% national."), and a pipe-delimited evidence box that reads
   like a template. This pass rewrites those to natural analyst language that
   varies by zone and reflects actual zone dynamics.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
import zipfile

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv18.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv19.pptx"

EMU_IN = 914400
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ── helpers ───────────────────────────────────────────────────────────────────

def _runs(shape):
    """Return all a:r elements across all paragraphs in a text frame."""
    if not shape.has_text_frame:
        return []
    return shape.text_frame._txBody.findall(f".//{NS_A}r")


def patch_shape(shape, new_text: str, label: str = ""):
    """Replace text in the first run of a shape, preserving all formatting."""
    runs = _runs(shape)
    if not runs:
        print(f"  [SKIP no-runs] {label!r}")
        return False
    t_el = runs[0].find(f"{NS_A}t")
    if t_el is None:
        print(f"  [SKIP no-t] {label!r}")
        return False
    old = t_el.text or ""
    # Clear subsequent runs so they don't append garbage after the first run
    for r in runs[1:]:
        t2 = r.find(f"{NS_A}t")
        if t2 is not None:
            t2.text = ""
    t_el.text = new_text
    print(f"  [{label}] {old[:60]!r} → {new_text[:60]!r}")
    return True


def find_by_text(slide, needle: str, partial=True):
    """Return first shape whose text contains needle (or exact match)."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if partial and needle in txt:
                return sh
            elif not partial and txt.strip() == needle.strip():
                return sh
    return None


def audit_overflows(prs):
    slide_h = prs.slide_height / EMU_IN
    issues = []
    for si, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            bot = (sh.top + sh.height) / EMU_IN
            if bot > slide_h + 0.015:
                txt = sh.text_frame.text[:40] if sh.has_text_frame else ""
                issues.append(f"  S{si+1} '{txt}' bot={bot:.3f}\"")
    if issues:
        print(f"[OVERFLOW] {len(issues)} issue(s):")
        for x in issues:
            print(x)
    else:
        print("[OVERFLOW] Clean.")


def validate_xml(path):
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            if name.endswith(".xml") or name.endswith(".rels"):
                try:
                    etree.fromstring(z.read(name))
                except etree.XMLSyntaxError as e:
                    print(f"  XML ERROR {name}: {e}")
                    return False
    print("[XML] All valid.")
    return True


# ── zone config ──────────────────────────────────────────────────────────────

ZONES = [
    dict(
        name="West", si=4,
        # ── headline & management priority ──
        headline="West slips to #2 — Reliance is the one gap between here and the top",
        mgmt_priority=(
            "Reliance West 54.5% conv: map top-20 EAN-store failures by 30-Aug  •  "
            "DMart 94.2% sustains — maintain hero-EAN OSA to protect it  •  "
            "Hold Wellness Forever loading until Aug offtake confirms timing"
        ),
        # ── data-bug fixes: insight 03 and 04 ──
        old_conv_text="Flow conversion is 82.3% for July.",
        new_conv_text="84.5% flow conversion — second-best nationally. Only South-1 at 89.5% does better.",
        old_gap_text="The same-period flow gap is ₹1.78 Cr.",
        new_gap_text="Gap is ₹1.56 Cr, almost entirely Reliance. DMart's 94% leaves almost nothing on the table.",
        # ── humanized insight 01, 02, 05 ──
        old_ins01="West delivers ₹8.49Cr offtake (497K units), 20.9% national.",
        new_ins01="West held ₹8.49 Cr in July — 497K units — keeping close to #1 despite South-1's rise.",
        old_ins02="Primary +99.0% YoY: ₹10.05Cr (597K units) vs offtake ₹8.49Cr (497K).",
        new_ins02="Primary doubled YoY at ₹10.05 Cr. Most of the 597K units billed are already sold — DMart proves it.",
        old_ins05="West ranks #2 by national offtake.",
        new_ins05="#2 nationally. Close Reliance West and the ranking re-opens.",
        # ── evidence box ──
        old_ev="Off: ₹8.49Cr(497K) | DMart+Reliance+WF top 3 | 84.5% conv | ₹1.56Cr gap",
        new_ev=(
            "₹8.49 Cr from West in July, 497K units. DMart at 94% leads the zone; "
            "Reliance at 55% accounts for nearly all of the ₹1.56 Cr gap. "
            "Wellness Forever at 135% is a timing effect — hold loading."
        ),
    ),
    dict(
        name="South-1", si=5,
        headline="South-1 takes the national lead — Karnataka at ₹4.10 Cr is the one thing to protect",
        mgmt_priority=(
            "Karnataka anchors 50% of zone value — protect it before pushing Tamil Nadu  •  "
            "Apollo at 81%: fix top-5 declining EAN-store pairs before month-end  •  "
            "Lulu timing vs supply ambiguous — flag for Aug data check"
        ),
        old_conv_text="Flow conversion is 83.6% for July.",
        new_conv_text="89.5% flow — best converting zone nationally. Karnataka absorbs half of what's loaded here.",
        old_gap_text="The same-period flow gap is ₹1.61 Cr.",
        new_gap_text="Smallest zone gap nationally at ₹1.03 Cr. Entirely recoverable with routine store-level fixes.",
        old_ins01="South-1 delivers ₹8.77Cr offtake (461K units), 21.6% national.",
        new_ins01="South-1 is now the #1 zone at ₹8.77 Cr — overtook West on the back of Apollo's 188% YoY spike.",
        old_ins02="Primary +63.6% YoY: ₹9.8Cr (581K units) vs offtake ₹8.77Cr (461K).",
        new_ins02="Primary ₹9.80 Cr grew 63.6% YoY with the cleanest conversion in the country at 89.5%.",
        old_ins05="South-1 ranks #1 by national offtake.",
        new_ins05="#1 nationally for July — first time South-1 has led. Lose Karnataka and this ranking reverses.",
        old_ev="Off: ₹8.77Cr(461K) | Apollo+DMart+Lulu top 3 | 89.5% conv | ₹1.03Cr gap",
        new_ev=(
            "South-1 led all zones in July at ₹8.77 Cr, 461K units — 89.5% conversion, best in the country. "
            "Karnataka at ₹4.10 Cr is the load-bearing pillar. Apollo dipped 3pp below its benchmark; "
            "fix EAN-store gaps before it widens."
        ),
    ),
    dict(
        name="North", si=6,
        headline="North billed the most nationally — but 30% of that primary is sitting unsold in stores",
        mgmt_priority=(
            "Reliance North 44.9% conv: check Aug offtake by 05-Aug before any re-loading decision  •  "
            "Apollo at 98.3% confirms genuine demand — treat it as the zone benchmark  •  "
            "Delhi NCR + Rajasthan + Punjab = 74% of zone; all need weekly hero-EAN OSA tracking"
        ),
        old_conv_text="Flow conversion is 58.5% for July.",
        new_conv_text="69.6% conversion — second-weakest zone. Reliance at 44.9% is the single driver of underperformance.",
        old_gap_text="The same-period flow gap is ₹4.97 Cr.",
        new_gap_text="₹3.63 Cr gap. If Reliance were at Apollo-level (98%), the gap would shrink to under ₹0.50 Cr.",
        old_ins01="North delivers ₹8.32Cr offtake (442K units), 20.5% national.",
        new_ins01="North jumped from ₹6.99 to ₹8.32 Cr — 134% YoY primary load — but conversion tells a different story.",
        old_ins02="Primary +134.3% YoY: ₹11.95Cr (649K units) vs offtake ₹8.32Cr (442K).",
        new_ins02="Primary ₹11.95 Cr is the heaviest zone load nationally. At 69.6% flow, ₹3.63 Cr is sitting unsold.",
        old_ins05="North ranks #3 by national offtake.",
        new_ins05="#3 nationally — but primary load ranks #1. Conversion gap is the only thing holding North back.",
        old_ev="Off: ₹8.32Cr(442K) | Rel+DMart+Apollo top 3 | 69.6% conv | ₹3.63Cr gap",
        new_ev=(
            "North billed ₹11.95 Cr primary but only ₹8.32 Cr (442K units) sold through. "
            "Reliance holds ₹5.35 Cr of that load and converts 44.9% — roughly ₹2.95 Cr of stock is in-store. "
            "Apollo at 98.3% proves demand is real when OSA is clean."
        ),
    ),
    dict(
        name="South-2", si=7,
        headline="South-2 DMart is the largest single-account recovery play this month",
        mgmt_priority=(
            "DMart S-2 at 45% conv — run store-level audit by 28-Aug to identify top blocked EANs  •  "
            "Apollo >100%: opening-stock timing effect; quarantine primary comparison until reconciled  •  "
            "Reliance at 82.5%: healthy — monitor only, no action needed"
        ),
        old_conv_text="Flow conversion is 71.3% for July.",
        new_conv_text="77.0% zone conversion — solid, but DMart at 45% pulls the average down single-handedly.",
        old_gap_text="The same-period flow gap is ₹1.98 Cr.",
        new_gap_text="₹1.59 Cr gap. Most of it is DMart — the ₹1.95 Cr DMart load converting at 45% leaves ₹1.07 Cr on the table.",
        old_ins01="South-2 delivers ₹5.3Cr offtake (299K units), 13.0% national.",
        new_ins01="South-2 rose from ₹4.91 to ₹5.30 Cr — DMart tripled its Jul-25 level and is now the story to resolve.",
        old_ins02="Primary +69.3% YoY: ₹6.89Cr (343K units) vs offtake ₹5.3Cr (299K).",
        new_ins02="Primary grew 69.3% YoY to ₹6.89 Cr; at 77.0% overall conv, the zone is steady except for the DMart gap.",
        old_ins05="South-2 ranks #4 by national offtake.",
        new_ins05="#4 nationally. Close the DMart gap and South-2 moves comfortably past East.",
        old_ev="Off: ₹5.30Cr(299K) | DMart+Apollo+Rel top 3 | 77.0% conv | ₹1.59Cr gap",
        new_ev=(
            "₹5.30 Cr offtake from South-2, 299K units. DMart grew 3× YoY but converts at 45% — "
            "pipeline issue, not demand. Apollo at 148% is an opening-stock timing effect. "
            "Run the DMart store audit before drawing any loading conclusions."
        ),
    ),
    dict(
        name="East", si=8,
        headline="East at 62% conversion — Reliance stock is in-store; moratorium holds",
        mgmt_priority=(
            "Non-hero loading moratorium stands until conversion crosses 70%  •  "
            "Validate Aug-26 Reliance East offtake by 05-Aug — if still below 70%, pull EAN-store list  •  "
            "Apollo East is NEW: protect launch availability; NPI at 10.2% is a secondary exposure to manage"
        ),
        old_conv_text="Flow conversion is 45.3% for July.",
        new_conv_text="61.9% conversion — weakest zone nationally. Reliance holds 44% of the load and converts barely half.",
        old_gap_text="The same-period flow gap is ₹4.28 Cr.",
        new_gap_text="₹2.99 Cr gap — second-largest nationally after North. Reliance East alone accounts for ~₹1.90 Cr of that.",
        old_ins01="East delivers ₹4.85Cr offtake (245K units), 11.9% national.",
        new_ins01="East moved from ₹3.55 to ₹4.85 Cr once pending chains were counted — still the weakest converting zone.",
        old_ins02="Primary +115.7% YoY: ₹7.83Cr (409K units) vs offtake ₹4.85Cr (245K).",
        new_ins02="Primary ₹7.83 Cr (115.7% YoY) is disproportionate to offtake. Close to ₹3 Cr of stock hasn't moved.",
        old_ins05="East ranks #5 by national offtake.",
        new_ins05="#5 nationally. Non-hero loading moratorium is the right call until conversion crosses 70%.",
        old_ev="Off: ₹4.85Cr(245K) | Rel+Apollo+RBC top 3 | 61.9% conv | ₹2.99Cr gap",
        new_ev=(
            "₹4.85 Cr offtake, 245K units — weakest conversion nationally at 61.9%. "
            "Reliance East holds ₹4.08 Cr primary and converts 53%: roughly ₹1.90 Cr of stock is sitting. "
            "Non-hero moratorium active. Validate Aug offtake before any re-loading."
        ),
    ),
    dict(
        name="Central", si=9,
        headline="Central opens clean at 107% flow — protect this rate before the pipeline normalises",
        mgmt_priority=(
            "DMart Central 95.3% is real demand — document EAN cadence as the national DMart benchmark  •  "
            "Reliance 51%: confirm whether this is opening stock or a true demand gap before re-ordering  •  "
            "Apollo NEW: validate opening stock first; do not load beyond confirmed sell-through rate"
        ),
        old_conv_text="Flow conversion is 78.8% for July.",
        new_conv_text=(
            "107.1% conversion — offtake exceeds primary this month. "
            "Prior channel inventory is clearing. DMart at 95% is genuine; Reliance at 51% is not."
        ),
        old_gap_text="The same-period flow gap is ₹0.57 Cr.",
        new_gap_text=(
            "₹0.19 Cr 'surplus' is a data artifact — normalise over two months before "
            "treating Central as gap-free. Reliance 51% is a watch."
        ),
        old_ins01="Central delivers ₹2.88Cr offtake (171K units), 7.1% national.",
        new_ins01="Central is new in the data; ₹2.88 Cr offtake already exceeds ₹2.69 Cr primary — prior stock clearing.",
        old_ins02="Primary +202.2% YoY: ₹2.69Cr (160K units) vs offtake ₹2.88Cr (171K).",
        new_ins02="202% YoY primary growth reflects the zone's first proper July in the MT count, not organic demand.",
        old_ins05="Central ranks #6 by national offtake.",
        new_ins05="#6 nationally, smallest zone. Set OSA and ordering benchmarks now while stores are manageable.",
        old_ev="Off: ₹2.88Cr(171K) | DMart+Rel+Apollo top 3 | 107% conv (timing surplus)",
        new_ev=(
            "₹2.88 Cr offtake, 171K units. Offtake > primary this month — prior inventory clearing. "
            "DMart 95.3% is a real and repeatable rate. Reliance 51% is suspect. "
            "Document DMart's EAN cadence as the national benchmark before loading more."
        ),
    ),
]


# ── main ─────────────────────────────────────────────────────────────────────

prs = Presentation(SRC)

# ── FIX 1: Slide 1 shape[13] — Sales Flow sub-label still shows old values ──
print("\n=== FIX 1: Slide 1 shape[13] — Sales Flow sub-label ===")
slide1 = prs.slides[0]
sh = find_by_text(slide1, "76.8% conversion")
if sh:
    patch_shape(sh, "86.5% flow conversion | ₹6.35 Cr gap", "S1:sales-flow-label")
else:
    print("  shape not found by text — already patched or moved")

# ── FIX 2: Per-zone data bugs + humanization ─────────────────────────────────
for z in ZONES:
    name, si = z["name"], z["si"]
    slide = prs.slides[si]
    print(f"\n=== {name} (slide {si+1}) ===")

    # Slide headline
    sh = find_by_text(slide, name + ":", partial=True)
    if sh and "Watch and convert" in sh.text_frame.text or \
       sh and "Urgent gap closure" in sh.text_frame.text:
        patch_shape(sh, z["headline"], f"{name}:headline")
    else:
        # Try by index 1 (always the zone headline)
        sh2 = slide.shapes[1]
        if sh2.has_text_frame and name in sh2.text_frame.text:
            patch_shape(sh2, z["headline"], f"{name}:headline")

    # Management Priority (shape after the MANAGEMENT PRIORITY label)
    sh = find_by_text(slide, "Close DMart exceptions") or \
         find_by_text(slide, "Close Apollo exceptions") or \
         find_by_text(slide, "Close Reliance exceptions")
    if sh:
        patch_shape(sh, z["mgmt_priority"], f"{name}:mgmt-priority")

    # Insight 01 — humanize
    sh = find_by_text(slide, z["old_ins01"], partial=False)
    if sh:
        patch_shape(sh, z["new_ins01"], f"{name}:ins01")

    # Insight 02 — humanize
    sh = find_by_text(slide, z["old_ins02"][:50])
    if sh:
        patch_shape(sh, z["new_ins02"], f"{name}:ins02")

    # Insight 03 — DATA BUG: stale conversion % + humanize
    sh = find_by_text(slide, z["old_conv_text"], partial=False)
    if sh:
        patch_shape(sh, z["new_conv_text"], f"{name}:ins03-conv")
    else:
        print(f"  [{name}:ins03] not found: {z['old_conv_text']!r}")

    # Insight 04 — DATA BUG: stale gap + humanize
    sh = find_by_text(slide, z["old_gap_text"], partial=False)
    if sh:
        patch_shape(sh, z["new_gap_text"], f"{name}:ins04-gap")
    else:
        print(f"  [{name}:ins04] not found: {z['old_gap_text']!r}")

    # Insight 05 — humanize
    sh = find_by_text(slide, z["old_ins05"], partial=False)
    if sh:
        patch_shape(sh, z["new_ins05"], f"{name}:ins05")

    # Evidence box (EVIDENCE section text)
    sh = find_by_text(slide, z["old_ev"], partial=False)
    if sh:
        patch_shape(sh, z["new_ev"], f"{name}:evidence")
    else:
        print(f"  [{name}:ev] not found by exact match — trying partial")
        sh = find_by_text(slide, "Off: ₹")
        if sh:
            patch_shape(sh, z["new_ev"], f"{name}:evidence(partial)")

# ── also humanize South-2 offtake display: fix ₹5.3 → ₹5.30 ─────────────────
print("\n=== South-2 offtake display fix (₹5.3 → ₹5.30) ===")
slide_s2 = prs.slides[7]
sh = find_by_text(slide_s2, "₹5.3 Cr", partial=False)
if sh:
    patch_shape(sh, "₹5.30 Cr", "S2:off-display")

# ── OVERFLOW AUDIT ────────────────────────────────────────────────────────────
print()
audit_overflows(prs)

# ── SAVE ─────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved → {DST}")

# ── XML VALIDATION ────────────────────────────────────────────────────────────
validate_xml(DST)

# ── SPOT CHECKS ───────────────────────────────────────────────────────────────
print("\n=== SPOT CHECKS ===")
p2 = Presentation(DST)

# Slide 1: sales flow label
s1 = p2.slides[0]
sh = find_by_text(s1, "86.5% flow conversion")
print(f"  S1 flow label: {'OK' if sh else 'MISSING'}")

# West conversion
sw = p2.slides[4]
sh = find_by_text(sw, "84.5% flow conversion")
print(f"  West ins03:    {'OK' if sh else 'MISSING'}")

# North gap
sn = p2.slides[6]
sh = find_by_text(sn, "₹3.63 Cr gap")
print(f"  North ins04:   {'OK' if sh else 'MISSING'}")

# East conversion
se = p2.slides[8]
sh = find_by_text(se, "61.9% conversion")
print(f"  East ins03:    {'OK' if sh else 'MISSING'}")

# Central timing surplus
sc = p2.slides[9]
sh = find_by_text(sc, "107.1% conversion")
print(f"  Central ins03: {'OK' if sh else 'MISSING'}")

# South-1 no longer shows old rank
ss1 = p2.slides[5]
sh_old = find_by_text(ss1, "ranks #2 by national")
sh_new = find_by_text(ss1, "#1 nationally")
print(f"  S1 rank:       old={'PRESENT (BAD)' if sh_old else 'gone'} | new={'OK' if sh_new else 'MISSING'}")
