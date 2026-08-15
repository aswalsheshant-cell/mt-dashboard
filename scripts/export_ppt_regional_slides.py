#!/usr/bin/env python3
"""
Export regional (zone-by-distribution, city-tier) breakdown slides to PowerPoint.

Appends new slides to an existing leadership deck:
  1. Zone Performance Matrix (share, distribution, SAH by zone)
  2. City-Tier Distribution Breakdown (distribution points, share by tier)

Usage:
    python export_ppt_regional_slides.py \\
        --data dashboard/data.js \\
        --leadership Final_MT_Offtake_May26_Leadership_slide.pptx \\
        --out Final_MT_Leadership_WITH_REGIONAL_DEPTH_v3.pptx

Requires:
    python-pptx, pandas
"""

import argparse
import json
import sys
from pathlib import Path
from dataclasses import dataclass
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("ERROR: python-pptx not installed.", file=sys.stderr)
    sys.exit(1)

try:
    import pandas as pd
except ImportError:
    print("ERROR: pandas not installed.", file=sys.stderr)
    sys.exit(1)


# ============================================================================
# Brand palette (matching rebuild_mt_offtake_ppt.py)
# ============================================================================

TEAL = RGBColor(0x2D, 0x9B, 0x7F)
WARM_BG = RGBColor(0xFA, 0xF7, 0xF2)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x33)
GROWTH_GREEN = RGBColor(0x1E, 0x8E, 0x3E)
DECLINE_RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Aptos"
FONT_FALLBACK = "Calibri"

MARGIN_LR = Inches(0.3)
MARGIN_TOP = Inches(0.25)
MARGIN_BOTTOM = Inches(0.25)

# City-tier mapping (Tier-1 metros, Tier-2 cities, Tier-3 towns)
CITY_TIER_MAP = {
    # Tier-1 (metros)
    "Delhi": 1, "NCR": 1, "Mumbai": 1, "Bangalore": 1, "Hyderabad": 1, "Chennai": 1,
    # Tier-2 (major cities)
    "Pune": 2, "Ahmedabad": 2, "Jaipur": 2, "Lucknow": 2, "Chandigarh": 2,
    "Kolkata": 2, "Indore": 2, "Kochi": 2, "Visakhapatnam": 2,
    # Default to Tier-3 for unmapped cities
}


# ============================================================================
# Data extraction helpers
# ============================================================================

def extract_zone_metrics(data: dict) -> dict:
    """Extract zone-level metrics from data.js offtake block."""
    zones = {}

    if 'offtake' not in data:
        return zones

    # Zone-level data
    for zone_entry in data['offtake'].get('by_zone', []):
        zone_name = zone_entry.get('name', 'Unknown')
        zones[zone_name] = {
            'fy25': zone_entry.get('fy25', 0),
            'fy26': zone_entry.get('fy26', 0),
            'fy27': zone_entry.get('fy27', 0),
            'yoy': zone_entry.get('yoy', 0),
            'states': [],
            'total_share': 0,
            'sah': 0,
        }

    # Add state details
    for state_entry in data['offtake'].get('by_state', []):
        state = state_entry.get('state', 'Unknown')
        zone = state_entry.get('zone', 'Unknown')
        if zone in zones:
            zones[zone]['states'].append(state)

    return zones


def extract_city_tier_metrics(data: dict) -> dict:
    """
    Extract city-tier distribution metrics from detail_records.
    Returns {tier: {dist_points, share, sah, count}}.
    """
    tier_stats = defaultdict(lambda: {
        'dist_points': 0,
        'share': 0.0,
        'sah': 0.0,
        'count': 0,
        'cities': set(),
    })

    if 'detail_records' not in data:
        return {}

    # Group by state/city to extract distinct distribution points
    state_city_groups = defaultdict(lambda: {'nsv': 0, 'count': 0})
    city_tier_map = {}

    for record in data['detail_records']:
        state = record.get('State', 'Unknown')
        city = record.get('City', state)  # Fallback to state if city missing

        # Classify city to tier
        if city not in city_tier_map:
            # Try exact match or state-level default
            city_tier_map[city] = CITY_TIER_MAP.get(city,
                                                     CITY_TIER_MAP.get(state, 3))

        tier = city_tier_map[city]
        tier_key = f"Tier-{tier}"

        # Accumulate metrics
        key = (state, city)
        state_city_groups[key]['nsv'] += record.get('NSV', 0) or 0
        state_city_groups[key]['count'] += 1

        # Track unique cities per tier
        tier_stats[tier_key]['cities'].add(city)

    # Aggregate to tier level
    for (state, city), metrics in state_city_groups.items():
        tier = city_tier_map[city]
        tier_key = f"Tier-{tier}"
        tier_stats[tier_key]['dist_points'] += 1  # Each store = 1 distribution point
        tier_stats[tier_key]['share'] += metrics['nsv']
        tier_stats[tier_key]['count'] += metrics['count']

    # Convert city sets to lists for JSON serialization
    return {
        k: {
            'dist_points': v['dist_points'],
            'share': v['share'],
            'sah': v['sah'],
            'count': v['count'],
            'cities': len(v['cities']),
        }
        for k, v in tier_stats.items()
    }


# ============================================================================
# Slide builders
# ============================================================================

def add_zone_performance_slide(prs: Presentation, zone_metrics: dict):
    """Add Zone Performance Matrix slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # blank layout (index 1)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WARM_BG

    # Title
    title_box = slide.shapes.add_textbox(MARGIN_LR, MARGIN_TOP, Inches(7), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Zone Performance Matrix"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEAL

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(MARGIN_LR, Inches(0.85), Inches(7), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Share trend × Distribution gap analysis by zone"
    subtitle_frame.paragraphs[0].font.size = Pt(14)
    subtitle_frame.paragraphs[0].font.color.rgb = DARK_TEXT

    # Create a simple 2x2 matrix table
    zones_list = sorted(zone_metrics.keys())

    # Table: 3 cols (Zone, Share YoY, FY27), 5 rows (header + 4 zones + 1 padding)
    rows, cols = len(zones_list) + 1, 3
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(6.5), Inches(4)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    for col_idx in range(cols):
        table.columns[col_idx].width = Inches(2.1)

    # Header row
    headers = ['Zone', 'YoY %', 'FY27 Value']
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        # Format text
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.size = Pt(12)

    # Data rows
    for row_idx, zone_name in enumerate(zones_list, 1):
        zone_data = zone_metrics[zone_name]

        # Zone name
        cell = table.cell(row_idx, 0)
        cell.text = zone_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
        cell.text_frame.paragraphs[0].font.bold = True

        # YoY %
        cell = table.cell(row_idx, 1)
        yoy = zone_data.get('yoy', 0)
        cell.text = f"{yoy:.1f}%"
        color = GROWTH_GREEN if yoy > 0 else DECLINE_RED
        cell.text_frame.paragraphs[0].font.color.rgb = color
        cell.text_frame.paragraphs[0].font.bold = True

        # FY27 Value
        cell = table.cell(row_idx, 2)
        fy27 = zone_data.get('fy27', 0)
        cell.text = f"₹{fy27:,.0f}L"
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    return slide


def add_city_tier_breakdown_slide(prs: Presentation, tier_metrics: dict):
    """Add City-Tier Distribution Breakdown slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # blank layout (index 1)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WARM_BG

    # Title
    title_box = slide.shapes.add_textbox(MARGIN_LR, MARGIN_TOP, Inches(7), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "City-Tier Distribution Breakdown"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEAL

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(MARGIN_LR, Inches(0.85), Inches(7), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Distribution points and market reach by city tier"
    subtitle_frame.paragraphs[0].font.size = Pt(14)
    subtitle_frame.paragraphs[0].font.color.rgb = DARK_TEXT

    # Sort tiers
    tiers = sorted(tier_metrics.keys(), key=lambda x: int(x.split('-')[1]))

    # Table: 4 cols (Tier, Dist Points, Cities, Share)
    rows, cols = len(tiers) + 1, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(6.5), Inches(3)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    for col_idx, w in enumerate([1.5, 1.8, 1.5, 1.5]):
        table.columns[col_idx].width = Inches(w)

    # Header
    headers = ['Tier', 'Distribution Points', 'Unique Cities', 'Share (₹L)']
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.size = Pt(11)

    # Data rows
    for row_idx, tier in enumerate(tiers, 1):
        tier_data = tier_metrics[tier]

        # Tier name
        cell = table.cell(row_idx, 0)
        cell.text = tier
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
        cell.text_frame.paragraphs[0].font.bold = True

        # Distribution points
        cell = table.cell(row_idx, 1)
        cell.text = str(tier_data['dist_points'])
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True

        # Unique cities
        cell = table.cell(row_idx, 2)
        cell.text = str(tier_data['cities'])
        cell.text_frame.paragraphs[0].font.size = Pt(12)

        # Share
        cell = table.cell(row_idx, 3)
        share = tier_data['share']
        cell.text = f"₹{share/100000:,.1f}L"  # Convert to Lakhs
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add key insight at bottom
    insight_box = slide.shapes.add_textbox(MARGIN_LR, Inches(4.8), Inches(6.5), Inches(0.8))
    insight_frame = insight_box.text_frame
    insight_frame.word_wrap = True
    insight_frame.text = "Key Insight: Tier-1 metros drive majority of volume. Tier-2/3 expansion opportunity: identify high-PDO, low-distribution gaps."
    insight_frame.paragraphs[0].font.size = Pt(10)
    insight_frame.paragraphs[0].font.italic = True
    insight_frame.paragraphs[0].font.color.rgb = DARK_TEXT

    return slide


# ============================================================================
# Main export function
# ============================================================================

def export_regional_slides(data_file: Path, leadership_pptx: Path, output_pptx: Path):
    """Load data, build slides, append to existing deck."""

    # Load data.js
    print(f"Loading data from {data_file}...")
    with open(data_file, 'r') as f:
        content = f.read()
        start = content.find('{')
        end = content.rfind('}') + 1
        data = json.loads(content[start:end])

    # Extract metrics
    print("Extracting zone metrics...")
    zone_metrics = extract_zone_metrics(data)

    print("Extracting city-tier metrics...")
    tier_metrics = extract_city_tier_metrics(data)

    # Load existing presentation
    print(f"Loading leadership deck from {leadership_pptx}...")
    prs = Presentation(str(leadership_pptx))

    # Add new slides
    print("Adding Zone Performance slide...")
    add_zone_performance_slide(prs, zone_metrics)

    print("Adding City-Tier Breakdown slide...")
    add_city_tier_breakdown_slide(prs, tier_metrics)

    # Save
    print(f"Saving to {output_pptx}...")
    prs.save(str(output_pptx))

    print(f"✓ Regional depth slides exported successfully to {output_pptx}")
    print(f"  - {len(zone_metrics)} zones analyzed")
    print(f"  - {len(tier_metrics)} city-tiers analyzed")


# ============================================================================
# CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Export zone-by-distribution and city-tier breakdown slides to PowerPoint."
    )
    parser.add_argument('--data', type=Path, default=Path('dashboard/data.js'),
                        help='Path to data.js (default: dashboard/data.js)')
    parser.add_argument('--leadership', type=Path, required=True,
                        help='Path to existing leadership PowerPoint deck')
    parser.add_argument('--out', type=Path, required=True,
                        help='Output PowerPoint file path')

    args = parser.parse_args()

    if not args.data.exists():
        print(f"ERROR: data file not found: {args.data}", file=sys.stderr)
        sys.exit(1)

    if not args.leadership.exists():
        print(f"ERROR: leadership deck not found: {args.leadership}", file=sys.stderr)
        sys.exit(1)

    export_regional_slides(args.data, args.leadership, args.out)


if __name__ == '__main__':
    main()
