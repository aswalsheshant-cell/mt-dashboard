"""
Fix overlapping text in slides 1-5 by adjusting text frame heights and spacing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def fix_overlapping_text(input_file, output_file):
    """Fix overlapping text elements in presentation."""
    prs = Presentation(input_file)

    # Process slides 1-5 (indices 0-4)
    for slide_idx in range(min(5, len(prs.slides))):
        slide = prs.slides[slide_idx]
        print(f"\nFixing Slide {slide_idx + 1}...")

        # Collect all text shapes with their positions
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text.strip():
                text_shapes.append(shape)

        # For KPI card pattern: label (small), value (large), change (small)
        # Sort by vertical position to identify cards
        for shape in text_shapes:
            if not hasattr(shape, 'text_frame'):
                continue

            text_frame = shape.text_frame

            # Reduce height of label/change text (9-9.5pt) to prevent overlap
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    font_size = run.font.size
                    if font_size and font_size <= Pt(10):
                        # These are labels or small text - reduce frame margin
                        text_frame.margin_bottom = Inches(0.01)
                        text_frame.margin_top = Inches(0.01)

            # Adjust text frame word wrap and fit
            text_frame.word_wrap = True
            text_frame.auto_size = None  # Don't auto-shrink

        # Post-process: check for remaining overlaps and adjust vertically
        overlaps_found = []
        for i, s1 in enumerate(text_shapes):
            for j, s2 in enumerate(text_shapes[i+1:], i+1):
                if not (hasattr(s1, 'left') and hasattr(s2, 'left')):
                    continue

                left1 = s1.left / 914400
                top1 = s1.top / 914400
                width1 = s1.width / 914400
                height1 = s1.height / 914400

                left2 = s2.left / 914400
                top2 = s2.top / 914400
                width2 = s2.width / 914400
                height2 = s2.height / 914400

                # Check for overlap
                if (left1 < left2 + width2 and left1 + width1 > left2 and
                    top1 < top2 + height2 and top1 + height1 > top2):
                    # If s2 is below s1, push s2 down
                    if top2 > top1:
                        new_top = top1 + height1 + 0.05  # Add 0.05" gap
                        s2.top = int(new_top * 914400)
                        overlaps_found.append((s1.text[:30], s2.text[:30]))

        if overlaps_found:
            print(f"  Fixed {len(overlaps_found)} overlaps:")
            for t1, t2 in overlaps_found:
                print(f"    - '{t1}' → '{t2}'")

    # Save fixed presentation
    prs.save(output_file)
    print(f"\n✓ Saved fixed presentation: {output_file}")


if __name__ == "__main__":
    input_file = "/home/user/mt-dashboard/MT_July26_Final_UPDATED_with_All3_Insights_v1.pptx"
    output_file = "/home/user/mt-dashboard/MT_July26_Final_FIXED_No_Overlap.pptx"

    fix_overlapping_text(input_file, output_file)
