#!/usr/bin/env python3
"""Rasterise a .pptx to PNGs straight from the shape tree, for visual QA.

LibreOffice is not usable in the build container, so this walks the python-pptx
object model and draws it with Pillow. It is a layout proof, not a pixel-perfect
renderer: fonts fall back to Liberation Sans and native charts are drawn from
their plot data. Good enough to catch clipping, collisions, bad wraps and
colour mistakes before a deck goes out.

    python scripts/preview_deck.py deck.pptx out_dir [--dpi 100] [--slides 1,2]
"""
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
FONT_DIR = "/usr/share/fonts/truetype/liberation"
REG = os.path.join(FONT_DIR, "LiberationSans-Regular.ttf")
BLD = os.path.join(FONT_DIR, "LiberationSans-Bold.ttf")
_CACHE = {}


def font(size_pt, bold, dpi):
    px = max(6, int(round(size_pt * dpi / 72.0)))
    key = (px, bold)
    if key not in _CACHE:
        _CACHE[key] = ImageFont.truetype(BLD if bold else REG, px)
    return _CACHE[key]


def rgb_of(color, default=None):
    try:
        if color and color.type is not None and color.rgb is not None:
            return tuple(color.rgb)
    except (AttributeError, ValueError, TypeError):
        pass
    return default


def fill_of(shape, default=None):
    try:
        f = shape.fill
        if f.type is None:
            return default
        if str(f.type).startswith("BACKGROUND"):
            return None
        return rgb_of(f.fore_color, default)
    except (AttributeError, ValueError, TypeError, NotImplementedError):
        return default


def line_of(shape):
    try:
        ln = shape.line
        if ln.fill.type is not None and str(ln.fill.type).startswith("BACKGROUND"):
            return None, 0
        c = rgb_of(ln.color)
        w = ln.width.pt if ln.width else 0.75
        return (c, w) if c else (None, 0)
    except (AttributeError, ValueError, TypeError, NotImplementedError):
        return None, 0


def wrap(draw, txt, fnt, max_w):
    """Greedy wrap that mirrors PowerPoint's word wrapping closely enough."""
    out = []
    for hard in txt.split("\n"):
        words, line = hard.split(" "), ""
        for w in words:
            trial = (line + " " + w).strip()
            if draw.textlength(trial, font=fnt) <= max_w or not line:
                line = trial
            else:
                out.append(line)
                line = w
        out.append(line)
    return out


def draw_text_frame(draw, sh, box, dpi, warn):
    x0, y0, w, h = box
    tf = sh.text_frame
    try:
        ml = tf.margin_left.inches if tf.margin_left is not None else 0.1
        mr = tf.margin_right.inches if tf.margin_right is not None else 0.1
        mt = tf.margin_top.inches if tf.margin_top is not None else 0.05
    except (AttributeError, ValueError):
        ml = mr = 0.1
        mt = 0.05
    ix = x0 + ml * dpi
    iw = max(4, w - (ml + mr) * dpi)

    lines = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            lines.append(([], 0, p.alignment))
            continue
        size = runs[0].font.size.pt if runs[0].font.size else 11
        bold = bool(runs[0].font.bold)
        full = "".join(r.text for r in runs)
        fnt = font(size, bold, dpi)
        ls = getattr(p, "line_spacing", None)
        mult = ls if isinstance(ls, float) else 1.0
        wrapped = wrap(draw, full, fnt, iw) if tf.word_wrap is not False else [full]
        for ln in wrapped:
            lines.append(([(ln, runs)], size * dpi / 72.0 * 1.20 * mult, p.alignment))

    total = sum(lh for _, lh, _ in lines)
    anchor = str(tf.vertical_anchor or "")
    if "MIDDLE" in anchor:
        cy = y0 + (h - total) / 2
    elif "BOTTOM" in anchor:
        cy = y0 + h - total
    else:
        cy = y0 + mt * dpi

    if total > h + 2:
        warn.append(f"text taller than its box ({total/dpi:.2f}in vs "
                    f"{h/dpi:.2f}in): {tf.text[:46]!r}")

    for parts, lh, align in lines:
        if not parts:
            cy += lh
            continue
        seg, runs = parts[0]
        size = runs[0].font.size.pt if runs[0].font.size else 11
        bold = bool(runs[0].font.bold)
        col = rgb_of(runs[0].font.color, (31, 41, 51))
        fnt = font(size, bold, dpi)
        tw = draw.textlength(seg, font=fnt)
        a = str(align or "")
        if "CENTER" in a:
            tx = ix + (iw - tw) / 2
        elif "RIGHT" in a:
            tx = ix + iw - tw
        else:
            tx = ix
        draw.text((tx, cy), seg, font=fnt, fill=col)
        cy += lh


def draw_chart(draw, gf, box, dpi):
    """Approximate a native chart from its plot data — enough to judge layout."""
    x0, y0, w, h = box
    draw.rectangle([x0, y0, x0 + w, y0 + h], outline=(220, 216, 208))
    try:
        ch = gf.chart
        ctype = str(ch.chart_type)
    except (AttributeError, ValueError):
        return
    if "DOUGHNUT" in ctype or "PIE" in ctype:
        plot = ch.plots[0]
        vals = list(plot.series[0].values)
        cols = []
        for pt in plot.series[0].points:
            cols.append(rgb_of(pt.format.fill.fore_color, (150, 150, 150)))
        tot = sum(vals) or 1
        d = min(w, h) * 0.86
        # legend sits to the right, so the ring is centred in the left portion
        cx, cy = x0 + w * 0.31, y0 + h / 2
        bb = [cx - d / 2, cy - d / 2, cx + d / 2, cy + d / 2]
        ang = -90
        for v, c in zip(vals, cols):
            sweep = 360 * v / tot
            draw.pieslice(bb, ang, ang + sweep, fill=c, outline=(255, 255, 255))
            ang += sweep
        hole = d * 0.62
        draw.ellipse([cx - hole / 2, cy - hole / 2, cx + hole / 2, cy + hole / 2],
                     fill=(244, 241, 236))
        ly = y0 + h / 2 - len(vals) * 9
        cats = [str(c) for c in ch.plots[0].categories]
        for cat, c in zip(cats, cols):
            lx = x0 + w * 0.64
            draw.rectangle([lx, ly + 3, lx + 9, ly + 12], fill=c)
            draw.text((lx + 14, ly), cat, font=font(8.5, False, dpi),
                      fill=(31, 41, 51))
            ly += 18


def render(path, out_dir, dpi=100, only=None):
    prs = Presentation(path)
    W = int(prs.slide_width / EMU_IN * dpi)
    H = int(prs.slide_height / EMU_IN * dpi)
    os.makedirs(out_dir, exist_ok=True)
    made, warnings = [], []

    for idx, slide in enumerate(prs.slides, 1):
        if only and idx not in only:
            continue
        bg = (255, 255, 255)
        try:
            b = rgb_of(slide.background.fill.fore_color)
            if b:
                bg = b
        except (AttributeError, ValueError, TypeError, NotImplementedError):
            pass
        img = Image.new("RGB", (W, H), bg)
        draw = ImageDraw.Draw(img)
        warn = []

        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            x0 = sh.left / EMU_IN * dpi
            y0 = sh.top / EMU_IN * dpi
            w = (sh.width or 0) / EMU_IN * dpi
            h = (sh.height or 0) / EMU_IN * dpi
            box = (x0, y0, w, h)
            st = str(sh.shape_type or "")
            tag = sh._element.tag.split("}")[-1]

            if sh.has_chart if hasattr(sh, "has_chart") else False:
                draw_chart(draw, sh, box, dpi)
                continue

            if "TABLE" in st or (hasattr(sh, "has_table") and sh.has_table):
                tbl = sh.table
                ty = y0
                for r_i, row in enumerate(tbl.rows):
                    rh = row.height / EMU_IN * dpi
                    tx = x0
                    for c_i, col in enumerate(tbl.columns):
                        cw = col.width / EMU_IN * dpi
                        cell = tbl.cell(r_i, c_i)
                        cf = rgb_of(cell.fill.fore_color, (255, 255, 255))
                        draw.rectangle([tx, ty, tx + cw, ty + rh], fill=cf)
                        p = cell.text_frame.paragraphs[0]
                        if p.runs:
                            r = p.runs[0]
                            f = font(r.font.size.pt if r.font.size else 9,
                                     bool(r.font.bold), dpi)
                            c = rgb_of(r.font.color, (31, 41, 51))
                            tw = draw.textlength(cell.text, font=f)
                            a = str(p.alignment or "")
                            if "CENTER" in a:
                                cx = tx + (cw - tw) / 2
                            elif "RIGHT" in a:
                                cx = tx + cw - tw - 6
                            else:
                                cx = tx + 6
                            fh = f.size if hasattr(f, "size") else 10
                            draw.text((cx, ty + (rh - fh) / 2), cell.text,
                                      font=f, fill=c)
                        tx += cw
                    ty += rh
                continue

            # connectors are <p:cxnSp> and report shape_type LINE, not CONNECTOR
            if tag == "cxnSp":
                c, lw = line_of(sh)
                if c:
                    xf = sh._element.spPr.xfrm
                    flip_h = xf is not None and xf.get("flipH") == "1"
                    flip_v = xf is not None and xf.get("flipV") == "1"
                    ax, bx = (x0 + w, x0) if flip_h else (x0, x0 + w)
                    ay, by = (y0 + h, y0) if flip_v else (y0, y0 + h)
                    draw.line([ax, ay, bx, by], fill=c,
                              width=max(1, int(lw * dpi / 72.0)))
                continue

            f = fill_of(sh)
            lc, lw = line_of(sh)
            if f or lc:
                try:
                    ast = str(sh.auto_shape_type or "")
                except (AttributeError, ValueError):
                    ast = ""
                st = st + " " + ast
                if "OVAL" in st:
                    draw.ellipse([x0, y0, x0 + w, y0 + h], fill=f, outline=lc,
                                 width=max(1, int(lw * dpi / 72.0)) if lc else 1)
                else:
                    rad = 0
                    if "ROUNDED" in st:
                        rad = int(min(w, h) * 0.13)
                    if rad > 1:
                        draw.rounded_rectangle([x0, y0, x0 + w, y0 + h], rad,
                                               fill=f, outline=lc,
                                               width=max(1, int(lw * dpi / 72.0)) if lc else 1)
                    else:
                        draw.rectangle([x0, y0, x0 + w, y0 + h], fill=f, outline=lc,
                                       width=max(1, int(lw * dpi / 72.0)) if lc else 1)

            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_text_frame(draw, sh, box, dpi, warn)

        p = os.path.join(out_dir, f"slide{idx:02d}.png")
        img.save(p)
        made.append(p)
        for wmsg in warn:
            warnings.append(f"S{idx}: {wmsg}")

    return made, warnings


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(2)
    deck, out = sys.argv[1], sys.argv[2]
    dpi = 100
    only = None
    if "--dpi" in sys.argv:
        dpi = int(sys.argv[sys.argv.index("--dpi") + 1])
    if "--slides" in sys.argv:
        only = {int(v) for v in sys.argv[sys.argv.index("--slides") + 1].split(",")}
    made, warns = render(deck, out, dpi, only)
    print(f"rendered {len(made)} slide(s) -> {out}")
    if warns:
        print(f"{len(warns)} text-fit warning(s):")
        for w in warns:
            print("  -", w)
    else:
        print("no text-fit warnings")


if __name__ == "__main__":
    main()
