"""Geometry QA for the reworked deck.

Flags the defect classes found in the source pack:
  - shapes running past the page edge
  - content colliding with the footer rail or the source line
  - cards whose children spill out of them
"""
import os
import sys
from pptx import Presentation

EMU = 914400
PH, PW = 13.333, 7.5
FOOT_Y, SRC_Y = 12.44, 13.00

DECK = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
    os.path.dirname(os.path.abspath(__file__)), '..',
    'July_MT_Command_Centre_REWORKED.pptx')

p = Presentation(DECK)
issues = []

for i, s in enumerate(p.slides, 1):
    shapes = []
    for sh in s.shapes:
        if sh.left is None or sh.top is None:
            continue
        x, y = sh.left / EMU, sh.top / EMU
        w, h = (sh.width or 0) / EMU, (sh.height or 0) / EMU
        txt = sh.text_frame.text.strip()[:44] if sh.has_text_frame else ''
        shapes.append((x, y, w, h, txt, sh))

    for x, y, w, h, txt, sh in shapes:
        b, r = y + h, x + w
        if b > PH + 0.01:
            issues.append(f'S{i}: past page bottom ({b:.2f}in) — {txt!r}')
        if r > PW + 0.01:
            issues.append(f'S{i}: past page right ({r:.2f}in) — {txt!r}')
        if x < -0.01:
            issues.append(f'S{i}: past page left ({x:.2f}in) — {txt!r}')
        # anything that is not the footer rail itself must clear it
        is_footer = y >= FOOT_Y - 0.02
        if not is_footer and b > FOOT_Y + 0.02:
            issues.append(f'S{i}: collides with footer rail (bottom {b:.2f}in) — {txt!r}')
        if not is_footer and h > 0 and y < SRC_Y < b:
            issues.append(f'S{i}: crosses the source line — {txt!r}')

    # card containment: rounded rects wider than 1in are cards; their text children
    # should not extend past the bottom edge
    cards = [t for t in shapes if t[2] > 1.0 and t[3] > 0.8 and not t[4]]
    for cx, cy, cw, ch, _, _ in cards:
        for x, y, w, h, txt, sh in shapes:
            if not txt or h == 0:
                continue
            inside_x = x >= cx - 0.02 and x + w <= cx + cw + 0.06
            starts_in = cy < y < cy + ch
            if inside_x and starts_in and y + h > cy + ch + 0.06:
                issues.append(
                    f'S{i}: text spills below its card '
                    f'(card ends {cy+ch:.2f}, text ends {y+h:.2f}) — {txt!r}')

if issues:
    seen, out = set(), []
    for it in issues:
        if it not in seen:
            seen.add(it); out.append(it)
    print(f'{len(out)} geometry issue(s):')
    for it in out:
        print('  -', it)
    sys.exit(1)
print(f'Geometry OK — {len(p.slides)} slides, no overflow or collisions.')
