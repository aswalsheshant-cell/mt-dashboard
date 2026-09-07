#!/usr/bin/env python3
"""Layout kit for the portrait MT leadership decks (7.5in x 13.33in).

Houses the shared design system so individual deck builders only carry data and
copy. Palette and geometry follow the existing MT decks; the visual grammar
(one takeaway headline, one anchor visual, a short so-what block) follows the
reference chart pack.

Validate any deck built with this using scripts/check_deck_geometry.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------
# palette — taken from the existing MT / Honasa decks
INK    = RGBColor(0x1F, 0x29, 0x33)   # headline near-black
DEEP   = RGBColor(0x18, 0x3B, 0x39)   # deep teal, section blocks
TEAL   = RGBColor(0x2D, 0x9B, 0x7F)   # primary accent
TEAL_D = RGBColor(0x11, 0x6F, 0x68)   # darker teal
GREEN  = RGBColor(0x1E, 0x8E, 0x3E)   # growth
RED    = RGBColor(0xC0, 0x39, 0x2B)   # decline / alert
AMBER  = RGBColor(0xC7, 0x7D, 0x17)   # caution
BLUE   = RGBColor(0x2B, 0x6C, 0xB0)   # secondary series
MUTED  = RGBColor(0x6B, 0x76, 0x82)   # subhead, source
WARM   = RGBColor(0xF4, 0xF1, 0xEC)   # page ground
RULE   = RGBColor(0xE4, 0xE0, 0xD8)   # hairline
MINT   = RGBColor(0xDF, 0xF2, 0xED)   # soft panel
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY_L = RGBColor(0xED, 0xEF, 0xF1)

FONT = "Calibri"

# geometry (inches) — mirrors check_deck_geometry.py expectations
PW, PH = 7.5, 13.333
ML, MR = 0.42, 0.42
CW = PW - ML - MR                      # 6.66in content width
FOOT_Y, SRC_Y = 12.44, 13.00
BODY_BOTTOM = 12.35


_FONT_DIR = "/usr/share/fonts/truetype/liberation"
_measure_cache = {}


def line_count(txt, size_pt, width_in, bold=False):
    """How many lines this text wraps to at `size_pt` in a `width_in` box.

    Measured with Liberation Sans, which is a little wider than Calibri, so the
    count errs on the safe side and headlines never clip. Falls back to a
    character estimate if the font is unavailable.
    """
    key = (txt, size_pt, round(width_in, 3), bold)
    if key in _measure_cache:
        return _measure_cache[key]
    n = None
    try:
        from PIL import ImageFont, ImageDraw, Image
        path = f"{_FONT_DIR}/LiberationSans-{'Bold' if bold else 'Regular'}.ttf"
        f = ImageFont.truetype(path, max(6, int(round(size_pt * 96 / 72.0))))
        d = ImageDraw.Draw(Image.new("RGB", (8, 8)))
        limit = width_in * 96
        n = 0
        for hard in txt.split("\n"):
            line = ""
            n += 1
            for w in hard.split(" "):
                trial = (line + " " + w).strip()
                if d.textlength(trial, font=f) <= limit or not line:
                    line = trial
                else:
                    n += 1
                    line = w
    except (ImportError, OSError, ValueError):
        n = None
    if n is None:
        per_line = max(8, int(width_in * 72 / (size_pt * 0.50)))
        n = max(1, -(-len(txt) // per_line)) + txt.count("\n")
    _measure_cache[key] = n
    return n


def text_height(txt, size_pt, width_in, bold=False, line_spacing=1.0, pad=0.06):
    """Box height that fits `txt` without clipping."""
    return line_count(txt, size_pt, width_in, bold) * \
        size_pt * 1.20 * line_spacing / 72.0 + pad


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(PW)
    prs.slide_height = Inches(PH)
    return prs


def _blank(prs, ground=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if ground is not None:
        set_ground(s, ground)
    return s


def set_ground(slide, rgb):
    """Paint the real slide background, not a full-bleed shape — a shape would
    read as content to check_deck_geometry.py and mask every overlap test."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def fits(y, where=""):
    """Guard: body content must clear the footer rail."""
    if y > BODY_BOTTOM + 0.005:
        raise AssertionError(
            f"body overruns the footer rail: ends {y:.2f}in "
            f"(limit {BODY_BOTTOM:.2f}in) {where}")
    return y


# --------------------------------------------------------------------------
# primitives

def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE,
         adj=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if adj is not None:
        try:
            sh.adjustments[0] = adj
        except (IndexError, AttributeError):
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def text(slide, x, y, w, h, runs, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None, wrap=True):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, o in runs:
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = FONT
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.color.rgb = o.get("color", color)
        sp = o.get("spacing", spacing)
        if sp:
            # character spacing is not exposed by python-pptx; set it on the XML
            r.font._rPr.set("spc", str(int(sp * 100)))
    return box


def para(box, runs, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_before=0, line_spacing=None):
    """Append a paragraph to an existing textbox created by text()."""
    p = box.text_frame.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, o in runs:
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = FONT
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.color.rgb = o.get("color", color)
    return p


def hairline(slide, y, x=ML, w=CW, color=RULE, weight=1.1):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Emu(int(weight * 12700)))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    ln.text_frame.text = ""
    return ln


# --------------------------------------------------------------------------
# page furniture

def page(prs, eyebrow, headline, subhead=None, page_no=None, total=None,
         source=None, ground=WARM, head_size=23):
    """Standard content page. Returns (slide, y) where y is the body top."""
    s = _blank(prs, ground)
    text(s, ML, 0.40, CW, 0.22, eyebrow.upper(), size=8.5, bold=True,
         color=TEAL, spacing=1.4)

    # headline and subhead boxes are measured, not guessed, so nothing clips
    hl_h = text_height(headline, head_size, CW, bold=True, line_spacing=0.94)
    text(s, ML, 0.68, CW, hl_h, headline, size=head_size, bold=True, color=INK,
         line_spacing=0.94)
    y = 0.68 + hl_h + 0.10

    if subhead:
        sh_h = text_height(subhead, 10.5, CW, line_spacing=1.06)
        text(s, ML, y, CW, sh_h, subhead, size=10.5, color=MUTED, line_spacing=1.06)
        y += sh_h + 0.10

    hairline(s, y)
    y += 0.20

    # footer rail — offsets must match RAIL_YS in check_deck_geometry.py
    hairline(s, FOOT_Y, color=RULE)
    if page_no is not None:
        text(s, ML, FOOT_Y + 0.09, CW, 0.22,
             f"{page_no:02d}" + (f" / {total:02d}" if total else ""),
             size=8.5, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    if source:
        text(s, ML, SRC_Y, CW, 0.24, source, size=7.5, color=MUTED)
    return s, y


def divider(prs, kicker, title, blurb=None, page_no=None, total=None):
    s = _blank(prs, DEEP)
    text(s, ML, 4.86, CW, 0.26, kicker.upper(), size=9.5, bold=True,
         color=TEAL, spacing=1.6)
    rect(s, ML, 5.30, 1.05, 0.055, fill=TEAL)
    text(s, ML, 5.58, CW, 1.70, title, size=30, bold=True,
         color=WHITE, line_spacing=0.96)
    if blurb:
        text(s, ML, 7.42, CW - 0.40, 1.00, blurb, size=11.5,
             color=RGBColor(0xC8, 0xDC, 0xD7), line_spacing=1.24)
    if page_no is not None:
        text(s, ML, FOOT_Y + 0.09, CW, 0.22,
             f"{page_no:02d}" + (f" / {total:02d}" if total else ""),
             size=8.5, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    return s


# --------------------------------------------------------------------------
# components

def kpi_row(slide, y, cards, h=1.06, gap=0.11):
    """cards: list of (label, value, delta or None, tone) — tone in +/-/n."""
    n = len(cards)
    w = (CW - gap * (n - 1)) / n
    for i, (label, value, delta, tone) in enumerate(cards):
        x = ML + i * (w + gap)
        rect(slide, x, y, w, h, fill=WHITE, line=RULE, lw=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.07)
        # accent sits inside the rounded corners, not across them
        rect(slide, x + 0.11, y, w - 0.22, 0.05, fill=TEAL)
        text(slide, x + 0.10, y + 0.16, w - 0.20, 0.20, label.upper(), size=7.5,
             bold=True, color=MUTED, spacing=0.8)
        vsize = 17 if len(value) <= 8 else (14.5 if len(value) <= 11 else 12.5)
        text(slide, x + 0.10, y + 0.38, w - 0.20, 0.36, value, size=vsize,
             bold=True, color=INK)
        if delta:
            col = GREEN if tone == "+" else (RED if tone == "-" else MUTED)
            text(slide, x + 0.10, y + 0.78, w - 0.20, 0.20, delta, size=8.5,
                 bold=True, color=col)
    return y + h


def bars(slide, y, rows, h_row=0.40, gap=0.09, label_w=1.42, val_w=0.92,
         color=TEAL, hi_color=None, hi=(), note=None, max_val=None):
    """Horizontal bars, reference-pack style: label left, bar, value at the end.

    rows: list of (label, value, display) — value drives bar length.
    """
    vals = [abs(v) for _, v, _ in rows] or [1]
    mx = max_val or max(vals) or 1
    track = CW - label_w - val_w - 0.16
    for i, (label, val, disp) in enumerate(rows):
        ry = y + i * (h_row + gap)
        is_hi = label in hi
        col = (hi_color or TEAL_D) if is_hi else color
        text(slide, ML, ry + 0.05, label_w - 0.10, h_row - 0.06, label,
             size=9.5, bold=is_hi, color=INK if is_hi else MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
        bw = max(0.05, track * (abs(val) / mx))
        rect(slide, ML + label_w, ry + 0.055, bw, h_row - 0.11, fill=col)
        text(slide, ML + label_w + bw + 0.10, ry + 0.05, val_w, h_row - 0.06,
             disp, size=9.5, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    y2 = y + len(rows) * (h_row + gap) - gap
    if note:
        y2 += 0.12
        text(slide, ML, y2, CW, 0.24, note, size=8, color=MUTED)
        y2 += 0.24
    return y2


def bridge(slide, y, start, steps, end, h=2.85, start_label="", end_label="",
           unit="₹ L", floor=None):
    """Waterfall: start bar, +/- floating steps, end bar.

    `floor` sets the axis base. Small deltas on a large base are invisible from
    a zero baseline, so pass a floor just under the smallest bar — the note
    under the chart states it, so the truncation is never silent.
    """
    n = len(steps) + 2
    gap = 0.10
    bw = (CW - gap * (n - 1)) / n
    # running tops
    running = start
    levels = [(0, start, start_label or "Start", INK)]
    for lbl, delta in steps:
        levels.append((running, running + delta, lbl, GREEN if delta >= 0 else RED))
        running += delta
    levels.append((0, end, end_label or "End", TEAL_D))

    lo = floor if floor is not None else min(
        [min(a, b) for a, b, _, _ in levels] + [0])
    hi = max([max(a, b) for a, b, _, _ in levels])
    span = (hi - lo) or 1
    plot_h = h - 0.86          # room for labels above and below
    base_y = y + 0.42 + plot_h

    def ypos(v):
        return base_y - (v - lo) / span * plot_h

    for i, (a, b, lbl, col) in enumerate(levels):
        x = ML + i * (bw + gap)
        # clamp into the visible band: with a floored axis the start/end bars
        # run from zero and would otherwise extend far below the plot
        cl_hi = min(hi, max(lo, max(a, b)))
        cl_lo = min(hi, max(lo, min(a, b)))
        top, bot = ypos(cl_hi), ypos(cl_lo)
        bh = max(0.05, bot - top)
        rect(slide, x, top, bw, bh, fill=col)
        delta = b - a
        if i == 0 or i == len(levels) - 1:
            disp = f"{b:,.0f}"
        elif abs(delta) < 10:
            disp = f"{delta:+,.2f}"
        else:
            disp = f"{delta:+,.0f}"
        text(slide, x - 0.04, top - 0.24, bw + 0.08, 0.22, disp, size=8.5,
             bold=True, color=col, align=PP_ALIGN.CENTER)
        text(slide, x - 0.04, base_y + 0.08, bw + 0.08, 0.34, lbl, size=7.5,
             color=MUTED, align=PP_ALIGN.CENTER)
    hairline(slide, base_y, color=MUTED, weight=0.75)
    note = unit
    if floor is not None:
        note += f"  ·  axis starts at {floor:,.0f}, not zero"
    text(slide, ML, y + 0.02, CW, 0.20, note, size=7.5, color=MUTED)
    return y + h


def data_table(slide, y, header, rows, widths, hi_rows=(), row_h=0.30,
               head_h=0.30, size=8.5, aligns=None):
    """Compact table. widths are fractions of CW."""
    from pptx.util import Emu as _E
    n_rows = len(rows) + 1
    tbl_h = head_h + len(rows) * row_h
    gf = slide.shapes.add_table(n_rows, len(header), Inches(ML), Inches(y),
                                Inches(CW), Inches(tbl_h))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    for i, frac in enumerate(widths):
        tbl.columns[i].width = Emu(int(Inches(CW) * frac))
    tbl.rows[0].height = Inches(head_h)
    for i in range(len(rows)):
        tbl.rows[i + 1].height = Inches(row_h)
    aligns = aligns or ([PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (len(header) - 1))

    def _cell(c, val, bold, color, fill, align):
        c.fill.solid()
        c.fill.fore_color.rgb = fill
        c.margin_left = c.margin_right = Inches(0.06)
        c.margin_top = c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(val)
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    for j, hcell in enumerate(header):
        _cell(tbl.cell(0, j), hcell, True, WHITE, DEEP, aligns[j])
    for i, row in enumerate(rows):
        hi = row[0] in hi_rows
        fill = MINT if hi else (WHITE if i % 2 == 0 else GREY_L)
        for j, v in enumerate(row):
            col = INK
            sv = str(v)
            if sv.startswith("▲") or sv.startswith("+"):
                col = GREEN
            elif sv.startswith("▼") or sv.startswith("−") or sv.startswith("-"):
                col = RED
            _cell(tbl.cell(i + 1, j), v, hi or j == 0, col, fill, aligns[j])
    return y + tbl_h


def insight(slide, y, title, items, h_item=0.50, accent=TEAL):
    """Titled block of short 'label — sentence' points."""
    text(slide, ML, y, CW, 0.22, title.upper(), size=8.5, bold=True,
         color=accent, spacing=1.2)
    y += 0.30
    for label, body in items:
        rect(slide, ML, y + 0.03, 0.035, h_item - 0.10, fill=accent)
        box = text(slide, ML + 0.16, y, CW - 0.16, h_item,
                   [(label, {"bold": True, "color": INK, "size": 9.5}),
                    ("  —  ", {"color": RULE, "size": 9.5}),
                    (body, {"color": MUTED, "size": 9.5})],
                   line_spacing=1.18)
        box.text_frame.word_wrap = True
        y += h_item
    return y


def sowhat(slide, y, items, h=None, accent=DEEP):
    """The closing 'so what' strip — verb-led calls."""
    n = len(items)
    gap = 0.10
    w = (CW - gap * (n - 1)) / n
    h = h or 1.16
    for i, (verb, body) in enumerate(items):
        x = ML + i * (w + gap)
        rect(slide, x, y, w, h, fill=MINT, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.06)
        text(slide, x + 0.12, y + 0.12, w - 0.24, 0.20, verb.upper(), size=8,
             bold=True, color=TEAL_D, spacing=1.0)
        b = text(slide, x + 0.12, y + 0.36, w - 0.24, h - 0.46, body, size=9,
                 color=INK, line_spacing=1.16)
        b.text_frame.word_wrap = True
    return y + h


def legend(slide, y, entries, size=8.5):
    """entries: list of (label, color)."""
    x = ML
    for label, col in entries:
        rect(slide, x, y + 0.045, 0.13, 0.13, fill=col)
        w = 0.10 + len(label) * 0.062
        text(slide, x + 0.19, y, w, 0.22, label, size=size, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
        x += 0.19 + w + 0.16
    return y + 0.24


# --------------------------------------------------------------------------
# charts

def nice_ceiling(v):
    """Round an axis maximum up to a readable 1/2/2.5/5 x 10^n value."""
    import math
    if v <= 0:
        return 1
    mag = 10 ** math.floor(math.log10(v))
    for step in (1, 1.25, 1.5, 2, 2.5, 3, 4, 5, 6, 7.5, 8, 10):
        if v <= step * mag:
            return step * mag
    return 10 * mag


def trend(slide, y, x_labels, series, h=3.10, unit="₹ L", y_max=None, y_min=0,
          gridlines=4, label_every=3, markers_every=0, band=None):
    """Multi-series line chart drawn from connectors, so it stays editable.

    series: list of dicts {name, values (None = gap), color, width, dash}
    band:   optional (start_idx, end_idx, label) shaded region
    """
    from pptx.enum.shapes import MSO_CONNECTOR

    plot_l = ML + 0.56
    plot_w = CW - 0.56
    plot_t = y + 0.30
    plot_h = h - 0.94
    plot_b = plot_t + plot_h

    flat = [v for s in series for v in s["values"] if v is not None]
    hi = y_max or nice_ceiling(max(flat) if flat else 1)
    lo = y_min if y_min is not None else (min(flat) if flat else 0)
    span = (hi - lo) or 1

    def px(i):
        n = max(1, len(x_labels) - 1)
        return plot_l + plot_w * i / n

    def py(v):
        return plot_b - (v - lo) / span * plot_h

    if band:
        b0, b1, blabel = band
        bx0, bx1 = px(b0), px(b1)
        rect(slide, bx0, plot_t, max(0.04, bx1 - bx0), plot_h, fill=GREY_L)
        if blabel:
            text(slide, bx0, plot_t - 0.20, bx1 - bx0 + 0.9, 0.18, blabel,
                 size=7, bold=True, color=MUTED)

    # gridlines + y labels
    for g in range(gridlines + 1):
        v = lo + span * g / gridlines
        gy = py(v)
        hairline(slide, gy, x=plot_l, w=plot_w,
                 color=RULE if g else MUTED, weight=0.75)
        text(slide, ML - 0.02, gy - 0.10, 0.52, 0.20, f"{v:,.0f}", size=7,
             color=MUTED, align=PP_ALIGN.RIGHT)

    # x labels
    for i, lab in enumerate(x_labels):
        if i % label_every and i != len(x_labels) - 1:
            continue
        text(slide, px(i) - 0.30, plot_b + 0.08, 0.60, 0.20, lab, size=6.5,
             color=MUTED, align=PP_ALIGN.CENTER)

    # series
    for s in series:
        col = s.get("color", TEAL)
        wpt = s.get("width", 2.0)
        vals = s["values"]
        for i in range(len(vals) - 1):
            a, b = vals[i], vals[i + 1]
            if a is None or b is None:
                continue
            cn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(px(i)), Inches(py(a)),
                Inches(px(i + 1)), Inches(py(b)))
            cn.line.color.rgb = col
            cn.line.width = Pt(wpt)
            if s.get("dash"):
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                d = s["dash"]
                cn.line.dash_style = (getattr(MSO_LINE_DASH_STYLE, d)
                                      if isinstance(d, str) else d)
        step = markers_every or 0
        real = [i for i, v in enumerate(vals) if v is not None]

        def _dot(i, size=0.085):
            rect(slide, px(i) - size / 2, py(vals[i]) - size / 2, size, size,
                 fill=col, shape=MSO_SHAPE.OVAL)

        if step:
            for i in real:
                if i % step == 0:
                    _dot(i)
        # a point whose neighbours are both gaps draws no line segment, so it
        # would vanish entirely — mark it, and always mark the series end.
        for i in real:
            prev_gap = i == 0 or vals[i - 1] is None
            next_gap = i == len(vals) - 1 or vals[i + 1] is None
            if prev_gap and next_gap:
                _dot(i, 0.10)
        if real:
            _dot(real[-1], 0.10)
    # end-point labels go on last so no later line is drawn over them
    for s in series:
        vals = s["values"]
        seen = [i for i, v in enumerate(vals) if v is not None]
        if seen and s.get("end_label"):
            i = seen[-1]
            dy = 0.10 if s.get("end_label_below") else -0.25
            text(slide, px(i) - 0.90, py(vals[i]) + dy, 0.88, 0.20,
                 s["end_label"], size=8, bold=True, color=s.get("color", TEAL),
                 align=PP_ALIGN.RIGHT)

    text(slide, ML, y, 1.6, 0.20, unit, size=7.5, color=MUTED)
    ly = legend(slide, plot_b + 0.32, [(s["name"], s.get("color", TEAL))
                                       for s in series])
    return max(ly, y + h)


def donut(slide, x, y, w, h, pairs, colors, centre_value=None, centre_label=None,
          hole=62):
    """Native doughnut chart (stays editable in PowerPoint) + centre callout."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    cd = CategoryChartData()
    cd.categories = [p[0] for p in pairs]
    cd.add_series("share", [p[1] for p in pairs])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(8.5)
    ch.legend.font.name = FONT
    ch.legend.font.color.rgb = INK
    plot = ch.plots[0]
    try:
        plot.doughnut_hole_size = hole
    except (AttributeError, ValueError):
        pass
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1.25)
    if centre_value:
        # the plot area sits in the left ~62% of the frame once the legend is out
        cx = x + w * 0.31
        text(slide, cx - 0.75, y + h / 2 - 0.30, 1.50, 0.34, centre_value,
             size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
        if centre_label:
            # keep the caption inside the ring's hole, not over the ring itself
            text(slide, cx - 0.62, y + h / 2 + 0.07, 1.24, 0.20, centre_label,
                 size=7, color=MUTED, align=PP_ALIGN.CENTER)
    return y + h


def col_compare(slide, y, groups, series_names, colors, h=2.60, unit="₹ L",
                value_fmt="{:,.0f}"):
    """Grouped vertical columns — for a small number of side-by-side compares.

    groups: list of (label, [v1, v2, ...])
    """
    plot_t = y + 0.34
    plot_h = h - 1.00
    plot_b = plot_t + plot_h
    flat = [v for _, vs in groups for v in vs if v is not None]
    hi = max(flat) if flat else 1
    n_g = len(groups)
    n_s = len(series_names)
    slot = CW / n_g
    barw = min(0.58, (slot - 0.24) / n_s)
    for gi, (label, vals) in enumerate(groups):
        gx = ML + gi * slot + (slot - barw * n_s) / 2
        for si, v in enumerate(vals):
            if v is None:
                continue
            bh = max(0.04, plot_h * v / hi)
            bx = gx + si * barw
            rect(slide, bx, plot_b - bh, barw - 0.045, bh,
                 fill=colors[si % len(colors)])
            text(slide, bx - 0.10, plot_b - bh - 0.23, barw + 0.16, 0.20,
                 value_fmt.format(v), size=7.5, bold=True,
                 color=colors[si % len(colors)], align=PP_ALIGN.CENTER)
        text(slide, ML + gi * slot, plot_b + 0.09, slot, 0.34, label, size=8,
             color=MUTED, align=PP_ALIGN.CENTER)
    hairline(slide, plot_b, color=MUTED, weight=0.75)
    text(slide, ML, y, 1.6, 0.20, unit, size=7.5, color=MUTED)
    ly = legend(slide, plot_b + 0.44, list(zip(series_names, colors)))
    return max(ly, y + h)
