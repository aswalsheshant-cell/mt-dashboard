"""Build V13 from V12 — apply verified Jul'26 national data corrections + normAutofit fixes."""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
import copy, re

SRC = "/root/.claude/uploads/55682d72-6b69-5e1a-b307-5d5ff59817ad/3f5aae9c-MT_Jul26_Honasa_Finalv12.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv13.pptx"

prs = Presentation(SRC)

def get_text(shape):
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def patch_text(shape, old_fragment, new_fragment, label=""):
    """Replace old_fragment with new_fragment anywhere in shape's runs.
    Handles multi-run text by concatenating runs first, then patching.
    Falls back to in-place run substitution for simple cases.
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    full_text = tf.text

    if old_fragment not in full_text:
        return False

    # Simple case: try to replace within individual runs first
    replaced = False
    for para in tf.paragraphs:
        for run in para.runs:
            if old_fragment in run.text:
                run.text = run.text.replace(old_fragment, new_fragment)
                replaced = True

    if replaced:
        if label:
            print(f"    [{label}] '{old_fragment}' → '{new_fragment}'")
        return True

    # Multi-run case: rebuild the first paragraph's runs as a single run
    for para in tf.paragraphs:
        para_text = "".join(r.text for r in para.runs)
        if old_fragment in para_text:
            new_para_text = para_text.replace(old_fragment, new_fragment)
            if para.runs:
                para.runs[0].text = new_para_text
                for r in para.runs[1:]:
                    r.text = ""
            replaced = True
            if label:
                print(f"    [{label}] (multi-run) '{old_fragment}' → '{new_fragment}'")
            break

    return replaced


def apply_normautofit(shape):
    """Enable auto-shrink on overflow text boxes."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


# ─── Slide 1 (index 0): National KPIs ─────────────────────────────────────────
print("=== Slide 1 — National KPI corrections ===")
s1 = prs.slides[0]

fixes_s1 = [
    # (shape_idx, old_fragment, new_fragment, label)
    (5,  '73.4% flow conversion',              '76.8% flow conversion',              'Conv headline'),
    (6,  '₹49.2 Cr primary',                  '₹47.02 Cr primary',                  'Primary'),
    (6,  '₹13.11 Cr gap',                      '₹10.92 Cr gap',                      'Gap in headline'),
    (13, '73.4% conversion | ₹13.11 Cr gap',  '76.8% conversion | ₹10.92 Cr gap',  'Sales flow KPI'),
    (53, '₹24.49 Cr',                          '₹24.55 Cr',                          'ME offtake'),
    (53, '73.4% conversion',                   '76.8% conversion',                   'ME conv label'),
    (55, '₹11.03 Cr',                          '₹11.08 Cr',                          'TDC offtake'),
]
for idx, old, new, label in fixes_s1:
    shape = s1.shapes[idx]
    ok = patch_text(shape, old, new, f"S1:{idx} {label}")
    if not ok:
        print(f"  WARNING S1 shape {idx}: '{old}' not found in '{get_text(shape)[:80]}'")

# ─── Slide 13 (index 12): eB2B — update MT avg reference ──────────────────────
print("\n=== Slide 13 — eB2B MT avg baseline ===")
s13 = prs.slides[12]
# 99.4 - 73.4 = 26.0 pp  →  99.4 - 76.8 = 22.6 pp
patch_text(s13.shapes[18], '+26.0 pp vs MT avg (73.4%)', '+22.6 pp vs MT avg (76.8%)', 'S13:18 eB2B pp')

# ─── Slide 14 (index 13): NPI slide ───────────────────────────────────────────
print("\n=== Slide 14 — NPI corrections ===")
s14 = prs.slides[13]

# Shape 1: title/subtitle
patch_text(s14.shapes[1], '₹2.82 Cr', '₹2.68 Cr', 'S14:1 title')
# Shape 7: KPI card value
patch_text(s14.shapes[7], '₹2.82 Cr', '₹2.68 Cr', 'S14:7 KPI card')
# Shape 12: NPI% of MT
patch_text(s14.shapes[12], '7.82%', '7.4%', 'S14:12 NPI%')
# Shape 47: evidence text (long — Reliance/DMart NPI sub-values recalculated)
# Reliance NPI: 46.3% × 2.68 = ₹1.24 Cr; DMart NPI: 34.6% × 2.68 = ₹0.93 Cr
patch_text(s14.shapes[47], '₹2.82 Cr NPI',  '₹2.68 Cr NPI',  'S14:47 evidence')
patch_text(s14.shapes[47], '(₹1.30 Cr)',    '(₹1.24 Cr)',    'S14:47 Reliance NPI')
patch_text(s14.shapes[47], '(₹0.98 Cr)',    '(₹0.93 Cr)',    'S14:47 DMart NPI')

# ─── Slide 15 (index 14): Chain deep-dive — update evidence text ────────────────
print("\n=== Slide 15 — Chain deep-dive ===")
s15 = prs.slides[14]
# Shape 50 long evidence text (already correct chain gap figures; only national gap ref changes)
patch_text(s15.shapes[50], '₹13.11',  '₹10.92',  'S15:50 gap ref')

# ─── Slide 16 (index 15): Brand breakdown ─────────────────────────────────────
print("\n=== Slide 16 — Brand breakdown ===")
s16 = prs.slides[15]
# Shape 13
patch_text(s16.shapes[13], '67.8% of offtake and converts at 73.4%', '68.0% of offtake and converts at 73.4%', 'S16:13 ME%')
# Shape 15
patch_text(s16.shapes[15], '₹11.03 Cr', '₹11.08 Cr', 'S16:15 TDC offtake')
# Shape 24
patch_text(s16.shapes[24], '₹24.49 Cr at 73.4% conv', '₹24.55 Cr at 73.4% conv', 'S16:24 ME detail')
patch_text(s16.shapes[24], '₹11.03 Cr at 72.6% conv', '₹11.08 Cr at 72.6% conv', 'S16:24 TDC detail')

# ─── Slide 17 (index 16): Trend/history ───────────────────────────────────────
print("\n=== Slide 17 — Trend history ===")
s17 = prs.slides[16]
patch_text(s17.shapes[62], '₹13.11 Cr gap', '₹10.92 Cr gap', 'S17:62 trend text')

# ─── Slide 18 (index 17): Audit card ──────────────────────────────────────────
print("\n=== Slide 18 — Audit card ===")
s18 = prs.slides[17]
patch_text(s18.shapes[10], '₹49.21 Cr', '₹47.02 Cr', 'S18:10 PRIMARY total')
patch_text(s18.shapes[30], '73.4%',     '76.8%',      'S18:30 FLOW CONVERSION')

# ─── normAutofit across all slides for overflow text ──────────────────────────
print("\n=== Applying normAutofit to text-heavy shapes ===")
overflow_fixed = 0
for si, slide in enumerate(prs.slides):
    for shi, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        try:
            total_chars = len(tf.text)
            box_h_pt = shape.height / 12700
            box_w_pt = shape.width / 12700
            if box_h_pt > 0 and total_chars > 0:
                lines_available = max(1, box_h_pt / 12)
                chars_per_line = max(40, box_w_pt / 6)
                lines_needed = total_chars / chars_per_line
                if lines_needed > lines_available * 1.3:
                    apply_normautofit(shape)
                    overflow_fixed += 1
        except Exception:
            pass

print(f"  normAutofit applied to {overflow_fixed} shape(s)")

# ─── Save ──────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved → {DST}")

# ─── Validate ─────────────────────────────────────────────────────────────────
def validate(path):
    p = Presentation(path)
    print(f"\n=== VALIDATE {path} ===")
    print(f"  Slides: {len(p.slides)}")
    checks = [
        (0,  5,  '76.8%'),
        (0,  6,  '₹47.02 Cr'),
        (0,  6,  '₹10.92 Cr'),
        (0,  13, '76.8% conversion'),
        (0,  53, '₹24.55 Cr'),
        (0,  55, '₹11.08 Cr'),
        (12, 18, '76.8%'),
        (13, 7,  '₹2.68 Cr'),
        (13, 12, '7.4%'),
        (15, 13, '68.0%'),
        (15, 15, '₹11.08 Cr'),
        (17, 10, '₹47.02 Cr'),
        (17, 30, '76.8%'),
    ]
    all_ok = True
    for si, shi, expected in checks:
        shape = p.slides[si].shapes[shi]
        txt = get_text(shape)
        ok = expected in txt
        status = "OK" if ok else "FAIL"
        if not ok:
            all_ok = False
        print(f"  S{si+1}:{shi} [{status}] expected '{expected}' → got '{txt[:80]}'")
    return all_ok

ok = validate(DST)
print(f"\n{'ALL CHECKS PASSED' if ok else 'SOME CHECKS FAILED — see above'}")
