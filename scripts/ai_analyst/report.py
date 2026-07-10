"""
Phase 5 — offline report export.

Assembles analysis findings (free text, KPI tiles, tables, EDA profiles, query
results, document summaries) into a Report that renders to several formats:

  * Markdown / HTML / CSV — pure standard library, fully offline and tested here.
    The HTML is self-contained and styled to match the dashboard palette, so it
    opens by double-click and prints cleanly to PDF from any browser.
  * XLSX  — via openpyxl (pluggable; clear error if the package is absent).
  * PPTX  — via python-pptx (already used by this repo's slide builders), so
    report decks reuse the same dependency the PowerBI/ scripts rely on.

Nothing here touches the network; a report is built from data already in hand.
"""

from __future__ import annotations

import csv
import html
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Sequence, Tuple

from ai_analyst.profiler import TableProfile
from ai_analyst.nl2sql import QueryResult


class ReportDependencyError(Exception):
    """A backend library is required for this export format but not installed."""


def _num(x: Optional[float]) -> str:
    if x is None:
        return "-"
    if abs(x - round(x)) < 1e-9:
        return str(int(round(x)))
    return f"{x:.2f}"


def _growth_value(v) -> Optional[float]:
    """Parse a possibly-formatted growth cell (e.g. '+12.4%') to a float."""
    if v is None:
        return None
    try:
        return float(str(v).replace("%", "").replace("+", "").replace(",", "").strip())
    except (TypeError, ValueError):
        return None


def _growth_color(v) -> Optional[str]:
    n = _growth_value(v)
    if n is None or n == 0:
        return None
    return "#1E8E3E" if n > 0 else "#C0392B"  # dashboard green / red


def _growth_marker(v) -> str:
    n = _growth_value(v)
    if n is None or n == 0:
        return ""
    return "▲ " if n > 0 else "▼ "


# dashboard palette (kept in sync with dashboard/index.html :root)
_TEAL = "#2D9B7F"
_TEAL_D = "#1f7a63"
_BG = "#F4F1EC"
_INK = "#1F2933"
_LINE = "#e4e0d8"


@dataclass
class Section:
    kind: str  # 'text' | 'kpis' | 'table'
    title: str = ""
    text: str = ""
    columns: List[str] = field(default_factory=list)
    rows: List[Sequence] = field(default_factory=list)
    kpis: List[Tuple[str, str]] = field(default_factory=list)
    note: str = ""
    growth_col: Optional[int] = None  # index of a column to colour green(+)/red(-)


class Report:
    def __init__(self, title: str, subtitle: str = "", classification: str = ""):
        self.title = title
        self.subtitle = subtitle
        self.classification = classification  # e.g. "Confidential - MT Internal"
        self.sections: List[Section] = []

    # -- builders ----------------------------------------------------------
    def add_text(self, text: str, title: str = "") -> "Report":
        self.sections.append(Section("text", title=title, text=text))
        return self

    def add_kpis(self, kpis: Sequence[Tuple[str, str]], title: str = "Key figures") -> "Report":
        self.sections.append(Section("kpis", title=title, kpis=[(str(k), str(v)) for k, v in kpis]))
        return self

    def add_table(self, columns: Sequence[str], rows: Sequence[Sequence],
                  title: str = "", note: str = "", growth_col: Optional[int] = None) -> "Report":
        self.sections.append(Section("table", title=title, columns=list(columns),
                                     rows=[list(r) for r in rows], note=note, growth_col=growth_col))
        return self

    def add_query_result(self, res: QueryResult, title: str = "") -> "Report":
        t = title or f"Q: {res.question}"
        note = f"SQL ({res.provider}): {res.sql}" if res.sql else (res.error or "")
        if res.ok:
            self.add_table(res.columns, res.rows, title=t, note=note)
        else:
            self.add_text(f"Query failed: {res.error}", title=t)
        return self

    def add_profile(self, profile: TableProfile, title: str = "") -> "Report":
        t = title or f"EDA profile — {profile.table}"
        head = f"{profile.rows_total} rows × {len(profile.columns)} columns"
        if profile.duplicates:
            head += f"; {profile.duplicates} duplicate rows"
        cols = ["column", "kind", "nulls", "distinct", "summary"]
        rows = []
        for c in profile.columns:
            if c.kind == "numeric":
                summary = f"min={_num(c.minimum)} max={_num(c.maximum)} mean={_num(c.mean)}"
            elif c.kind == "categorical":
                summary = ", ".join(f"{v} ({n})" for v, n in c.top_values[:3])
            else:
                summary = "(empty)"
            rows.append([c.name, c.kind, f"{c.nulls} ({c.null_pct:.0f}%)", c.distinct, summary])
        self.add_table(cols, rows, title=t, note=head)
        return self

    # -- renderers ---------------------------------------------------------
    def to_markdown(self) -> str:
        out = []
        if self.classification:
            out.append(f"> **{self.classification}**")
            out.append("")
        out.append(f"# {self.title}")
        if self.subtitle:
            out.append(f"_{self.subtitle}_")
        for s in self.sections:
            out.append("")
            if s.title:
                out.append(f"## {s.title}")
            if s.kind == "text":
                out.append(s.text)
            elif s.kind == "kpis":
                for k, v in s.kpis:
                    out.append(f"- **{k}:** {v}")
            elif s.kind == "table":
                if s.note:
                    out.append(f"_{s.note}_\n")
                out.append("| " + " | ".join(str(c) for c in s.columns) + " |")
                out.append("| " + " | ".join("---" for _ in s.columns) + " |")
                for r in s.rows:
                    cells = []
                    for ci, v in enumerate(r):
                        txt = "" if v is None else str(v)
                        if s.growth_col is not None and ci == s.growth_col:
                            txt = _growth_marker(v) + txt
                        cells.append(txt)
                    out.append("| " + " | ".join(cells) + " |")
        return "\n".join(out) + "\n"

    def to_html(self) -> str:
        e = html.escape
        parts = [
            "<main style=\"font-family:'Segoe UI',system-ui,Arial,sans-serif;"
            f"max-width:1000px;margin:0 auto;padding:24px;color:{_INK};background:{_BG}\">",
            f"<h1 style='color:{_TEAL_D};margin:0 0 4px'>{e(self.title)}</h1>",
        ]
        if self.classification:
            parts.insert(1,
                "<div style='background:#7a1f1f;color:#fff;font-weight:700;letter-spacing:.5px;"
                "text-align:center;padding:6px;border-radius:6px;margin:0 0 14px;"
                f"font-size:12px'>{e(self.classification)}</div>")
        if self.subtitle:
            parts.append(f"<p style='color:#6b7682;margin:0 0 18px'>{e(self.subtitle)}</p>")
        for s in self.sections:
            parts.append("<section style='margin:22px 0'>")
            if s.title:
                parts.append(f"<h2 style='font-size:17px;color:{_INK};"
                             f"border-bottom:2px solid {_TEAL};padding-bottom:4px'>{e(s.title)}</h2>")
            if s.kind == "text":
                parts.append(f"<p style='line-height:1.55'>{e(s.text)}</p>")
            elif s.kind == "kpis":
                parts.append("<div style='display:flex;flex-wrap:wrap;gap:12px'>")
                for k, v in s.kpis:
                    parts.append(
                        "<div style='background:#fff;border:1px solid " + _LINE +
                        ";border-radius:12px;padding:12px 16px;min-width:140px'>"
                        f"<div style='font-size:11px;text-transform:uppercase;color:#6b7682'>{e(k)}</div>"
                        f"<div style='font-size:24px;font-weight:800;color:{_TEAL_D}'>{e(v)}</div></div>")
                parts.append("</div>")
            elif s.kind == "table":
                if s.note:
                    parts.append(f"<p style='font-size:12px;color:#6b7682'>{e(s.note)}</p>")
                parts.append("<table style='width:100%;border-collapse:collapse;font-size:13px'>")
                parts.append("<thead><tr>" + "".join(
                    f"<th style='background:{_TEAL};color:#fff;text-align:left;padding:8px 10px'>{e(str(c))}</th>"
                    for c in s.columns) + "</tr></thead><tbody>")
                for r in s.rows:
                    tds = []
                    for ci, v in enumerate(r):
                        style = f"padding:8px 10px;border-bottom:1px solid {_LINE}"
                        if s.growth_col is not None and ci == s.growth_col:
                            col = _growth_color(v)
                            if col:
                                style += f";color:{col};font-weight:700"
                        tds.append(f"<td style='{style}'>{e('' if v is None else str(v))}</td>")
                    parts.append("<tr>" + "".join(tds) + "</tr>")
                parts.append("</tbody></table>")
            parts.append("</section>")
        parts.append("</main>")
        return "\n".join(parts)

    def first_table(self) -> Optional[Section]:
        return next((s for s in self.sections if s.kind == "table"), None)

    def to_csv(self) -> str:
        s = self.first_table()
        if s is None:
            raise ValueError("No table section to export as CSV.")
        buf = io.StringIO()
        w = csv.writer(buf)
        w.writerow(s.columns)
        for r in s.rows:
            w.writerow(["" if v is None else v for v in r])
        return buf.getvalue()

    # -- save --------------------------------------------------------------
    def save(self, path) -> str:
        p = Path(path)
        ext = p.suffix.lower()
        p.parent.mkdir(parents=True, exist_ok=True)
        if ext in (".md", ".markdown"):
            p.write_text(self.to_markdown(), encoding="utf-8")
        elif ext in (".html", ".htm"):
            p.write_text(self.to_html(), encoding="utf-8")
        elif ext == ".csv":
            p.write_text(self.to_csv(), encoding="utf-8")
        elif ext == ".xlsx":
            self._save_xlsx(p)
        elif ext == ".pptx":
            self._save_pptx(p)
        else:
            raise ValueError(f"Unsupported report format: {ext!r} "
                             "(use .md, .html, .csv, .xlsx, .pptx)")
        return str(p)

    def _save_xlsx(self, p: Path) -> None:
        try:
            import openpyxl  # type: ignore
        except ImportError:
            raise ReportDependencyError("XLSX export needs openpyxl: `pip install openpyxl`.")
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        if self.classification:
            cover = wb.create_sheet("Cover")
            cover["A1"] = self.classification
            cover["A2"] = self.title
            if self.subtitle:
                cover["A3"] = self.subtitle
        i = 0
        for s in self.sections:
            if s.kind != "table":
                continue
            i += 1
            ws = wb.create_sheet((s.title or f"Table {i}")[:31] or f"Table {i}")
            ws.append([str(c) for c in s.columns])
            for r in s.rows:
                ws.append(["" if v is None else v for v in r])
        if i == 0:
            wb.create_sheet("Report").append([self.title])
        wb.save(str(p))

    def _save_pptx(self, p: Path) -> None:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
        except ImportError:
            raise ReportDependencyError("PPTX export needs python-pptx: `pip install python-pptx`.")
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = self.title
        sub = self.subtitle
        if self.classification:
            sub = f"{self.classification}\n{sub}" if sub else self.classification
        if sub:
            title_slide.placeholders[1].text = sub
        for s in self.sections:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = s.title or s.kind
            if s.kind == "table" and s.rows:
                rows, cols = min(len(s.rows) + 1, 12), len(s.columns)
                tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5),
                                             Inches(9), Inches(0.4 * rows)).table
                for ci, c in enumerate(s.columns):
                    tbl.cell(0, ci).text = str(c)
                for ri, r in enumerate(s.rows[:rows - 1], start=1):
                    for ci, v in enumerate(r):
                        tbl.cell(ri, ci).text = "" if v is None else str(v)
            else:
                body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                if body is not None:
                    body.text = s.text or "\n".join(f"{k}: {v}" for k, v in s.kpis)
        prs.save(str(p))
