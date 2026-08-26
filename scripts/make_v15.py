"""Build V15 from V13 — Qty/YoY decomposition with layout-safe text and shape resizing.

Fixes vs V14:
- Shape [8] (primary sublabel 1.25"×0.23"): use compact "+XX% YoY | XXXXKu" text + expand height
- Shape [101]/[100] (evidence 1.56"×0.38"): concise chain YoY summary (no long chain_note)
- Shape [58]/[57] (insight 02): shorter text that fits in 1.67"×0.47"
- normAutofit on every modified shape
- audit_overflows check before save
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu
from lxml import etree

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv13.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv15.pptx"

# ── Zone data — all from correct sources (ShipTo Jul'25 + article Jul'26) ──────
# pri25/pri26 = ₹ Cr primary NSV; pq25/pq26 = thousand units
# off26 = ₹ Cr offtake NSV (authoritative from PPT); oq26 = thousand units
# ins01_idx / ins02_idx = insight boxes (inside 2.17"×0.56" containers)
# ev_idx = evidence box (EVIDENCE column, 1.56"×0.38")
ZONE = {
    'West': {
        'slide_idx': 4,
        'pri25': 5.05, 'pq25': 315.0,
        'pri26': 10.05, 'pq26': 597.8,
        'pri_yoy_pct': 99.0, 'pq_yoy_pct': 89.8,
        'off26': 8.28, 'oq26': 466.6,
        'mix_pct': 22.9,
        'ins01_idx': 55, 'ins02_idx': 58, 'ev_idx': 101,
        # chains: (name, pri25, pri26, yoy%) — from ShipTo Jul'25 file
        'chains': [
            ('DMart',            3.27, 6.12,  87.0),
            ('Reliance',         1.02, 2.26, 121.6),
            ('Apollo',           0.26, 0.62, 138.5),
        ],
        # Concise evidence text — must fit in 1.56" W × 0.38" H (≈50–60 chars)
        'ev_brief': "Prim YoY: DMart +87% | Reliance +122% | Apollo +139% | WF NEW",
        # Action insight for chain labels (already in PPT insight boxes)
        'chain_note_long': (
            "Primary YoY (ShipTo source): DMart ₹3.27→₹6.12Cr (+87.0%); "
            "Reliance ₹1.02→₹2.26Cr (+121.6%); Apollo ₹0.26→₹0.62Cr (+138.5%); "
            "Wellness Forever NEW ₹0.49Cr. "
            "Reliance West conversion 54.5% vs primary +121.6% → stock build watch."
        ),
    },
    'South-1': {
        'slide_idx': 5,
        'pri25': 5.99, 'pq25': 349.0,
        'pri26': 9.80, 'pq26': 581.6,
        'pri_yoy_pct': 63.6, 'pq_yoy_pct': 66.6,
        'off26': 8.19, 'oq26': 435.3,
        'mix_pct': 22.7,
        'ins01_idx': 55, 'ins02_idx': 58, 'ev_idx': 101,
        'chains': [
            ('Apollo',    1.22, 3.51, 187.7),
            ('DMart',     1.60, 3.09,  93.1),
            ('Reliance',  1.14, 2.24,  96.5),
        ],
        'ev_brief': "Prim YoY: Apollo +188% | DMart +93% | Reliance +97%",
        'chain_note_long': (
            "Primary YoY: Apollo ₹1.22→₹3.51Cr (+187.7%); "
            "DMart ₹1.60→₹3.09Cr (+93.1%); Reliance ₹1.14→₹2.24Cr (+96.5%). "
            "Apollo S-1 surge: validate offtake velocity before next loading."
        ),
    },
    'North': {
        'slide_idx': 6,
        'pri25': 5.10, 'pq25': 306.7,
        'pri26': 11.95, 'pq26': 649.1,
        'pri_yoy_pct': 134.3, 'pq_yoy_pct': 111.6,
        'off26': 6.99, 'oq26': 377.6,
        'mix_pct': 19.4,
        'ins01_idx': 55, 'ins02_idx': 58, 'ev_idx': 101,
        'chains': [
            ('Reliance',  1.78, 5.35, 200.6),
            ('DMart',     1.16, 3.24, 179.3),
            ('Apollo',    0.47, 1.17, 148.9),
        ],
        'ev_brief': "Prim YoY: Reliance +201% | DMart +179% | Apollo +149%",
        'chain_note_long': (
            "Primary YoY: Reliance ₹1.78→₹5.35Cr (+200.6%); "
            "DMart ₹1.16→₹3.24Cr (+179.3%); Apollo ₹0.47→₹1.17Cr (+148.9%). "
            "PAISA VASOOL EFFECT: Reliance North +200.6% mirrors LY Jul'25 buildup — "
            "validate Aug'26 offtake by 05-Aug before holding or loading."
        ),
    },
    'South-2': {
        'slide_idx': 7,
        'pri25': 4.07, 'pq25': 223.4,
        'pri26': 6.89, 'pq26': 343.8,
        'pri_yoy_pct': 69.3, 'pq_yoy_pct': 53.9,
        'off26': 4.91, 'oq26': 281.1,
        'mix_pct': 13.6,
        'ins01_idx': 54, 'ins02_idx': 57, 'ev_idx': 100,
        'chains': [
            ('DMart',     1.17, 4.32, 269.2),
            ('Apollo',    1.10, 1.10,   0.0),
            ('Reliance',  0.64, 0.81,  26.6),
        ],
        'ev_brief': "Prim YoY: DMart +269% | Apollo flat | Reliance +27%",
        'chain_note_long': (
            "Primary YoY: DMart ₹1.17→₹4.32Cr (+269.2%); "
            "Apollo ₹1.10→₹1.10Cr (flat); Reliance ₹0.64→₹0.81Cr (+26.6%). "
            "DMart S-2 surge +269%: validate offtake before next loading."
        ),
    },
    'East': {
        'slide_idx': 8,
        'pri25': 3.63, 'pq25': 211.3,
        'pri26': 7.83, 'pq26': 409.4,
        'pri_yoy_pct': 115.7, 'pq_yoy_pct': 93.8,
        'off26': 3.55, 'oq26': 183.9,
        'mix_pct': 9.8,
        'ins01_idx': 55, 'ins02_idx': 58, 'ev_idx': 101,
        'chains': [
            ('Reliance',  1.87, 4.09, 118.7),
            ('Vmart',     0.00, 2.00,  None),
            ('Apollo',    0.40, 0.66,  65.0),
        ],
        'ev_brief': "Prim YoY: Reliance +119% | Vmart NEW | Apollo +65%",
        'chain_note_long': (
            "Primary YoY: Reliance ₹1.87→₹4.09Cr (+118.7%); "
            "Vmart NEW ₹2.00Cr; Apollo ₹0.40→₹0.66Cr (+65.0%). "
            "PAISA VASOOL: East 45.3% conv → over-stock signal; check Aug'26 by 05-Aug."
        ),
    },
    'Central': {
        'slide_idx': 9,
        'pri25': 0.89, 'pq25': 55.4,
        'pri26': 2.69, 'pq26': 160.9,
        'pri_yoy_pct': 202.2, 'pq_yoy_pct': 190.4,
        'off26': 2.12, 'oq26': 160.6,
        'mix_pct': 5.9,
        'ins01_idx': 54, 'ins02_idx': 57, 'ev_idx': 100,
        'chains': [
            ('DMart',     0.55, 1.48, 169.1),
            ('Reliance',  0.23, 0.91, 295.7),
            ('Apollo',    0.00, 0.14,  None),
        ],
        'ev_brief': "Prim YoY: DMart +169% | Reliance +296% | Apollo NEW",
        'chain_note_long': (
            "Primary YoY: DMart ₹0.55→₹1.48Cr (+169.1%); "
            "Reliance ₹0.23→₹0.91Cr (+295.7%); Apollo NEW ₹0.14Cr. "
            "Central NEW zone (+202.2% base effect); protect 95.3% DMart conversion."
        ),
    },
}

NAT_PRI25_NSV = 24.73
NAT_PRI25_QTY = 1460.8
NAT_PRI26_NSV = 49.21
NAT_PRI26_QTY = 2742.6
NAT_OFF26_NSV = 36.10
NAT_OFF26_QTY = 2114.2


# ── Helpers ────────────────────────────────────────────────────────────────────
def get_text(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def patch(shape, old_frag: str, new_frag: str, label="") -> bool:
    if not shape.has_text_frame:
        return False
    full = shape.text_frame.text
    if old_frag not in full:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_frag in run.text:
                run.text = run.text.replace(old_frag, new_frag)
                if label:
                    print(f"    [{label}] '{old_frag[:30]}' → '{new_frag[:40]}'")
                return True
    # multi-run fallback
    for para in shape.text_frame.paragraphs:
        pt = "".join(r.text for r in para.runs)
        if old_frag in pt:
            if para.runs:
                para.runs[0].text = pt.replace(old_frag, new_frag)
                for r in para.runs[1:]:
                    r.text = ""
            if label:
                print(f"    [{label}] (multi-run) '{old_frag[:30]}'")
            return True
    return False


def set_text_safe(shape, new_text: str) -> None:
    """Replace all text in shape preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    runs = first_para._p.findall(f'{ns}r')
    if runs:
        runs[0].find(f'{ns}t').text = new_text
        for r in runs[1:]:
            t = r.find(f'{ns}t')
            if t is not None:
                t.text = ""
    else:
        run = first_para.add_run()
        run.text = new_text
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def apply_normautofit(shape):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


def expand_shape_height(shape, new_h_emu: int) -> None:
    """Increase shape height if new_h_emu > current height (never shrink)."""
    sp = shape._element
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    xfrm = sp.find(f'.//{ns}xfrm')
    if xfrm is None:
        # try spPr
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            return
        xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        return
    cur_h = int(ext.get('cy', 0))
    if new_h_emu > cur_h:
        ext.set('cy', str(new_h_emu))


def expand_shape_width(shape, new_w_emu: int) -> None:
    """Increase shape width if new_w_emu > current width."""
    sp = shape._element
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    xfrm = sp.find(f'.//{ns}xfrm')
    if xfrm is None:
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            return
        xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        return
    cur_w = int(ext.get('cx', 0))
    if new_w_emu > cur_w:
        ext.set('cx', str(new_w_emu))


def audit_overflows(prs, label=""):
    """Print shapes whose text or bottom edge overflows slide bounds."""
    slide_h = prs.slide_height / 914400
    slide_w = prs.slide_width / 914400
    issues = []
    for si, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            bot = (sh.top + sh.height) / 914400
            right = (sh.left + sh.width) / 914400
            if bot > slide_h + 0.02 or right > slide_w + 0.02:
                issues.append(
                    f"  S{si+1} shape[{list(slide.shapes).index(sh)}] '{sh.name}': "
                    f"bot={bot:.3f}\" right={right:.3f}\" (slide {slide_h:.2f}\"×{slide_w:.2f}\")"
                )
    if issues:
        print(f"\n[OVERFLOW AUDIT{' '+label if label else ''}] {len(issues)} issue(s):")
        for i in issues:
            print(i)
    else:
        print(f"\n[OVERFLOW AUDIT{' '+label if label else ''}] Clean — no overflows.")
    return len(issues) == 0


# ── Load ───────────────────────────────────────────────────────────────────────
prs = Presentation(SRC)

# EMU constants
EMU_IN = 914400
H_SUBLABEL_EXPANDED = int(0.42 * EMU_IN)   # expanded primary sublabel height
H_EV_EXPANDED      = int(0.46 * EMU_IN)   # expanded evidence box height (stays within slide_h 13.33")
W_SUBLABEL_PRI     = int(3.30 * EMU_IN)   # widen primary sublabel across KPI area
# (spans from L=0.58" to nearly end of OFFTAKE at R=3.55" → width 2.97"; use 3.3" to reach GAP edge area)

# ── Slide 1 (national) — compact Qty + YoY ─────────────────────────────────────
print("=== Slide 1 — National Qty + YoY ===")
s1 = prs.slides[0]
patch(s1.shapes[6],
      '₹47.02 Cr primary  •  ₹10.92 Cr gap',
      f'₹47.02 Cr primary (+98.9% YoY) | {int(NAT_PRI26_QTY)}K units  •  ₹10.92 Cr gap',
      'S1:6 primary+YoY+qty')
patch(s1.shapes[5],
      '₹36.1 Cr July offtake',
      f'₹36.1 Cr July offtake | {int(NAT_OFF26_QTY)}K units',
      'S1:5 offtake+qty')
apply_normautofit(s1.shapes[5])
apply_normautofit(s1.shapes[6])
expand_shape_height(s1.shapes[5], int(0.45 * EMU_IN))
expand_shape_height(s1.shapes[6], int(0.40 * EMU_IN))

# ── Zone slides 5–10 ──────────────────────────────────────────────────────────
for zone_name, Z in ZONE.items():
    si = Z['slide_idx']
    slide = prs.slides[si]
    shapes = slide.shapes
    print(f"\n=== Slide {si+1} ({zone_name}) ===")

    # 1. Primary sublabel (shape 8): compact "+YoY% | QtyK units"
    #    Original text "July billing" → replace with compact two-part label
    pri_sub = f"+{Z['pri_yoy_pct']}% YoY | {int(Z['pq26'])}K units"
    ok = patch(shapes[8], "July billing", pri_sub, f"S{si+1}:8 pri sub")
    if not ok:
        print(f"  WARN S{si+1}:8 'July billing' not found; current='{get_text(shapes[8])[:40]}'")
    # Expand width to span across primary column (1.25"→1.7") and height (0.23"→0.40")
    expand_shape_width(shapes[8], int(1.70 * EMU_IN))
    expand_shape_height(shapes[8], int(0.38 * EMU_IN))
    apply_normautofit(shapes[8])

    # 2. Offtake sublabel (shape 13): keep compact — just add Qty if not already there
    old_off_sub = f"{Z['mix_pct']}% mix"
    new_off_sub = f"{Z['mix_pct']}% mix | {int(Z['oq26'])}K units"
    ok = patch(shapes[13], old_off_sub, new_off_sub, f"S{si+1}:13 off sub")
    if not ok:
        print(f"  WARN S{si+1}:13 '{old_off_sub}' not found; current='{get_text(shapes[13])[:40]}'")
    expand_shape_height(shapes[13], int(0.38 * EMU_IN))
    apply_normautofit(shapes[13])

    # 3. Chain labels [33][34][35]: append compact YoY tag
    chain_shape_idxs = [33, 34, 35]
    for ci, (chain_idx, (ch_name, p25, p26, yoy)) in enumerate(
            zip(chain_shape_idxs, Z['chains'])):
        sh = shapes[chain_idx]
        cur = get_text(sh)
        tag = f" |+{yoy:.0f}%" if yoy is not None else " |NEW"
        if "|+" not in cur and "|NEW" not in cur:
            set_text_safe(sh, cur + tag)
            print(f"    [{zone_name} chain{ci+1}] '{cur[:35]}' → appended '{tag}'")
        apply_normautofit(sh)
        expand_shape_height(sh, int(0.35 * EMU_IN))

    # 4. Insight 01: add offtake qty — compact version
    ins01_old = f"₹{Z['off26']} Cr offtake, {Z['mix_pct']}% of national value."
    ins01_new = f"₹{Z['off26']}Cr offtake ({int(Z['oq26'])}K units), {Z['mix_pct']}% national."
    ok = patch(shapes[Z['ins01_idx']], ins01_old, ins01_new, f"S{si+1}:ins01")
    if not ok:
        # Try without trailing period variant
        ins01_old2 = f"₹{Z['off26']} Cr offtake, {Z['mix_pct']}% of national value"
        patch(shapes[Z['ins01_idx']], ins01_old2, ins01_new, f"S{si+1}:ins01 v2")
    apply_normautofit(shapes[Z['ins01_idx']])

    # 5. Insight 02: compact primary+offtake YoY — fits in 1.67"×0.47"
    ins02_new = (
        f"Primary +{Z['pri_yoy_pct']}% YoY: ₹{Z['pri26']}Cr ({int(Z['pq26'])}K units) "
        f"vs offtake ₹{Z['off26']}Cr ({int(Z['oq26'])}K)."
    )
    # Try both "9.8 Cr" and "9.80 Cr" style (Python float formatting varies)
    ins02_old = f"Primary is ₹{Z['pri26']} Cr versus ₹{Z['off26']} Cr offtake."
    ok = patch(shapes[Z['ins02_idx']], ins02_old, ins02_new, f"S{si+1}:ins02")
    if not ok:
        # Try formatted variant with trailing zero (e.g. 9.80 instead of 9.8)
        ins02_old2 = f"Primary is ₹{Z['pri26']:.2f} Cr versus ₹{Z['off26']:.2f} Cr offtake."
        ok = patch(shapes[Z['ins02_idx']], ins02_old2, ins02_new, f"S{si+1}:ins02 (2dp)")
    if not ok:
        # Final fallback: patch on "Primary is" prefix (unique per slide)
        cur = get_text(shapes[Z['ins02_idx']])
        if cur.startswith("Primary is "):
            set_text_safe(shapes[Z['ins02_idx']], ins02_new)
            print(f"    [S{si+1}:ins02 fallback] replaced '{cur[:40]}'")
        else:
            print(f"  WARN S{si+1}:ins02 not patched; current='{cur[:50]}'")
    apply_normautofit(shapes[Z['ins02_idx']])

    # 6. Evidence box: replace EVIDENCE text with concise chain YoY brief
    ev_sh = shapes[Z['ev_idx']]
    cur_ev = get_text(ev_sh)
    if "[PRIMARY YoY]" not in cur_ev and "Prim YoY:" not in cur_ev:
        # Prepend to existing evidence text (keep first ~60 chars of existing)
        existing_short = cur_ev[:55].rstrip() + ("…" if len(cur_ev) > 55 else "")
        new_ev = f"{Z['ev_brief']}. {existing_short}"
        set_text_safe(ev_sh, new_ev)
        print(f"    [{zone_name} evidence] set to: '{new_ev[:70]}'")
    elif "[PRIMARY YoY]" in cur_ev:
        # Already has old-style prefix — replace with clean brief
        patch(ev_sh, cur_ev, Z['ev_brief'], f"{zone_name} ev replace")
    apply_normautofit(ev_sh)
    # Expand evidence box height so it doesn't overflow the container
    expand_shape_height(ev_sh, H_EV_EXPANDED)

# ── Overflow audit before save ─────────────────────────────────────────────────
audit_overflows(prs, "pre-save")

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved → {DST}")


# ── Validate ──────────────────────────────────────────────────────────────────
def validate(path):
    import zipfile
    from lxml import etree as _et

    # ZIP + XML integrity
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            data = z.read(name)
            if name.endswith('.xml') or name.endswith('.rels'):
                try:
                    _et.fromstring(data)
                except _et.XMLSyntaxError as e:
                    print(f"  XML ERROR in {name}: {e}")
                    return False

    p = Presentation(path)
    checks = [
        (0,  5,  '2114K units'),
        (0,  6,  '+98.9% YoY'),
        (4,  8,  '+99.0% YoY'),
        (4, 13,  '466K units'),
        (5,  8,  '+63.6% YoY'),
        (5, 33,  '|+188%'),
        (6,  8,  '+134.3% YoY'),
        (6, 33,  '|+201%'),
        (7,  8,  '+69.3% YoY'),
        (7, 33,  '|+269%'),
        (8,  8,  '+115.7% YoY'),
        (8, 33,  '|+119%'),
        (9,  8,  '+202.2% YoY'),
        (9, 33,  '|+169%'),
    ]
    print(f"\n=== VALIDATE {path} ===")
    ok = True
    for si, shi, exp in checks:
        sh = p.slides[si].shapes[shi]
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        status = "OK" if exp in txt else "FAIL"
        if status == "FAIL":
            ok = False
        print(f"  S{si+1}:{shi} [{status}] expect '{exp}' → got '{txt[:70]}'")
    return ok


valid = validate(DST)
print(f"\n{'ALL CHECKS PASSED' if valid else 'SOME CHECKS FAILED'}")
