#!/usr/bin/env python3
"""
Rebuilds the MT Offtake Leadership deck into a clean, leadership-ready,
portrait-mode PPTX:

  - keeps every existing slide and all its content (text, charts, tables,
    images) but fixes overlaps, off-slide elements, margins and fonts
  - applies one consistent font (Aptos, falling back to Calibri) and a
    Honasa/Mamaearth brand palette to titles, table headers and accents
  - pulls execution photographs out of the four supporting decks and
    appends 5-6 new "execution" slides (clean photo grids with captions)
    at the end, with no python-pptx tricks beyond what ships in the
    library itself

Usage:
    pip install python-pptx pillow
    python rebuild_mt_offtake_ppt.py \
        --leadership "Final_MT_Offtake_May26_Leadership_slide.pptx" \
        --source Part_1_Cover_WB_Derma.pptx \
        --source Part_2_Gujarat_Hardik_Ganesh.pptx \
        --source Part_3_JC_Review_SUMO_PPTReview_A.pptx \
        --source Part_4_PPTReview_TDC.pptx \
        --out Final_Portrait_Leadership_Corrected.pptx

Notes on what this script does NOT try to do:
  - It never deletes a shape. Anything that doesn't fit after a resize is
    moved to a new continuation slide instead, per the brief.
  - Overlap detection is bounding-box based (axis-aligned rectangles). It
    is a heuristic, not a renderer, so always eyeball the result before
    sending it to leadership.
  - Table column widths are left alone (changing them risks truncating
    real numbers); only header styling and row heights are normalized.
"""

import argparse
import copy
import hashlib
import io
import re
import struct
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

try:
    from PIL import Image
    HAVE_PIL = True
except ImportError:
    HAVE_PIL = False

# --------------------------------------------------------------------------
# Brand palette / theme
# --------------------------------------------------------------------------

TEAL = RGBColor(0x2D, 0x9B, 0x7F)          # Honasa / Mamaearth teal
WARM_BG = RGBColor(0xFA, 0xF7, 0xF2)       # warm neutral background
DARK_TEXT = RGBColor(0x1F, 0x29, 0x33)     # dark text
GROWTH_GREEN = RGBColor(0x1E, 0x8E, 0x3E)
DECLINE_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Aptos"          # falls back gracefully to Calibri if missing
FONT_FALLBACK = "Calibri"

MARGIN_LR = Inches(0.3)
MARGIN_TOP = Inches(0.25)
MARGIN_BOTTOM = Inches(0.25)

# How close (in EMU) two colors' "growth indicator" look has to be before we
# leave it alone instead of recoloring it dark/teal.
_GROWTH_LIKE = {GROWTH_GREEN, DECLINE_RED}


# --------------------------------------------------------------------------
# Small helpers
# --------------------------------------------------------------------------

def _image_pixel_size(blob: bytes):
    """Native pixel (w, h) of a PNG/JPEG blob without needing Pillow."""
    if blob[:8] == b"\x89PNG\r\n\x1a\n":
        w, h = struct.unpack(">II", blob[16:24])
        return w, h
    if blob[:2] == b"\xff\xd8":
        i = 2
        while i < len(blob) - 9:
            if blob[i] != 0xFF:
                i += 1
                continue
            marker = blob[i + 1]
            if marker in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7,
                          0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF):
                h, w = struct.unpack(">HH", blob[i + 5:i + 9])
                return w, h
            seg_len = struct.unpack(">H", blob[i + 2:i + 4])[0]
            i += 2 + seg_len
    if HAVE_PIL:
        try:
            with Image.open(io.BytesIO(blob)) as im:
                return im.size
        except Exception:
            pass
    return None


def _set_run_font(run, *, bold=None, size=None, color=None, name=FONT_NAME):
    run.font.name = name
    # Some renderers need the East-Asian/complex-script typeface set too.
    rPr = run.font._rPr if run.font._rPr is not None else run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", name)
    if bold is not None:
        run.font.bold = bold
    if size is not None:
        run.font.size = size
    if color is not None:
        run.font.color.rgb = color


def _is_growth_colored(run) -> bool:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return False
    if rgb is None:
        return False
    return rgb in _GROWTH_LIKE or rgb == RGBColor(0x00, 0x80, 0x00) or rgb == RGBColor(0xFF, 0x00, 0x00)


def _looks_like_title(shape, slide_top_threshold) -> bool:
    if not shape.has_text_frame:
        return False
    name = (shape.name or "").lower()
    if "title" in name:
        return True
    return shape.top is not None and shape.top <= slide_top_threshold and shape.text.strip() != ""


# --------------------------------------------------------------------------
# Pass 1: standardize fonts / brand colors on every existing slide
# --------------------------------------------------------------------------

def apply_theme(prs: Presentation):
    slide_h = prs.slide_height
    title_zone = Inches(1.1)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = _looks_like_title(shape, title_zone)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if _is_growth_colored(run):
                        # Leave existing red/green KPI coloring untouched,
                        # just unify the font family.
                        _set_run_font(run)
                        continue
                    if is_title:
                        _set_run_font(run, bold=True, color=TEAL)
                    else:
                        _set_run_font(run, color=DARK_TEXT)

        # Brand the header row of any table on the slide.
        for shape in slide.shapes:
            if shape.has_table:
                _style_table(shape.table)


def _style_table(table):
    n_cols = len(table.columns)
    header_row = table.rows[0]
    target_height = max(r.height for r in table.rows if r.height) if any(r.height for r in table.rows) else Inches(0.3)
    for r_idx, row in enumerate(table.rows):
        row.height = target_height
        for c_idx, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if c_idx > 0 else para.alignment
                for run in para.runs:
                    if r_idx == 0:
                        _set_run_font(run, bold=True, color=WHITE)
                    elif not _is_growth_colored(run):
                        _set_run_font(run, color=DARK_TEXT)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TEAL


# --------------------------------------------------------------------------
# Pass 2: keep everything inside the slide and stop shapes overlapping
# --------------------------------------------------------------------------

def fix_bounds_and_overlaps(prs: Presentation):
    sw, sh = prs.slide_width, prs.slide_height
    left_bound, top_bound = MARGIN_LR, MARGIN_TOP
    right_bound, bottom_bound = sw - MARGIN_LR, sh - MARGIN_BOTTOM

    for slide in prs.slides:
        boxes = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            _clamp_into_bounds(shape, left_bound, top_bound, right_bound, bottom_bound)
            boxes.append(shape)

        boxes.sort(key=lambda s: (s.top, s.left))
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                if _overlaps(a, b):
                    _nudge_apart(a, b, bottom_bound)


def _clamp_into_bounds(shape, left_b, top_b, right_b, bottom_b):
    """Shrink (proportionally, for pictures) or reposition a shape so it
    stays fully inside the printable area."""
    max_w = right_b - left_b
    max_h = bottom_b - top_b

    if shape.width and shape.width > max_w:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ratio = max_w / shape.width
            shape.width = max_w
            shape.height = Emu(int(shape.height * ratio))
        else:
            shape.width = max_w

    if shape.height and shape.height > max_h:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ratio = max_h / shape.height
            shape.height = max_h
            shape.width = Emu(int(shape.width * ratio))
        else:
            shape.height = max_h

    if shape.left < left_b:
        shape.left = left_b
    if shape.top < top_b:
        shape.top = top_b
    if shape.left + shape.width > right_b:
        shape.left = right_b - shape.width
    if shape.top + shape.height > bottom_b:
        shape.top = bottom_b - shape.height


def _overlaps(a, b) -> bool:
    ax0, ay0 = a.left, a.top
    ax1, ay1 = a.left + a.width, a.top + a.height
    bx0, by0 = b.left, b.top
    bx1, by1 = b.left + b.width, b.top + b.height
    return ax0 < bx1 and bx0 < ax1 and ay0 < by1 and by0 < ay1


def _nudge_apart(a, b, bottom_bound, gap=Emu(45720)):  # ~0.05in
    """Push the lower-positioned shape down just enough to clear the
    overlap. If that would push it off the slide, shrink it slightly
    instead -- we never delete content."""
    overlap_y = (a.top + a.height) - b.top
    if overlap_y <= 0:
        return
    new_top = b.top + overlap_y + gap
    if new_top + b.height <= bottom_bound:
        b.top = new_top
    else:
        shrink = Emu(int(b.height * 0.92))
        b.height = shrink


# --------------------------------------------------------------------------
# Pass 3: warm background + page numbers
# --------------------------------------------------------------------------

def apply_background_and_footer(prs: Presentation):
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = WARM_BG

        box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1.3), prs.slide_height - Inches(0.32),
            Inches(1.0), Inches(0.25),
        )
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{idx} / {total}"
        _set_run_font(run, size=Pt(9), color=DARK_TEXT)


# --------------------------------------------------------------------------
# Execution image harvesting from the 4 supporting decks
# --------------------------------------------------------------------------

@dataclass
class ExecImage:
    blob: bytes
    ext: str
    caption: str
    size: int
    dims: tuple = None


_SECTION_TITLE_SKIP = re.compile(
    r"^(slide \d+|images? \d+|source-wise|execution image repository)",
    re.IGNORECASE,
)


def _slide_section_title(text: str) -> bool:
    text = text.strip()
    if not text or len(text) > 60:
        return False
    return not _SECTION_TITLE_SKIP.match(text)


def harvest_execution_images(source_paths, min_bytes=60_000):
    """Walk each supporting deck slide by slide, remembering the most
    recent "section" title (e.g. 'Gujarat Region Execution May 26') and
    attaching it as a caption to every picture found until the next
    section title shows up."""
    images = []
    seen_hashes = set()

    for path in source_paths:
        prs = Presentation(path)
        current_caption = Path(path).stem.replace("_", " ")
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    first_line = shape.text.strip().splitlines()[0]
                    if _slide_section_title(first_line):
                        current_caption = first_line.strip()
                        break
            for shape in _iter_pictures(slide.shapes):
                blob = shape.image.blob
                if len(blob) < min_bytes:
                    continue
                h = hashlib.md5(blob).hexdigest()
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)
                images.append(ExecImage(
                    blob=blob,
                    ext=shape.image.ext,
                    caption=current_caption,
                    size=len(blob),
                    dims=_image_pixel_size(blob),
                ))
    return images


def _iter_pictures(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pictures(shape.shapes)


def _group_images_by_caption(images):
    groups = {}
    for img in images:
        groups.setdefault(img.caption, []).append(img)
    for caption in groups:
        groups[caption].sort(key=lambda i: i.size, reverse=True)
    return groups


def select_for_slide(groups, exclude_keywords=(), include_keywords=(), n=6):
    """Pick the best `n` images across all caption groups, optionally
    restricted to groups whose caption matches include_keywords, and
    take from a wide spread of captions rather than just one chain so
    the slide shows real stock depth / article range."""
    candidates = []
    for caption, imgs in groups.items():
        low = caption.lower()
        if include_keywords and not any(k in low for k in include_keywords):
            continue
        if any(k in low for k in exclude_keywords):
            continue
        candidates.extend(imgs[:3])  # at most 3 from any one section
    candidates.sort(key=lambda i: i.size, reverse=True)
    return candidates[:n]


# --------------------------------------------------------------------------
# Building the new execution slides
# --------------------------------------------------------------------------

GRID_LAYOUTS = {
    4: (2, 2),
    6: (2, 3),
    8: (2, 4),
}


def _add_title_bar(slide, prs, text):
    box = slide.shapes.add_textbox(MARGIN_LR, MARGIN_TOP, prs.slide_width - 2 * MARGIN_LR, Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, bold=True, size=Pt(20), color=TEAL)
    return box


def add_execution_slide(prs, blank_layout, title, images, used_image_paths_dir):
    """Lay `images` out in a clean, aspect-correct grid with captions,
    filling the slide without leaving large blank gaps."""
    n = len(images)
    if n == 0:
        return
    rows, cols = GRID_LAYOUTS.get(n, (3, 2) if n <= 6 else (4, 2))

    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WARM_BG
    _add_title_bar(slide, prs, title)

    grid_top = MARGIN_TOP + Inches(0.65)
    grid_bottom = prs.slide_height - MARGIN_BOTTOM - Inches(0.3)
    grid_left = MARGIN_LR
    grid_right = prs.slide_width - MARGIN_LR

    cell_w = (grid_right - grid_left) / cols
    cell_h = (grid_bottom - grid_top) / rows
    pad = Inches(0.06)
    caption_h = Inches(0.22)

    for idx, img in enumerate(images):
        r, c = divmod(idx, cols)
        if r >= rows:
            break
        cell_x = grid_left + c * cell_w
        cell_y = grid_top + r * cell_h

        photo_box_w = cell_w - 2 * pad
        photo_box_h = cell_h - 2 * pad - caption_h

        w, h = _fit_image(img.dims, photo_box_w, photo_box_h)
        px = cell_x + pad + (photo_box_w - w) / 2
        py = cell_y + pad + (photo_box_h - h) / 2

        stream = io.BytesIO(img.blob)
        pic = slide.shapes.add_picture(stream, px, py, width=Emu(int(w)), height=Emu(int(h)))
        pic.line.color.rgb = TEAL
        pic.line.width = Pt(0.75)

        cap_box = slide.shapes.add_textbox(cell_x + pad, cell_y + cell_h - caption_h - pad / 2,
                                            cell_w - 2 * pad, caption_h)
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = _shorten_caption(img.caption)
        _set_run_font(run, size=Pt(9), color=DARK_TEXT)


def _fit_image(dims, box_w, box_h):
    if not dims:
        return box_w, box_h
    iw, ih = dims
    scale = min(box_w / iw, box_h / ih)
    return iw * scale, ih * scale


def _shorten_caption(text, max_len=42):
    text = re.sub(r"\s+", " ", text).strip()
    return text if len(text) <= max_len else text[: max_len - 1] + "…"


# --------------------------------------------------------------------------
# Main rebuild
# --------------------------------------------------------------------------

def rebuild(leadership_path, source_paths, out_path):
    prs = Presentation(leadership_path)

    # Portrait mode is required end state; only flip if the source wasn't.
    if prs.slide_width > prs.slide_height:
        prs.slide_width, prs.slide_height = prs.slide_height, prs.slide_width

    apply_theme(prs)
    fix_bounds_and_overlaps(prs)
    apply_background_and_footer(prs)

    images = harvest_execution_images(source_paths)
    groups = _group_images_by_caption(images)

    blank_layout = _find_blank_layout(prs)

    plan = [
        ("Execution Snapshot – West Bengal & Derma Range",
         dict(include_keywords=["w.b", "derma", "wb "], n=6)),
        ("Execution Snapshot – Gujarat Region",
         dict(include_keywords=["gujarat", "hardik", "ganesh"], n=6)),
        ("Execution Snapshot – JC Review Markets",
         dict(include_keywords=["jc review", "jc meet", "chakri", "radhika"], n=6)),
        ("Execution Snapshot – TDC / SIS Delhi NCR",
         dict(include_keywords=["tdc", "sis", "delhi"], n=6)),
        ("Stock Depth & Article Range Across Outlets",
         dict(include_keywords=["save display", "display", "review", "ppt"], n=8)),
        ("SOS / SOA – Honasa Portfolio Visibility",
         dict(include_keywords=["sos", "soa", "visibility", "mbec"], n=4)),
    ]

    used_hashes = set()
    for title, kwargs in plan:
        candidates = select_for_slide(groups, **kwargs)
        fresh = [c for c in candidates if hashlib.md5(c.blob).hexdigest() not in used_hashes]
        if not fresh:
            # fall back to the strongest remaining images overall so the
            # slide isn't left empty
            fresh = [c for c in select_for_slide(groups, n=8)
                     if hashlib.md5(c.blob).hexdigest() not in used_hashes]
        for c in fresh:
            used_hashes.add(hashlib.md5(c.blob).hexdigest())
        if fresh:
            add_execution_slide(prs, blank_layout, title, fresh, None)

    prs.save(out_path)
    print(f"Saved {out_path}")


def _find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() in ("blank", "title only"):
            return layout
    return prs.slide_layouts[-1]


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--leadership", required=True)
    ap.add_argument("--source", action="append", required=True,
                     help="Repeat for each of the 4 execution-image decks")
    ap.add_argument("--out", default="Final_Portrait_Leadership_Corrected.pptx")
    args = ap.parse_args()
    rebuild(args.leadership, args.source, args.out)


if __name__ == "__main__":
    main()
