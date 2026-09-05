"""
Add 6 zone-specific slides to the existing 6-slide MT July26 final deck.
Extends MT_July26_Final_UPDATED_with_All3_Insights_v1.pptx from 6 to 12 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color palette (same as original deck)
NAVY = RGBColor(31, 41, 51)          # #1F2933
TEAL = RGBColor(45, 155, 127)        # #2D9B7F
GREY = RGBColor(107, 118, 130)       # #6B7682
GREEN = RGBColor(30, 142, 62)        # #1E8E3E
RED = RGBColor(192, 57, 43)          # #C0392B
ORANGE = RGBColor(199, 125, 23)      # #C77D17
BLUE = RGBColor(43, 108, 176)        # #2B6CB0
WHITE = RGBColor(255, 255, 255)

# Zone data extracted from File 2 (MT_Jul26_Honasa_Updated_v18.pptx, Slides 5-10)
ZONES_DATA = {
    "West": {
        "offtake": "₹8.28 Cr",
        "conversion": "82.3%",
        "gap": "₹1.78 Cr",
        "status": "WATCH",
        "status_color": ORANGE,
        "chains": [
            ("DMart", "₹5.76 Cr", "94.2%", "▲87% YoY"),
            ("Reliance", "₹1.23 Cr", "54.5%", "▲122% YoY"),
            ("Wellness Forever", "₹0.65 Cr", "134.9%", "NEW"),
        ],
        "categories": [
            ("Face Cleanser", "₹5.23 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹1.53 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.77 Cr", "Mamaearth + TDC", "▲stable"),
        ],
        "states": ("Maharashtra ₹3.43 Cr (41.4%)", "Gujarat ₹2.57 Cr (31.0%)", "Mumbai ₹2.28 Cr (27.5%)"),
        "insights": [
            "✓ West ranks #1 nationally (22.9% mix); DMart 94.2% proves discipline",
            "⚠ Reliance 54.5% conversion—₹1.78 Cr gap to close",
            "→ Primary +99% YoY validates positioning",
        ],
        "action_owner": "NKAM Reliance West + NKAM Wellness",
        "action_deadline": "30-Aug",
        "action_target": "Map Reliance EAN-store failures; restore top-10 articles",
    },
    "South-1": {
        "offtake": "₹8.19 Cr",
        "conversion": "83.6%",
        "gap": "₹1.61 Cr",
        "status": "WATCH",
        "status_color": ORANGE,
        "chains": [
            ("Apollo", "₹2.84 Cr", "81.0%", "▲188% YoY"),
            ("DMart", "₹2.31 Cr", "74.8%", "▲93% YoY"),
            ("Lulu", "₹1.22 Cr", "~75%", "▲96% YoY"),
        ],
        "categories": [
            ("Face Cleanser", "₹2.30 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹1.23 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.93 Cr", "Mamaearth + TDC", "▲stable"),
        ],
        "states": ("Karnataka ₹4.10 Cr (50.1%)", "Tamil Nadu ₹2.63 Cr (32.1%)", "Kerala ₹1.17 Cr (14.3%)"),
        "insights": [
            "✓ South-1 ranks #2 nationally (22.7% mix); Karnataka anchors 50.1%",
            "⚠ Apollo 3pp below benchmark; Karnataka loss > Tamil Nadu gain",
            "→ TDC Face Cleanser OSA underpins zone conversion",
        ],
        "action_owner": "NKAM Apollo S-1",
        "action_deadline": "28-Aug",
        "action_target": "Fix top-5 declining EAN-store pairs in Karnataka",
    },
    "North": {
        "offtake": "₹6.99 Cr",
        "conversion": "58.5%",
        "gap": "₹4.97 Cr",
        "status": "FIX",
        "status_color": RED,
        "chains": [
            ("DMart", "₹2.53 Cr", "77.9%", "▲201% YoY"),
            ("Reliance", "₹2.40 Cr", "44.9%", "▲179% YoY"),
            ("Apollo", "₹1.16 Cr", "98.3%", "▲149% YoY"),
        ],
        "categories": [
            ("Face Cleanser", "₹3.50 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹1.68 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.71 Cr", "Mamaearth + TDC", "▲stable"),
        ],
        "states": ("Delhi NCR ₹1.97 Cr (28.2%)", "Rajasthan ₹1.67 Cr (23.9%)", "Punjab ₹1.55 Cr (22.2%)"),
        "insights": [
            "⚠ LOWEST CONVERSION NATIONALLY (58.5%); ₹4.97 Cr gap",
            "⚠ Reliance 44.9% suggests Paisa Vasool loading—needs validation by 05-Aug",
            "✓ Primary +134% YoY (₹11.95 Cr); Apollo 98.3% shows real demand exists",
        ],
        "action_owner": "North ZSM + NKAM Reliance",
        "action_deadline": "05-Sep",
        "action_target": "Validate Reliance Aug conversion; if <80%, pull top-20 EAN-store pairs",
    },
    "South-2": {
        "offtake": "₹4.91 Cr",
        "conversion": "71.3%",
        "gap": "₹1.98 Cr",
        "status": "FIX",
        "status_color": RED,
        "chains": [
            ("DMart", "₹1.95 Cr", "45.1%", "▲269% YoY"),
            ("Apollo", "₹1.64 Cr", "148.5%", "FLAT—data risk"),
            ("Reliance", "₹0.67 Cr", "82.5%", "▲27% YoY"),
        ],
        "categories": [
            ("Face Cleanser", "₹1.60 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹0.94 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.34 Cr", "Mamaearth + TDC", "▲stable"),
        ],
        "states": ("Telangana ₹2.56 Cr (52.2%)", "Andhra Pradesh ₹2.35 Cr (47.9%)", "2-state zone = easier control"),
        "insights": [
            "⚠ DMart S-2 at 45.1% conversion—needs urgent recovery",
            "🔴 Apollo >100% conversion is a DATA RISK—likely loading into full pipeline",
            "→ Primary +69.3% YoY; Telangana/AP split (52/48) requires careful balance",
        ],
        "action_owner": "NKAM DMart S-2 + Analytics",
        "action_deadline": "28-Aug",
        "action_target": "DMart store-level audit by 28-Aug; quarantine Apollo primary until validated",
    },
    "East": {
        "offtake": "₹3.55 Cr",
        "conversion": "45.3%",
        "gap": "₹4.28 Cr",
        "status": "URGENT",
        "status_color": RED,
        "chains": [
            ("Reliance", "₹2.16 Cr", "52.9%", "▲119% YoY"),
            ("Apollo", "₹0.80 Cr", "121.1%", "NEW—data risk"),
            ("Vishal Mega Mart", "₹0.17 Cr", "~60%", "▲65% YoY"),
        ],
        "categories": [
            ("Face Cleanser", "₹1.81 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹1.18 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.28 Cr", "Mamaearth + TDC", "▲limited"),
        ],
        "states": ("West Bengal ₹1.56 Cr (43.8%)", "Odisha ₹0.60 Cr (16.9%)", "Bihar ₹0.49 Cr (13.8%)"),
        "insights": [
            "🔴 WORST CONVERSION NATIONALLY (45.3%); ₹4.28 Cr flow gap",
            "⚠ Over-loaded relative to demand; 10.2% NPI mix risks inventory buildup",
            "✓ Primary +115.7% YoY (₹7.83 Cr)—but conversion mismatch is critical",
        ],
        "action_owner": "East ZSM + NKAM Reliance East",
        "action_deadline": "05-Aug+",
        "action_target": "Validate Aug-26 Reliance East offtake by 05-Aug; if <70%, pull EAN-store list",
    },
    "Central": {
        "offtake": "₹2.12 Cr",
        "conversion": "78.8%",
        "gap": "₹0.57 Cr",
        "status": "WATCH",
        "status_color": ORANGE,
        "chains": [
            ("DMart", "₹1.41 Cr", "95.3%", "▲169% YoY"),
            ("Reliance", "₹0.46 Cr", "51.2%", "▲296% YoY"),
            ("Apollo", "₹0.19 Cr", "137.7%", "NEW"),
        ],
        "categories": [
            ("Face Cleanser", "₹1.20 Cr", "Mamaearth + TDC", "▲strong"),
            ("Shampoo", "₹0.40 Cr", "Mamaearth", "▲growth"),
            ("Sun Care", "₹0.27 Cr", "Mamaearth + TDC", "▲stable"),
        ],
        "states": ("Madhya Pradesh ₹1.68 Cr (79.3%)", "Chhattisgarh ₹0.44 Cr (20.8%)", "Clean new zone"),
        "insights": [
            "✓ Clean new zone; highest DMart conversion (95.3%); smallest gap (₹0.57 Cr)",
            "✓ Primary +202.2% YoY; maintains healthy 78.8% overall",
            "→ Protect early momentum; benchmark DMart EAN discipline for national rollout",
        ],
        "action_owner": "NKAM DMart Central + Analytics",
        "action_deadline": "31-Aug",
        "action_target": "Document Central DMart EAN ordering cadence; build state-level comparable data",
    },
}


def add_zone_slide(prs, zone_name, zone_data):
    """Add a single zone deep-dive slide to presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set background to light grey
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)

    # ============ SUBTITLE BAR ============
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.35))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = f"{zone_name} ZONE | Q1 FY27 | DEEP DIVE"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle background
    subtitle_fill = subtitle_box.fill
    subtitle_fill.solid()
    subtitle_fill.fore_color.rgb = TEAL

    # ============ HEADLINE (Zone Performance) ============
    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(6.5), Inches(0.5))
    headline_frame = headline_box.text_frame
    headline_frame.word_wrap = True
    p = headline_frame.paragraphs[0]
    p.text = f"{zone_data['offtake']} | {zone_data['conversion']} CONVERSION"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # Status badge (top right)
    status_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.75), Inches(2), Inches(0.4))
    status_frame = status_box.text_frame
    status_frame.word_wrap = False
    p = status_frame.paragraphs[0]
    p.text = zone_data['status']
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = zone_data['status_color']
    p.alignment = PP_ALIGN.RIGHT

    # Gap metric (secondary)
    gap_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(6.5), Inches(0.25))
    gap_frame = gap_box.text_frame
    p = gap_frame.paragraphs[0]
    p.text = f"Flow Gap: {zone_data['gap']}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    # ============ CHAINS SECTION ============
    chains_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(4.5), Inches(0.25))
    chains_frame = chains_title.text_frame
    p = chains_frame.paragraphs[0]
    p.text = "TOP CHAINS"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    chains_y = 1.95
    for i, (chain_name, offtake, conv, growth) in enumerate(zone_data['chains']):
        chain_box = slide.shapes.add_textbox(Inches(0.5), Inches(chains_y + i*0.3), Inches(4.5), Inches(0.28))
        chain_frame = chain_box.text_frame
        chain_frame.word_wrap = True

        p = chain_frame.paragraphs[0]
        p.text = f"{chain_name}"
        p.font.size = Pt(8.8)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT

        # Add metrics on same line
        metrics_box = slide.shapes.add_textbox(Inches(2.8), Inches(chains_y + i*0.3), Inches(2), Inches(0.28))
        metrics_frame = metrics_box.text_frame
        p = metrics_frame.paragraphs[0]
        p.text = f"{offtake} | {conv} conv"
        p.font.size = Pt(8.8)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT

        # Growth indicator
        growth_box = slide.shapes.add_textbox(Inches(4.8), Inches(chains_y + i*0.3), Inches(0.8), Inches(0.28))
        growth_frame = growth_box.text_frame
        p = growth_frame.paragraphs[0]
        p.text = growth
        p.font.size = Pt(8.8)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.LEFT

    # ============ CATEGORIES SECTION ============
    cat_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.65), Inches(4.3), Inches(0.25))
    cat_frame = cat_title.text_frame
    p = cat_frame.paragraphs[0]
    p.text = "CATEGORY PERFORMANCE"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    cat_y = 1.95
    for i, (cat_name, nsv, brands, growth) in enumerate(zone_data['categories']):
        cat_box = slide.shapes.add_textbox(Inches(5.2), Inches(cat_y + i*0.3), Inches(4.3), Inches(0.28))
        cat_frame = cat_box.text_frame
        cat_frame.word_wrap = True

        p = cat_frame.paragraphs[0]
        p.text = f"{cat_name}"
        p.font.size = Pt(8.8)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT

        # NSV
        nsv_box = slide.shapes.add_textbox(Inches(6.8), Inches(cat_y + i*0.3), Inches(1.2), Inches(0.28))
        nsv_frame = nsv_box.text_frame
        p = nsv_frame.paragraphs[0]
        p.text = nsv
        p.font.size = Pt(8.8)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT

        # Growth
        growth_box = slide.shapes.add_textbox(Inches(8.0), Inches(cat_y + i*0.3), Inches(1.5), Inches(0.28))
        growth_frame = growth_box.text_frame
        p = growth_frame.paragraphs[0]
        p.text = growth
        p.font.size = Pt(8.8)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.LEFT

    # ============ KEY STATES ============
    states_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.25))
    states_frame = states_title.text_frame
    p = states_frame.paragraphs[0]
    p.text = "KEY STATES"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    states_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.35))
    states_frame = states_box.text_frame
    states_frame.word_wrap = True
    p = states_frame.paragraphs[0]
    p.text = " | ".join(zone_data['states'])
    p.font.size = Pt(8.8)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # ============ INSIGHTS ============
    insights_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.95), Inches(9), Inches(0.25))
    insights_frame = insights_title.text_frame
    p = insights_frame.paragraphs[0]
    p.text = "INSIGHTS & WATCH-OUTS"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    insights_y = 4.25
    for i, insight in enumerate(zone_data['insights']):
        insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(insights_y + i*0.35), Inches(9), Inches(0.35))
        insight_frame = insight_box.text_frame
        insight_frame.word_wrap = True
        p = insight_frame.paragraphs[0]
        p.text = insight
        p.font.size = Pt(9)
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.LEFT

    # ============ ACTIONS (SO WHAT) ============
    action_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.35), Inches(9), Inches(0.25))
    action_frame = action_title.text_frame
    p = action_frame.paragraphs[0]
    p.text = "IMMEDIATE ACTIONS"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    action_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(5.35), Inches(9), Inches(0.25))
    action_bg.fill.solid()
    action_bg.fill.fore_color.rgb = TEAL
    action_bg.line.color.rgb = TEAL
    slide.shapes._spTree.remove(action_bg._element)
    slide.shapes._spTree.insert(2, action_bg._element)

    action_owner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.65), Inches(9), Inches(0.25))
    owner_frame = action_owner_box.text_frame
    owner_frame.word_wrap = True
    p = owner_frame.paragraphs[0]
    p.text = f"Owner: {zone_data['action_owner']}"
    p.font.size = Pt(8.8)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    action_deadline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.95), Inches(4.5), Inches(0.25))
    deadline_frame = action_deadline_box.text_frame
    p = deadline_frame.paragraphs[0]
    p.text = f"Deadline: {zone_data['action_deadline']}"
    p.font.size = Pt(8.8)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.LEFT

    action_target_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.25), Inches(9), Inches(0.5))
    target_frame = action_target_box.text_frame
    target_frame.word_wrap = True
    p = target_frame.paragraphs[0]
    p.text = f"Target: {zone_data['action_target']}"
    p.font.size = Pt(8.8)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # ============ PAGE NUMBER ============
    page_num_box = slide.shapes.add_textbox(Inches(9.2), Inches(6.85), Inches(0.5), Inches(0.2))
    page_frame = page_num_box.text_frame
    p = page_frame.paragraphs[0]
    p.text = f"{len(prs.slides)}/12"
    p.font.size = Pt(8)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.RIGHT


def main():
    """Load existing 6-slide deck and add 6 zone slides."""
    # Load existing deck
    input_file = "/home/user/mt-dashboard/MT_July26_Final_UPDATED_with_All3_Insights_v1.pptx"
    prs = Presentation(input_file)

    print(f"Loaded existing deck: {input_file}")
    print(f"Current slides: {len(prs.slides)}")

    # Add zone slides in order: West, South-1, North, South-2, East, Central
    zone_order = ["West", "South-1", "North", "South-2", "East", "Central"]
    for zone_name in zone_order:
        add_zone_slide(prs, zone_name, ZONES_DATA[zone_name])
        print(f"✓ Added {zone_name} zone slide ({len(prs.slides)}/12)")

    # Save expanded deck
    output_file = "/home/user/mt-dashboard/MT_July26_Final_with_Zone_DeepDives_v1.pptx"
    prs.save(output_file)

    print(f"\n✓ Saved expanded deck: {output_file}")
    print(f"Final slide count: {len(prs.slides)}/12")


if __name__ == "__main__":
    main()
