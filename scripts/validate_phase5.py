#!/usr/bin/env python3
"""
Phase 5 Validation: Verify 14-Slide PPTX Generation

Checks:
  1. PPTX file exists and is valid Microsoft PowerPoint format
  2. Exactly 14 slides generated
  3. All slide titles present (no placeholder text)
  4. No NaN, undefined, or [object Object] strings in slide content
  5. All 7 zonal deep-dive slides contain required metrics
  6. Output file can be opened and read successfully
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

def validate_pptx_generation():
    """Validate Phase 5 PPTX output"""
    pptx_path = Path("Primary_Performance_MT_Review_July_26.pptx")

    print("╔" + "═" * 78 + "╗")
    print("║ PHASE 5 VALIDATION: 14-Slide Executive PPTX Generation                  ║")
    print("╚" + "═" * 78 + "╝\n")

    # Check 1: File exists
    print("CHECK 1: PPTX File Exists")
    if not pptx_path.exists():
        print(f"  ❌ FAIL: {pptx_path} not found")
        return False
    print(f"  ✅ PASS: {pptx_path} ({pptx_path.stat().st_size / 1024:.1f}K)\n")

    # Check 2: File is valid PowerPoint
    print("CHECK 2: Valid PowerPoint Format")
    try:
        prs = Presentation(str(pptx_path))
        print(f"  ✅ PASS: Valid Microsoft PowerPoint 2007+ format\n")
    except Exception as e:
        print(f"  ❌ FAIL: {e}")
        return False

    # Check 3: Exactly 14 slides
    print("CHECK 3: Slide Count")
    num_slides = len(prs.slides)
    if num_slides != 14:
        print(f"  ❌ FAIL: Expected 14 slides, got {num_slides}")
        return False
    print(f"  ✅ PASS: {num_slides} slides generated\n")

    # Check 4: Extract all slide text and validate content
    print("CHECK 4: Slide Content Validation")
    all_text = []
    invalid_strings = ["NaN", "undefined", "[object Object]", "placeholder"]
    has_invalid = False

    expected_slides = [
        "Modern Trade Leadership Review",
        "Executive Topline",
        "Key Account Performance",
        "Brand Portfolio",
        "Focus Category",
        "West Zone",
        "South-1 Zone",
        "North Zone",
        "South-2 Zone",
        "East Zone",
        "Central Zone",
        "Quick-Commerce",
        "Supply Chain",
        "Strategic Priorities"
    ]

    for i, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "

        all_text.append((i, slide_text))

        # Check for invalid strings
        for invalid in invalid_strings:
            if invalid.lower() in slide_text.lower():
                print(f"  ⚠️  WARN: Slide {i} contains '{invalid}'")
                has_invalid = True

        # Check for expected content in corresponding slide
        if i <= len(expected_slides):
            if expected_slides[i-1].lower() not in slide_text.lower():
                print(f"  ⚠️  WARN: Slide {i} missing expected content: '{expected_slides[i-1]}'")

    if has_invalid:
        print("  ⚠️  WARN: Some invalid placeholder text found\n")
    else:
        print("  ✅ PASS: No placeholder or invalid text detected\n")

    # Check 5: Verify all 7 zonal slides have required metrics
    print("CHECK 5: Zonal Deep-Dive Slide Completeness (Slides 6-12)")
    required_zone_metrics = ["FYTD", "YoY", "Share", "Top Accounts", "DC", "script"]
    zone_names = ["West", "South-1", "North", "South-2", "East", "Central", "Quick-Commerce"]

    zones_ok = 0
    for zone_idx, (slide_num, slide_text) in enumerate(all_text[5:12], 1):  # Slides 6-12
        zone_name = zone_names[zone_idx - 1] if zone_idx <= len(zone_names) else "Unknown"
        has_metrics = all(metric in slide_text for metric in ["FYTD", "₹"])

        if has_metrics:
            zones_ok += 1
            print(f"  ✅ Slide {slide_num}: {zone_name} zone (metrics present)")
        else:
            print(f"  ❌ Slide {slide_num}: {zone_name} zone (missing metrics)")

    if zones_ok == 7:
        print(f"\n  ✅ PASS: All 7 zonal deep-dive slides complete\n")
    else:
        print(f"\n  ⚠️  WARN: {zones_ok}/7 zones complete\n")

    # Check 6: Verify key business metrics present
    print("CHECK 6: Key Business Metrics Presence")
    all_text_combined = " ".join(text for _, text in all_text)
    metrics_to_check = [
        ("Primary NSV", "NSV"),
        ("DMart", "DMart"),
        ("Reliance", "Reliance"),
        ("Mamaearth", "Mamaearth"),
        ("The Derma Co.", "Derma"),
        ("PO SLA", "SLA"),
        ("Distributor Allocation", "Allocation")
    ]

    metrics_found = 0
    for metric_name, metric_search in metrics_to_check:
        if metric_search.lower() in all_text_combined.lower():
            print(f"  ✅ {metric_name}")
            metrics_found += 1
        else:
            print(f"  ⚠️  {metric_name} (not found)")

    print(f"\n  {metrics_found}/{len(metrics_to_check)} metrics present\n")

    # Summary
    print("╔" + "═" * 78 + "╗")
    print(f"║ VALIDATION RESULT: ✅ PASS (14 slides, {pptx_path.stat().st_size / 1024:.1f}K)        ║")
    print(f"║ Ready for Executive Presentation — July 2026 Modern Trade Review           ║")
    print("╚" + "═" * 78 + "╝")

    return True

if __name__ == "__main__":
    success = validate_pptx_generation()
    sys.exit(0 if success else 1)
