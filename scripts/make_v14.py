"""Build V14 from V13 — add Qty + YoY decomposition to national + all 6 zone slides.

Changes per zone slide (slides 5–10):
  1. Primary sublabel → adds Jul'25/Jul'26 NSV YoY + Jul'26 Qty
  2. Offtake sublabel → adds Jul'26 Qty
  3. Insight #01 (overall delivery) → adds Qty
  4. Insight #02 (primary vs offtake) → adds Jul'25 primary + YoY Qty
  5. Chain evidence box → prepends chain-level primary YoY summary
  6. Chain top-3 labels → appended with primary YoY % tag

National slide 1: sublabel under Primary + Offtake KPIs updated with YoY and Qty.

Data source: Primary_Article_Monthly/primary_article_Jul_25.csv + Jul_26.csv
             Offtake_Monthly/offtake_store_article_Jul_26.csv
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv13.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv14.pptx"

# ── Pre-computed zone data (from build script analysis Jul'25 & Jul'26 CSVs) ──
# Primary: both years from CSVs; Offtake: Jul'26 only (no Jul'25 offtake CSV)
# qty_k = thousand units; nsv_cr = ₹ Crore
ZONE = {
    'West': {
        'slide_idx': 4,
        'pri25': 5.05, 'pq25': 315.0,
        'pri26': 10.05,'pq26': 597.8,
        'pri_yoy_pct': 99.0, 'pq_yoy_pct': 89.8,
        'off26': 8.28, 'oq26': 466.6,    # offtake NSV from PPT (authoritative)
        'mix_pct': 22.9,
        'ins01_idx': 55, 'ins02_idx': 58,
        'ev_idx': 101,
        'chains': [
            ('DMart',            3.27, 6.12, 87.0),   # chain, pri25, pri26, yoy% (ShipTo Jul'25)
            ('Reliance',         1.02, 2.26, 121.6),
            ('Apollo',           0.26, 0.62, 138.5),
        ],
        'chain_note': (
            "Primary YoY: DMart ₹3.27→₹6.12Cr (+87.0%); "
            "Reliance ₹1.02→₹2.26Cr (+121.6%); Apollo ₹0.26→₹0.62Cr (+138.5%); "
            "Wellness Forever NEW ₹0.49Cr. "
            "Reliance West conversion 54.5% vs primary +121.6% → stock build watch."
        ),
    },
    'South-1': {
        'slide_idx': 5,
        'pri25': 5.99, 'pq25': 349.0,
        'pri26': 9.80, 'pq26': 581.6,
        'pri_yoy_pct': 63.6, 'pq_yoy_pct': 66.6,
        'off26': 8.19, 'oq26': 435.3,
        'mix_pct': 22.7,
        'ins01_idx': 55, 'ins02_idx': 58,
        'ev_idx': 101,
        'chains': [
            ('Apollo',    1.22, 3.51, 187.7),   # ShipTo Jul'25 corrected
            ('DMart',     1.60, 3.09, 93.1),
            ('Reliance',  1.14, 2.24, 96.5),
        ],
        'chain_note': (
            "Primary YoY: Apollo ₹1.22→₹3.51Cr (+187.7%); "
            "DMart ₹1.60→₹3.09Cr (+93.1%); "
            "Reliance ₹1.14→₹2.24Cr (+96.5%). "
            "Apollo S-1 primary surge: validate offtake velocity before next loading."
        ),
    },
    'North': {
        'slide_idx': 6,
        'pri25': 5.10, 'pq25': 306.7,
        'pri26': 11.95,'pq26': 649.1,
        'pri_yoy_pct': 134.3, 'pq_yoy_pct': 111.6,
        'off26': 6.99, 'oq26': 377.6,
        'mix_pct': 19.4,
        'ins01_idx': 55, 'ins02_idx': 58,
        'ev_idx': 101,
        'chains': [
            ('Reliance',  1.78, 5.35, 200.6),
            ('DMart',     1.16, 3.24, 179.3),
            ('Apollo',    0.47, 1.17, 148.9),
        ],
        'chain_note': (
            "Primary YoY: Reliance ₹1.78→₹5.35Cr (+200.6%); "
            "DMart ₹1.16→₹3.24Cr (+179.3%); Apollo ₹0.47→₹1.17Cr (+148.9%). "
            "PAISA VASOOL EFFECT: Reliance North primary surged +200.6% LY Jul'25; "
            "LY Paisa Vasool event drove +46.4% offtake Aug'25 surge — "
            "this Jul'26 primary buildup likely repeating the pattern. "
            "Validate Reliance North Aug'26 offtake by 05-Aug before holding or loading."
        ),
    },
    'South-2': {
        'slide_idx': 7,
        'pri25': 4.07, 'pq25': 223.4,
        'pri26': 6.89, 'pq26': 343.8,
        'pri_yoy_pct': 69.3, 'pq_yoy_pct': 53.9,
        'off26': 4.91, 'oq26': 281.1,
        'mix_pct': 13.6,
        'ins01_idx': 54, 'ins02_idx': 57,   # South-2 has 110 shapes, indices shift
        'ev_idx': 100,
        'chains': [
            ('DMart',     1.17, 4.32, 269.2),  # ShipTo Jul'25 corrected
            ('Apollo',    1.10, 1.10, 0.0),    # flat
            ('Reliance',  0.64, 0.81, 26.6),
        ],
        'chain_note': (
            "Primary YoY: DMart ₹1.17→₹4.32Cr (+269.2%); "
            "Apollo ₹1.10→₹1.10Cr (flat); Reliance ₹0.64→₹0.81Cr (+26.6%). "
            "DMart S-2 primary surge +269.2%: Apollo flat and Reliance low growth → "
            "validate offtake velocity before next DMart loading."
        ),
    },
    'East': {
        'slide_idx': 8,
        'pri25': 3.63, 'pq25': 211.3,
        'pri26': 7.83, 'pq26': 409.4,
        'pri_yoy_pct': 115.7, 'pq_yoy_pct': 93.8,
        'off26': 3.55, 'oq26': 183.9,
        'mix_pct': 9.8,
        'ins01_idx': 55, 'ins02_idx': 58,
        'ev_idx': 101,
        'chains': [
            ('Reliance',  1.87, 4.09, 118.7),  # ShipTo Jul'25 corrected
            ('Vmart',     0.00, 2.00, None),    # NEW
            ('Apollo',    0.40, 0.66, 65.0),
        ],
        'chain_note': (
            "Primary YoY: Reliance ₹1.87→₹4.09Cr (+118.7%); "
            "Vmart NEW ₹2.00Cr; Apollo ₹0.40→₹0.66Cr (+65.0%). "
            "PAISA VASOOL EFFECT: Reliance East primary +118.7% mirrors LY Jul'25 pattern; "
            "LY Aug'25 Paisa Vasool drove massive sellthrough — "
            "East 45.3% conversion signals over-stock NOW, not demand. "
            "Moratorium on non-Hero EAN loading stands; check Aug'26 Reliance East offtake by 05-Aug."
        ),
    },
    'Central': {
        'slide_idx': 9,
        'pri25': 0.89, 'pq25': 55.4,
        'pri26': 2.69, 'pq26': 160.9,
        'pri_yoy_pct': 202.2, 'pq_yoy_pct': 190.4,
        'off26': 2.12, 'oq26': 160.6,
        'mix_pct': 5.9,
        'ins01_idx': 54, 'ins02_idx': 57,   # Central has 110 shapes
        'ev_idx': 100,
        'chains': [
            ('DMart',     0.55, 1.48, 169.1),
            ('Reliance',  0.23, 0.91, 295.7),
            ('Apollo',    0.00, 0.14, None),  # NEW
        ],
        'chain_note': (
            "Primary YoY: DMart ₹0.55→₹1.48Cr (+169.1%); "
            "Reliance ₹0.23→₹0.91Cr (+295.7%); Apollo NEW ₹0.14Cr. "
            "Central is a NEW zone (FY27 first full month); "
            "+202.2% primary YoY reflects base effect of zone launch — "
            "protect 95.3% DMart conversion as the national benchmark."
        ),
    },
}

# National primary Jul'25 total = ₹24.73 Cr; Jul'26 = ₹49.21 Cr; YoY +98.9%
# National qty: Jul'25 = 1460.8K; Jul'26 = 2742.6K; YoY +87.7%
NAT_PRI25_NSV  = 24.73
NAT_PRI25_QTY  = 1460.8
NAT_PRI26_NSV  = 49.21
NAT_PRI26_QTY  = 2742.6
NAT_OFF26_NSV  = 36.10   # from PPT (authoritative)
NAT_OFF26_QTY  = 2114.2  # from offtake CSV excl Reliance BC


# ── Helpers ────────────────────────────────────────────────────────────────────
def get_text(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def set_text(shape, new_text: str) -> None:
    """Replace all text in shape with new_text, preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Use first paragraph's first run for formatting; clear all others
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for r in first_para.runs[1:]:
            r.text = ""
    else:
        # No runs — add one
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _et
        r_el = _et.SubElement(first_para._p, _qn('a:r'))
        t_el = _et.SubElement(r_el, _qn('a:t'))
        t_el.text = new_text
    # Clear remaining paragraphs
    for para in tf.paragraphs[1:]:
        for r in para.runs:
            r.text = ""


def patch(shape, old_frag: str, new_frag: str, label=""):
    """In-place substitution anywhere in shape's runs."""
    if not shape.has_text_frame:
        return False
    full = shape.text_frame.text
    if old_frag not in full:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_frag in run.text:
                run.text = run.text.replace(old_frag, new_frag)
                if label:
                    print(f"    [{label}] patched '{old_frag}' → '{new_frag}'")
                return True
    # multi-run fallback
    for para in shape.text_frame.paragraphs:
        pt = "".join(r.text for r in para.runs)
        if old_frag in pt:
            if para.runs:
                para.runs[0].text = pt.replace(old_frag, new_frag)
                for r in para.runs[1:]:
                    r.text = ""
            if label:
                print(f"    [{label}] patched (multi-run)")
            return True
    return False


def apply_normautofit(shape):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


# ── Load ───────────────────────────────────────────────────────────────────────
prs = Presentation(SRC)

# ── Slide 1 (index 0): National — add Qty + YoY to KPI sublabels ───────────────
print("=== Slide 1 — National Qty + YoY ===")
s1 = prs.slides[0]
# Shape 8: "July billing" (primary sublabel) — actually shape[8]='July billing' is on zone slides
# Slide 1 primary sublabel: shape[8] says "Primary vs offtake" (chart header)
# Shape 6 = "₹47.02 Cr primary  •  ₹10.92 Cr gap  •  90.7% concentrated..."
# Add qty + YoY context to shape 6
patch(s1.shapes[6],
      '₹47.02 Cr primary  •  ₹10.92 Cr gap',
      f'₹47.02 Cr primary (+98.9% YoY) | {int(NAT_PRI26_QTY)}K units  •  ₹10.92 Cr gap',
      'S1:6 primary+qty+YoY')
# Shape 5: "₹36.1 Cr July offtake | 76.8% flow conversion" — add offtake qty
patch(s1.shapes[5],
      '₹36.1 Cr July offtake',
      f'₹36.1 Cr July offtake | {int(NAT_OFF26_QTY)}K units',
      'S1:5 offtake+qty')

apply_normautofit(s1.shapes[5])
apply_normautofit(s1.shapes[6])

# ── Zone slides 5–10 ──────────────────────────────────────────────────────────
for zone_name, Z in ZONE.items():
    si = Z['slide_idx']
    slide = prs.slides[si]
    shapes = slide.shapes
    print(f"\n=== Slide {si+1} ({zone_name}) ===")

    # 1. Primary sublabel (shape 8): "July billing" → add YoY + Qty
    pri_label = (
        f"Jul'25: ₹{Z['pri25']}Cr | Jul'26: ₹{Z['pri26']}Cr "
        f"(+{Z['pri_yoy_pct']}% YoY) | {int(Z['pq26'])}K units"
    )
    patch(shapes[8], "July billing", pri_label, f"S{si+1}:8 pri sublabel")
    apply_normautofit(shapes[8])

    # 2. Offtake sublabel (shape 13): "XX.X% mix" → add Jul'26 Qty
    old_off_sub = f"{Z['mix_pct']}% mix"
    new_off_sub = f"{Z['mix_pct']}% mix | {int(Z['oq26'])}K units (Jul'26)"
    patch(shapes[13], old_off_sub, new_off_sub, f"S{si+1}:13 off sublabel")
    apply_normautofit(shapes[13])

    # 3. Insight #01 — add qty
    ins01_old = f"₹{Z['off26']} Cr offtake, {Z['mix_pct']}% of national value."
    ins01_new = f"₹{Z['off26']} Cr offtake ({int(Z['oq26'])}K units), {Z['mix_pct']}% of national value."
    patch(shapes[Z['ins01_idx']], ins01_old, ins01_new, f"S{si+1}:ins01 qty")
    apply_normautofit(shapes[Z['ins01_idx']])

    # 4. Insight #02 — add primary qty + YoY
    ins02_old = f"Primary is ₹{Z['pri26']} Cr versus ₹{Z['off26']} Cr offtake."
    ins02_new = (
        f"Primary is ₹{Z['pri26']} Cr ({int(Z['pq26'])}K units, +{Z['pri_yoy_pct']}% YoY vs "
        f"₹{Z['pri25']}Cr Jul'25) versus ₹{Z['off26']} Cr offtake ({int(Z['oq26'])}K units)."
    )
    patch(shapes[Z['ins02_idx']], ins02_old, ins02_new, f"S{si+1}:ins02 YoY+qty")
    apply_normautofit(shapes[Z['ins02_idx']])

    # 5. Chain labels — append primary YoY tag
    chain_shapes = [33, 34, 35]
    for ci, (chain_shapes_idx, (ch_name, p25, p26, yoy)) in enumerate(
            zip(chain_shapes, Z['chains'])):
        sh = shapes[chain_shapes_idx]
        cur = get_text(sh)
        if yoy is not None:
            tag = f" | Prim: +{yoy:.0f}%"
        else:
            tag = " | Prim: NEW"
        if "| Prim:" not in cur:
            set_text(sh, cur + tag)
            print(f"    [{zone_name} chain{ci+1}] appended '{tag}' to '{cur[:40]}'")
        apply_normautofit(sh)

    # 6. Evidence box — prepend chain YoY summary
    ev_sh = shapes[Z['ev_idx']]
    cur_ev = get_text(ev_sh)
    prefix = f"[PRIMARY YoY] {Z['chain_note']} "
    if "[PRIMARY YoY]" not in cur_ev:
        ev_sh.text_frame.paragraphs[0].runs[0].text = prefix + cur_ev[:200] + "…"
        print(f"    [{zone_name} evidence] prepended chain YoY summary")
        apply_normautofit(ev_sh)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved → {DST}")


# ── Validate ──────────────────────────────────────────────────────────────────
def validate(path):
    p = Presentation(path)
    checks = [
        (0,  5,  '2114K units'),
        (0,  6,  '+98.9% YoY'),
        (4,  8,  "Jul'25: ₹5.05Cr"),
        (4, 13,  '466K units'),
        (5,  8,  "Jul'25: ₹5.99Cr"),
        (6,  8,  "Jul'25: ₹5.1Cr"),
        (6, Z_by_idx(6)['ev_idx'], '[PRIMARY YoY]'),
        (7,  8,  "Jul'25: ₹4.07Cr"),
        (8,  8,  "Jul'25: ₹3.63Cr"),
        (9,  8,  "Jul'25: ₹0.89Cr"),
    ]
    print(f"\n=== VALIDATE {path} ===")
    ok = True
    for si, shi, exp in checks:
        txt = p.slides[si].shapes[shi].text_frame.text.strip() if \
              p.slides[si].shapes[shi].has_text_frame else ""
        status = "OK" if exp in txt else "FAIL"
        if status == "FAIL":
            ok = False
        print(f"  S{si+1}:{shi} [{status}] expect '{exp}' → got '{txt[:80]}'")
    return ok


def Z_by_idx(idx):
    for z in ZONE.values():
        if z['slide_idx'] == idx:
            return z
    return {}


valid = validate(DST)
print(f"\n{'ALL CHECKS PASSED' if valid else 'SOME CHECKS FAILED'}")
