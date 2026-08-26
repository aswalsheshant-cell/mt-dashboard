"""Build V17 from V16 — deep QC fixes across all slides.

Issues fixed:
1. Slide 4: Table header shapes [6-12] hidden by full-width background shape [13]
   at same T=1.823" with higher Z-order → move [13] before [6] in spTree.
2. Slide 1: Shape [5] (offtake sublabel) H=0.45" overlaps shape [6] at T=1.729"
   → reduce H to 0.36" (bot=1.714", 0.015" clear).
3. Slide 2: Shape [109] footnote at T=12.85" overlaps evidence boxes at T=12.83"
   → move [109] up to T=12.50" (above evidence band).
4. Slide 2: Evidence content boxes [112,115,118,121] H=0.38" → expand to 0.50".
5. normAutofit on all resized text shapes.
6. Full overflow audit before save.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu
from lxml import etree
import copy, zipfile

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv16.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv17.pptx"

EMU_IN = 914400


def emu(inches: float) -> int:
    return int(inches * EMU_IN)


def get_xfrm(shape):
    sp = shape._element
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    return sp.find(f'.//{ns_a}xfrm')


def set_top(shape, top_emu: int):
    xfrm = get_xfrm(shape)
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    if off is not None:
        off.set('y', str(top_emu))


def set_height(shape, h_emu: int):
    xfrm = get_xfrm(shape)
    if xfrm is None:
        return
    ext = xfrm.find(qn('a:ext'))
    if ext is not None:
        ext.set('cy', str(h_emu))


def apply_normautofit(shape):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


def move_shape_before(slide, shape_to_move_idx: int, before_idx: int):
    """Move shape at shape_to_move_idx to appear just before shape at before_idx in spTree.

    Lower spTree position = lower Z-order (renders behind shapes added later).
    This makes the shape_to_move render behind the before_idx shape.
    """
    spTree = slide.shapes._spTree
    all_sps = list(spTree)

    # Find the actual XML elements
    shape_el = slide.shapes[shape_to_move_idx]._element
    before_el = slide.shapes[before_idx]._element

    # Remove shape_el from spTree
    spTree.remove(shape_el)

    # Find position of before_el in the updated spTree
    all_sps = list(spTree)
    before_pos = None
    for i, el in enumerate(all_sps):
        if el is before_el:
            before_pos = i
            break

    if before_pos is None:
        # Fallback: append
        spTree.append(shape_el)
        print("    WARNING: could not find before_el; appended to end")
        return

    # Insert before before_el
    before_el.addprevious(shape_el)
    print(f"    Moved shape[{shape_to_move_idx}] to appear before shape[{before_idx}] in spTree")


def audit_overflows(prs):
    slide_h = prs.slide_height / EMU_IN
    issues = []
    for si, slide in enumerate(prs.slides):
        for i, sh in enumerate(slide.shapes):
            bot = (sh.top + sh.height) / EMU_IN
            if bot > slide_h + 0.015:
                txt = sh.text_frame.text[:40] if sh.has_text_frame else ""
                issues.append(f"  S{si+1} shape[{i}] bot={bot:.3f}\" > {slide_h:.2f}\" | {txt!r}")
    if issues:
        print(f"[OVERFLOW] {len(issues)} issue(s):")
        for x in issues:
            print(x)
    else:
        print("[OVERFLOW] Clean — no overflows.")
    return len(issues) == 0


def validate_xml(path):
    with zipfile.ZipFile(path) as z:
        errors = []
        for name in z.namelist():
            if name.endswith('.xml') or name.endswith('.rels'):
                try:
                    etree.fromstring(z.read(name))
                except etree.XMLSyntaxError as e:
                    errors.append(f"  XML ERROR in {name}: {e}")
        if errors:
            for e in errors:
                print(e)
            return False
        print("[XML] All XML files valid.")
        return True


prs = Presentation(SRC)
slide_h_emu = int(prs.slide_height)
slide_h_in  = slide_h_emu / EMU_IN

print(f"Slide dimensions: {prs.slide_width/EMU_IN:.2f}\" W x {slide_h_in:.2f}\" H")
print(f"Total slides: {len(prs.slides)}")

# ─────────────────────────────────────────────────────────────────────────────
# FIX 1: Slide 1 (index 0) — shape [5] height regression
# Shape [5] (offtake sublabel) expanded in V15 to H=0.45", creating 0.075" overlap
# with shape [6] at T=1.729". Reduce back to H=0.36" → bot=1.714" (clears [6]).
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== FIX 1: Slide 1 — reduce shape[5] height to clear shape[6] ===")
slide1 = prs.slides[0]
sh5 = slide1.shapes[5]
sh6 = slide1.shapes[6]

t5 = sh5.top / EMU_IN
h5 = sh5.height / EMU_IN
t6 = sh6.top / EMU_IN
print(f"  Before: shape[5] T={t5:.3f}\" H={h5:.3f}\" bot={t5+h5:.3f}\"")
print(f"          shape[6] T={t6:.3f}\"")

new_h5 = emu(0.36)
set_height(sh5, new_h5)
apply_normautofit(sh5)

h5_new = sh5.height / EMU_IN
print(f"  After:  shape[5] H={h5_new:.3f}\" bot={t5+h5_new:.3f}\" (clears shape[6] at {t6:.3f}\")")

# ─────────────────────────────────────────────────────────────────────────────
# FIX 2: Slide 4 (index 3) — table header z-order
# Column header text shapes [6]-[12] at T=1.823" are hidden by blank background
# shape [13] at same T (full-width W=6.71", higher Z-order = renders on top).
# Fix: move shape [13] before shape [6] in spTree so header text renders on top.
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== FIX 2: Slide 4 — fix table header z-order ===")
slide4 = prs.slides[3]

# Inspect shapes before fix
for i in range(6, 15):
    if i >= len(slide4.shapes):
        break
    sh = slide4.shapes[i]
    t = sh.top / EMU_IN
    w = sh.width / EMU_IN
    txt = sh.text_frame.text.strip()[:30] if sh.has_text_frame else ""
    print(f"  shape[{i}] T={t:.3f}\" W={w:.3f}\" | '{txt}'")

# Move shape [13] (full-width background row) to appear before shape [6] in spTree
# This places it BELOW the header text shapes in rendering order
move_shape_before(slide4, shape_to_move_idx=13, before_idx=6)

print("  Header text shapes [6]-[12] now render on top of background shape [13]")

# ─────────────────────────────────────────────────────────────────────────────
# FIX 3: Slide 2 (index 1) — footnote shape overlaps evidence boxes
# Shape [109] "Diagnostic extensions are recommendations..." at T=12.85"
# overlaps evidence boxes at T=12.83". Move [109] up to T=12.50".
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== FIX 3: Slide 2 — move footnote above evidence band ===")
slide2 = prs.slides[1]

# Find shape [109] — the footnote disclaimer
if len(slide2.shapes) > 109:
    sh109 = slide2.shapes[109]
    t109 = sh109.top / EMU_IN
    txt109 = sh109.text_frame.text.strip()[:60] if sh109.has_text_frame else ""
    print(f"  shape[109] T={t109:.3f}\" | '{txt109}'")

    # Move to T=12.50" (above the 12.83" evidence band, with 0.33" clearance)
    set_top(sh109, emu(12.50))
    print(f"  Moved shape[109] from T={t109:.3f}\" → T=12.50\"")
else:
    print(f"  Slide 2 has only {len(slide2.shapes)} shapes — shape[109] not found, scanning for footnote...")
    # Fallback: find by text
    for i, sh in enumerate(slide2.shapes):
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if 'Diagnostic' in txt or 'extensions are recommendations' in txt:
                t = sh.top / EMU_IN
                print(f"  Found footnote at shape[{i}] T={t:.3f}\" | '{txt[:50]}'")
                set_top(sh, emu(12.50))
                print(f"  Moved to T=12.50\"")
                break

# ─────────────────────────────────────────────────────────────────────────────
# FIX 4: Slide 2 — expand evidence content boxes [112,115,118,121]
# These were missed in V16 (only zone slides 4-9 got evidence expansion).
# Expand from H=0.38" → 0.50" and apply normAutofit.
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== FIX 4: Slide 2 — expand evidence content boxes ===")
ev_indices = [112, 115, 118, 121]
expanded = 0
for idx in ev_indices:
    if idx >= len(slide2.shapes):
        print(f"  shape[{idx}] not found (slide has {len(slide2.shapes)} shapes)")
        continue
    sh = slide2.shapes[idx]
    t = sh.top / EMU_IN
    h = sh.height / EMU_IN
    txt = sh.text_frame.text.strip()[:40] if sh.has_text_frame else ""
    print(f"  shape[{idx}] T={t:.3f}\" H={h:.3f}\" | '{txt}'")
    if h < 0.52:
        set_height(sh, emu(0.50))
        apply_normautofit(sh)
        expanded += 1
        print(f"    → expanded to H=0.50\"")

# Also do a general sweep of slide 2 evidence boxes (in case indices shifted)
if expanded == 0:
    print("  Index-based approach found no shapes — scanning slide 2 for evidence content...")
    for i, sh in enumerate(slide2.shapes):
        t = sh.top / EMU_IN
        h = sh.height / EMU_IN
        w = sh.width / EMU_IN
        if t > 12.70 and w < 2.0 and h > 0.10 and h < 0.45:
            txt = sh.text_frame.text.strip()[:40] if sh.has_text_frame else ""
            print(f"  shape[{i}] T={t:.3f}\" H={h:.3f}\" | '{txt}'")
            set_height(sh, emu(0.50))
            apply_normautofit(sh)
            expanded += 1
            print(f"    → expanded to H=0.50\"")

print(f"  Expanded {expanded} evidence box(es) on slide 2")

# ─────────────────────────────────────────────────────────────────────────────
# SWEEP: normAutofit on ALL text shapes in evidence band across all slides
# (defensive pass to ensure no text is cut off)
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== SWEEP: normAutofit on evidence band shapes (T>12.50\") ===")
total_naf = 0
for si, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.top / EMU_IN
        if t > 12.50:
            apply_normautofit(sh)
            total_naf += 1
print(f"  Applied normAutofit to {total_naf} text shapes in evidence band")

# ─────────────────────────────────────────────────────────────────────────────
# OVERFLOW AUDIT
# ─────────────────────────────────────────────────────────────────────────────
print()
overflow_ok = audit_overflows(prs)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved → {DST}")

# ─────────────────────────────────────────────────────────────────────────────
# XML VALIDATION
# ─────────────────────────────────────────────────────────────────────────────
xml_ok = validate_xml(DST)

# ─────────────────────────────────────────────────────────────────────────────
# SPOT CHECKS
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== SPOT CHECKS ===")
p = Presentation(DST)

# Slide 1: shape[5] should no longer overlap shape[6]
sl1 = p.slides[0]
sh5 = sl1.shapes[5]
sh6 = sl1.shapes[6]
t5, h5 = sh5.top/EMU_IN, sh5.height/EMU_IN
t6 = sh6.top/EMU_IN
overlap = (t5 + h5) - t6
ok1 = overlap <= 0
print(f"  S1 shape[5] bot={t5+h5:.3f}\" vs shape[6] T={t6:.3f}\" → gap={-overlap:.3f}\" [{'OK' if ok1 else 'OVERLAP!'}]")

# Slide 4: check shape order — shape[13] background should be before text shapes [6]-[12]
sl4 = p.slides[3]
spTree4 = sl4.shapes._spTree
sp_elements = list(spTree4)
# Find positions of shape elements in spTree
shape13_el = sl4.shapes[13]._element if len(sl4.shapes) > 13 else None
shape6_el  = sl4.shapes[6]._element  if len(sl4.shapes) > 6  else None
pos13 = pos6 = None
for idx, el in enumerate(sp_elements):
    if el is shape13_el:
        pos13 = idx
    if el is shape6_el:
        pos6 = idx
ok4 = (pos13 is not None and pos6 is not None and pos13 < pos6)
print(f"  S4 spTree pos: shape[13]={pos13} shape[6]={pos6} → bg before text [{'OK' if ok4 else 'FAIL'}]")

# Slide 2: footnote shape should be above 12.83"
sl2 = p.slides[1]
# Find it
fn_ok = True
for i, sh in enumerate(sl2.shapes):
    if sh.has_text_frame and 'Diagnostic' in sh.text_frame.text:
        t = sh.top / EMU_IN
        fn_ok = t < 12.83
        print(f"  S2 footnote shape[{i}] T={t:.3f}\" [{'OK' if fn_ok else 'STILL OVERLAPS'}]")
        break

all_ok = ok1 and ok4 and fn_ok and overflow_ok and xml_ok
print(f"\n{'ALL CHECKS PASSED' if all_ok else 'SOME CHECKS FAILED'}")
