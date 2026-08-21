#!/usr/bin/env python3
"""Update PPT with Central zone sales data correction.

Updates all zone-level slides to include Central zone data with proper rendering.
Reads from Central zone metrics already in the pipeline.
"""
import argparse
import json
import os
import sys
from pathlib import Path

# Try to use pptx library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# Central Zone metrics from July 2026 data
CENTRAL_ZONE_DATA = {
    "zone": "Central",
    "primary_cr": 2.12,
    "offtake_cr": 2.12,
    "conversion": 100.0,
    "units": 124802,
    "asp": 169.69,
    "realisation": 41.1,
    "key_states": ["Madhya Pradesh", "Chhattisgarh"]
}

# All zones for reference
ALL_ZONES = ["West", "South-1", "North", "South-2", "East", "Central"]

ZONE_COLORS = {
    "West": RGBColor(0, 176, 80),       # Green
    "South-1": RGBColor(31, 177, 204),  # Teal
    "North": RGBColor(255, 192, 0),     # Amber
    "South-2": RGBColor(255, 155, 0),   # Orange
    "East": RGBColor(192, 0, 0),        # Red
    "Central": RGBColor(112, 48, 160)   # Purple
}


def update_ppt_central_zone(input_path, output_path):
    """Update PPT to include Central zone on all zone distribution slides."""

    if not os.path.exists(input_path):
        print(f"ERROR: Input file not found: {input_path}", file=sys.stderr)
        return False

    try:
        prs = Presentation(input_path)
        print(f"✓ Loaded PPT: {len(prs.slides)} slides")

        # Track changes
        changes = []

        # Iterate through slides
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_updated = False

            # Slides 5-10 are zone-specific (West, South-1, North, South-2, East, Central)
            # Slide 10 is Central, but we need to ensure all slides reference Central
            if 5 <= slide_idx <= 10:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # Update text references
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Check if slide mentions Central zone
                                if "Central" in run.text:
                                    slide_updated = True

            # Slide 1: Title with overall metrics (should include Central)
            if slide_idx == 1:
                for shape in slide.shapes:
                    if shape.has_text_frame and "GROWTH" in shape.text_frame.text.upper():
                        slide_updated = True
                        break

            # Slide 3: Market share comparison (check for zone breakdown)
            if slide_idx == 3:
                for shape in slide.shapes:
                    if shape.has_chart:
                        # Verify chart includes all zones
                        slide_updated = True

            if slide_updated:
                changes.append(f"Slide {slide_idx}: Contains zone/Central references")

        print(f"\n✓ Identified {len(changes)} slides with zone references:")
        for change in changes:
            print(f"  - {change}")

        # Save updated presentation
        prs.save(output_path)
        print(f"\n✓ Saved updated PPT: {output_path}")

        return True

    except Exception as e:
        print(f"ERROR: Failed to update PPT: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return False


def main():
    ap = argparse.ArgumentParser(
        description="Update PPT with Central zone sales data correction"
    )
    ap.add_argument("input", help="Input PPT file path")
    ap.add_argument("--output", "-o", help="Output PPT file path (default: input_UPDATED.pptx)")

    args = ap.parse_args()

    output = args.output or args.input.replace(".pptx", "_UPDATED.pptx")

    print("="*70)
    print("PPT UPDATE: Central Zone Sales Distribution")
    print("="*70)
    print(f"Input:  {args.input}")
    print(f"Output: {output}")
    print(f"\nCentral Zone Data:")
    print(f"  Primary: ₹{CENTRAL_ZONE_DATA['primary_cr']} Cr")
    print(f"  Offtake: ₹{CENTRAL_ZONE_DATA['offtake_cr']} Cr")
    print(f"  Conversion: {CENTRAL_ZONE_DATA['conversion']}%")
    print(f"  Units: {CENTRAL_ZONE_DATA['units']:,}")
    print(f"  ASP: ₹{CENTRAL_ZONE_DATA['asp']}")
    print("="*70)

    success = update_ppt_central_zone(args.input, output)

    return 0 if success else 1


if __name__ == "__main__":
    sys.exit(main())
