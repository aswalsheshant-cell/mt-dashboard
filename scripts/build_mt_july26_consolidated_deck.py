#!/usr/bin/env python3
"""
Build consolidated MT July'26 Master Deck (15 slides)
Merges File 1 (Q1 narrative) + File 2 (operational detail) + root-cause analysis
Output: MT_Jul26_Consolidated_v1.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# =====================================================================
# Colour scheme (Mamaearth brand)
# =====================================================================
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
TEAL_GROWTH = RGBColor(0x2A, 0x9D, 0x8F)
RED_DECLINE = RGBColor(0xE6, 0x39, 0x46)
ORANGE_CAUTION = RGBColor(0xF4, 0xA2, 0x61)
GREY_NEUTRAL = RGBColor(0x75, 0x75, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)

# =====================================================================
# Helper functions
# =====================================================================

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = TEAL_GROWTH

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Mamaearth Modern Trade Analytics | {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(10)
    p.font.color.rgb = GREY_NEUTRAL

    return slide

def add_content_slide(prs, headline, evidence, implication, action, owner_timeline):
    """Add a content slide with EIAO structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Headline (with background bar)
    headline_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.8))
    headline_shape.fill.solid()
    headline_shape.fill.fore_color.rgb = NAVY
    headline_shape.line.color.rgb = NAVY

    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Evidence box (with left border)
    ev_shape = slide.shapes.add_shape(1, Inches(0.3), Inches(1.3), Inches(0.08), Inches(1))
    ev_shape.fill.solid()
    ev_shape.fill.fore_color.rgb = RED_DECLINE
    ev_shape.line.color.rgb = RED_DECLINE

    ev_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    tf = ev_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"EVIDENCE: {evidence}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RED_DECLINE

    # Implication box
    imp_shape = slide.shapes.add_shape(1, Inches(0.3), Inches(2.6), Inches(0.08), Inches(0.8))
    imp_shape.fill.solid()
    imp_shape.fill.fore_color.rgb = ORANGE_CAUTION
    imp_shape.line.color.rgb = ORANGE_CAUTION

    imp_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.1))
    tf = imp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"IMPLICATION: {implication}"
    p.font.size = Pt(13)
    p.font.color.rgb = NAVY

    # Action box
    act_shape = slide.shapes.add_shape(1, Inches(0.3), Inches(3.8), Inches(0.08), Inches(1.2))
    act_shape.fill.solid()
    act_shape.fill.fore_color.rgb = TEAL_GROWTH
    act_shape.line.color.rgb = TEAL_GROWTH

    act_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.4))
    tf = act_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"ACTION: {action}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEAL_GROWTH

    # Owner & Timeline
    owner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.5))
    tf = owner_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Owner: {owner_timeline}"
    p.font.size = Pt(11)
    p.font.color.rgb = GREY_NEUTRAL
    p.font.italic = True

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Modern Trade Analytics | July 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = GREY_NEUTRAL

    return slide

def add_table_slide(prs, headline, table_data, column_widths):
    """Add a slide with a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Headline
    headline_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.3), Inches(10), Inches(0.6))
    headline_shape.fill.solid()
    headline_shape.fill.fore_color.rgb = NAVY
    headline_shape.line.color.rgb = NAVY

    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.5))
    tf = headline_box.text_frame
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(9)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for i, col_width in enumerate(column_widths):
        table_shape.columns[i].width = Inches(col_width)

    # Populate table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = str(cell_data)

            # Header row formatting
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.font.bold = True
                para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER
            else:
                # Alternate row colours
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREY
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.font.color.rgb = NAVY

    return slide

# =====================================================================
# Main: Build 15-slide deck
# =====================================================================

def build_consolidated_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Modern Trade July 2026", "Monthly Performance Review & Action Plan")

    # Slide 2: Executive Scorecard
    add_content_slide(
        prs,
        "July offtake ₹36.1 Cr — Q1 delivery +64% YoY",
        "July MT offtake ₹36.1 Cr; Q1 FY27 offtake ₹114.39 Cr (+64% YoY vs Q1 FY26).",
        "Strong Q1 momentum but July shows ₹10.92 Cr billed-not-sold gap — inventory risk.",
        "Reconcile July actual-vs-forecast by 8-Aug; flag any DC fill-rate drops.",
        "Sales Lead + Data Team | 8-Aug"
    )

    # Slide 3: Topline Bridge (Waterfall logic)
    add_content_slide(
        prs,
        "TDC +365% and Face Wash +33% offset Haircare decline",
        "TDC ₹29.5 Cr (+365% YoY); Face Wash ₹241.4 Cr MT category (+33% YoY); Haircare pressure: Shampoo ₹22 Cr, Sun Care ₹21.75 Cr.",
        "Portfolio mix shift: premium-tier brands (TDC, Aqualogica) driving growth; legacy core (Haircare) losing share to budget competitors.",
        "Increase FSDU for Face Wash by mid-Aug; launch Haircare defensive entry SKU (sub-₹300) by Oct.",
        "Category Head | 15-Aug (FSDU), 30-Oct (SKU launch)"
    )

    # Slide 4: Chain Performance Matrix
    add_content_slide(
        prs,
        "Reliance -15% MoM — single largest threat to MT growth",
        "Reliance fell ₹1.42 Cr MoM; conversion 51.4% vs DMart 76.5% vs Apollo 99.7%.",
        "Reliance is 40%+ of MT volume; -15% = ₹8–10 Cr recovery opportunity; low conversion signals planogram/shelf compliance issue.",
        "Audit Reliance hero-EAN (Shampoo/Sun Care) door-level OSA; reallocate trade spend to Haircare trial schemes.",
        "NKAM Reliance + Category Lead | 12-Aug checkpoint"
    )

    # Slide 5: Reliance Deep-Dive
    add_content_slide(
        prs,
        "Reliance Shampoo/Sun Care mix erosion — three drivers",
        "Why -15%: (a) Haircare weak across MT (consumer down-trading to budget HUL sub-₹350 shampoo); (b) Premium positioning not resonating at Reliance; (c) Conversion 51.4% = poor shelf visibility or planogram gap.",
        "If uncorrected, Reliance declines ₹2–3 Cr more by Aug-end; risk to Q2 guidance.",
        "Execute: (1) EAN-store pair placement audit by 10-Aug; (2) Deploy trial packs (₹100–150) from Aug 1; (3) Weekly conversion tracking.",
        "NKAM Reliance | Daily tracking; report by 12-Aug"
    )

    # Slide 6: Zone Performance Analysis
    add_content_slide(
        prs,
        "North 58.5%, East 45.3% — zones below benchmark need urgent recovery",
        "Zone conversions: North 58.5% (urgent), East 45.3% (worst), West 82.3% (healthy), South stable. Benchmark = 75%.",
        "East 45% signals severe supply/execution gap; North 58% = 16 pp gap to target. Risk: ₹2–3 Cr lost by month-end.",
        "North: increase trial schemes + DC fill-rate audit; East: early monsoon stock build + regional trade spend spike by Aug 1.",
        "North Regional Lead + East Regional Lead | 31-Aug recovery"
    )

    # Slide 7: Face Wash Uplift
    add_content_slide(
        prs,
        "Face Wash +33% YoY — premium tier growth + natural positioning momentum",
        "Face Wash ₹241.4 Cr MT category (+33% YoY); premium segment growing +15% YoY market-wide.",
        "Haircare weakness driving substitution to face care; Mamaearth/TDC natural positioning capturing trial.",
        "Increase FSDU placements by 30% across DMart/Reliance; expand eB2B (FSN+Nykaa at 99.4% conversion).",
        "Category Head + eB2B Lead | 15-Aug FSDU, 31-Aug eB2B 2x"
    )

    # Slide 8: Haircare Decline
    add_content_slide(
        prs,
        "Haircare losing to value tier — launch sub-₹300 SKU by Oct",
        "Shampoo ₹22 Cr, Sun Care ₹21.75 Cr facing -12% to -18% pressure; HUL sub-₹350 offensive accelerating consumer down-trading.",
        "Without defensive entry tier, Haircare share at risk of -500 bps by Q3; margin compression already at 120 bps.",
        "Launch sub-₹300 shampoo variant by Oct (9-week runway); support with aggressive trial (₹50 trial packs). Target: +2 pp share in budget tier.",
        "Brand Marketing + Category Lead | 30-Oct launch; target ₹3–4 Cr NSV by Q3"
    )

    # Slide 9: Brand Performance
    add_content_slide(
        prs,
        "TDC momentum strong (+365%), BBlunt underperforming (35.2% conversion)",
        "TDC ₹29.5 Cr (+365% YoY); Mamaearth ₹24.55 Cr at 73.4% conversion; Aqualogica 117% (timing flag); BBlunt 35.2% (underperform).",
        "TDC = new profit engine; BBlunt = drag on portfolio. Aqualogica 117% may signal overloading or timing mismatch.",
        "Double TDC allocation by 31-Aug; audit BBlunt SKU velocity — consolidate slow-moving SKUs; validate Aqualogica timing.",
        "Brand Heads (TDC, BBlunt, Aqualogica) | 31-Aug action update"
    )

    # Slide 10: Market Bifurcation
    add_content_slide(
        prs,
        "Premium +15% YoY, Budget -7% YoY — portfolio gap identified",
        "Market trend: Premium tier +15% YoY (opportunity); Budget tier -7% YoY (risk). Mamaearth portfolio skews premium; sub-₹300 entry tier gaps.",
        "Gap = ₹15–20 Cr opportunity in sub-₹300 if we own first-mover advantage; delay = share loss to HUL/P&G.",
        "Prioritise budget Face Wash + Shampoo SKU development (parallel to premium line extension). Timeline: Face Wash Sep, Shampoo Oct.",
        "Category Head + Brand Marketing | Sep-Oct SKU launch milestones"
    )

    # Slide 11: Digital & eB2B Channel
    add_content_slide(
        prs,
        "eB2B channels (FSN+Nykaa) at 99.4% conversion — 2x potential identified",
        "FSN+Nykaa ₹2.07 Cr at 99.4% conversion (benchmark best); current allocation = 5–7% of MT budget; 2x opportunity = +₹2 Cr potential.",
        "eB2B removes trade friction; captures online-ready consumers; efficiency = 99.4% conversion vs retail 51–76%.",
        "Double eB2B allocation by 31-Aug; increase SKU breadth (all Hero SKUs + new variants); target ₹3.5 Cr eB2B by month-end.",
        "eB2B Lead + Category Head | 31-Aug ₹3.5 Cr target"
    )

    # Slide 12: Operational Accountability Matrix
    table_data = [
        ["Action", "Owner", "Deadline", "KPI Target"],
        ["Reliance Shampoo/Sun Care recovery", "NKAM Reliance + Category Lead", "12-Aug", "+₹2 Cr offtake"],
        ["North Zone conversion ramp", "North Regional Lead", "30-Aug", "58.5% → 65%"],
        ["eB2B allocation 2x", "eB2B Lead", "31-Aug", "₹3.5 Cr"],
        ["Premium positioning campaign (TDC + bundle)", "Trade Marketing", "31-Aug", "+5 pp premium penetration"],
        ["East Zone recovery plan", "East Regional Lead", "31-Aug", "45.3% → 50%+ conversion"],
    ]
    add_table_slide(prs, "August Action Plan — 5 Key Levers", table_data, [2.5, 2.5, 1.5, 1.5])

    # Slide 13: August Recovery Trajectory
    add_content_slide(
        prs,
        "August recovery trajectory: Q1 baseline → Aug actions → Q3 target ₹400+ Cr",
        "July baseline ₹36.1 Cr; with Reliance recovery (+₹2 Cr) + North/East ramp (+₹2–3 Cr) + eB2B 2x (+₹2 Cr) = Aug target ₹42–45 Cr.",
        "Momentum carries to Q2 Q3 with festive loading (Sep-Oct) + budget SKU launch = ₹400+ Cr Q2 run-rate achievable.",
        "Weekly steering: Monday sales sync (offtake vs plan), Friday action audit (owner accountability). Escalate misses >3% by Tue-EOD.",
        "Sales VP + Category Head | Weekly cadence"
    )

    # Slide 14: Risk & Mitigation
    add_content_slide(
        prs,
        "Four headwinds: Supply-chain fill-rate, footfall, competitor push, budget cannibalization",
        "Risks: (a) DC fill-rate drops = OOS in North/East; (b) Footfall compression in tier-2 chains (Nature's Basket, Spencers); (c) HUL value-tier acceleration; (d) Budget SKU cannibalizing premium.",
        "Risk = ₹5–7 Cr Q3 if mitigations miss. Ongoing hedge: supply-chain daily audit, competitor pricing intelligence, portfolio cannibalization modeling.",
        "Assign: DC owner (fill-rate audit by 5-Aug), Trade (footfall survey by 8-Aug), Category (cannibalization model by 12-Aug).",
        "Supply Chain + Trade + Category Heads | Audit reports by stated dates"
    )

    # Slide 15: Next Steps & Q2 Outlook
    add_content_slide(
        prs,
        "Sep-Oct festive loading + budget SKU launch + zone recovery = ₹400+ Cr Q2 run-rate",
        "Q2 outlook: Sep-Oct festive schemes + budget Face Wash (Sep) + budget Shampoo (Oct) + North/East zone recovery = ₹110+ Cr Q2 offtake.",
        "Gates: (1) Reliance ₹2 Cr recovery by 12-Aug; (2) eB2B ₹3.5 Cr by 31-Aug; (3) Budget SKU readiness by 15-Sep.",
        "Steering cadence: Weekly sales sync (Mon), Friday action audit, monthly leadership review. Next review: Aug 30 (Q3 readiness).",
        "Sales VP | Next leadership review 30-Aug"
    )

    # Save
    output_path = "/home/user/mt-dashboard/MT_Jul26_Consolidated_Master_v1.pptx"
    prs.save(output_path)
    print(f"\n✅ CONSOLIDATED DECK CREATED: {output_path}")
    print(f"   Slides: 15 | Merged: File 1 narrative + File 2 operations + Root-cause analysis")
    return output_path

if __name__ == "__main__":
    path = build_consolidated_deck()
    print(f"\n✅ Ready for distribution. File size: ", end="")
    import os
    print(f"{os.path.getsize(path) / (1024*1024):.1f} MB")
