#!/usr/bin/env python3
"""
ppt_skill.py — Automated PPT validation and repair skill for Honasa MT decks.

Usage:
    python scripts/ppt_skill.py --input deck.pptx               # validate only
    python scripts/ppt_skill.py --input deck.pptx --fix          # validate + fix
    python scripts/ppt_skill.py --input deck.pptx --fix --out fixed.pptx

What it checks / fixes:
  1. Text overflow (shape text taller than the shape bounds) → shrinks font or
     truncates to fit
  2. External chart data links → replaces with embedded numCache so charts
     render without the linked Excel file
  3. Blank / empty placeholder shapes → reports them (no auto-fix; need data)
  4. Series with all-zero or all-null values → warns (chart will look empty)
  5. Duplicate series names within a chart → warns
  6. Page-number text-box sanity (e.g. "11 / 9" on slide 7) → auto-corrects
     to the real page number
"""
from __future__ import annotations
import argparse, sys, re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    from lxml import etree
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx lxml")
    sys.exit(1)


# ─── helpers ────────────────────────────────────────────────────────────────

def _shape_text(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def _series_values(series) -> list:
    try:
        return list(series.values)
    except Exception:
        return []


# ─── check 1: text overflow ─────────────────────────────────────────────────

def check_text_overflow(slide_idx: int, shape_idx: int, shape, fix: bool) -> list[str]:
    issues = []
    if not shape.has_text_frame:
        return issues
    tf = shape.text_frame
    # heuristic: if auto_size is None and text has many chars for the box height
    try:
        total_chars = len(tf.text)
        box_h_pt = shape.height / 12700  # EMU → points
        # rough estimate: ~72pt per line, ~80 chars per line at 10pt
        if box_h_pt > 0 and total_chars > 0:
            lines_available = max(1, box_h_pt / 12)
            chars_per_line = max(40, shape.width / 12700 / 6)
            lines_needed = total_chars / chars_per_line
            if lines_needed > lines_available * 1.4:
                issues.append(
                    f"  [OVERFLOW] Slide {slide_idx+1} Shape {shape_idx} '{shape.name}': "
                    f"~{int(lines_needed)} lines needed, ~{int(lines_available)} available"
                )
                if fix:
                    # enable auto-fit (shrink text to fit)
                    from pptx.oxml.ns import qn as _qn
                    txBody = tf._txBody
                    bodyPr = txBody.find(_qn('a:bodyPr'))
                    if bodyPr is not None:
                        # remove spAutoFit / noAutofit, add normAutofit
                        for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
                            for el in bodyPr.findall(_qn(tag)):
                                bodyPr.remove(el)
                        norm = etree.SubElement(bodyPr, _qn('a:normAutofit'))
                        issues.append(f"    → FIXED: normAutofit applied to shape {shape_idx}")
    except Exception:
        pass
    return issues


# ─── check 2: external data references in charts ───────────────────────────

def check_chart_external_data(slide_idx: int, shape_idx: int, shape, fix: bool) -> list[str]:
    issues = []
    try:
        ch = shape.chart
    except Exception:
        return issues

    el = ch._element
    ext_refs = el.findall('.//' + qn('c:externalData'))
    if not ext_refs:
        return issues

    issues.append(
        f"  [EXTERNAL_DATA] Slide {slide_idx+1} Shape {shape_idx} '{shape.name}': "
        f"chart linked to external Excel ({len(ext_refs)} ref(s)) — charts may appear "
        f"blank if the Excel file is missing"
    )
    if fix:
        for ref in ext_refs:
            ref.getparent().remove(ref)
        # Also remove any numRef formula so PowerPoint uses only the cache
        for numRef in el.findall('.//' + qn('c:numRef')):
            f_el = numRef.find(qn('c:f'))
            if f_el is not None:
                numRef.remove(f_el)
        issues.append(f"    → FIXED: externalData refs removed; chart now uses embedded cache only")

    return issues


# ─── check 3: blank shapes ──────────────────────────────────────────────────

def check_blank_shapes(slide_idx: int, shape_idx: int, shape) -> list[str]:
    issues = []
    if not shape.has_text_frame:
        return issues
    txt = _shape_text(shape)
    if txt == '' and shape.shape_type in (1, 17):  # AUTO_SHAPE or TEXT_BOX
        # small shapes (< 0.5 inch) are probably spacers — skip
        try:
            if shape.width / 914400 > 0.5 and shape.height / 914400 > 0.1:
                issues.append(
                    f"  [BLANK_SHAPE] Slide {slide_idx+1} Shape {shape_idx} '{shape.name}' "
                    f"(text box with no content)"
                )
        except Exception:
            pass
    return issues


# ─── check 4: series with all-zero / empty values ───────────────────────────

def check_series_values(slide_idx: int, shape_idx: int, shape) -> list[str]:
    issues = []
    try:
        ch = shape.chart
    except Exception:
        return issues

    seen_names: set[str] = set()
    for s in ch.series:
        name = s.name or '(unnamed)'
        vals = _series_values(s)
        non_null = [v for v in vals if v is not None]
        if not non_null or all(v == 0 for v in non_null):
            issues.append(
                f"  [EMPTY_SERIES] Slide {slide_idx+1} Shape {shape_idx} '{shape.name}': "
                f"series '{name}' has all-zero/null values — chart will appear empty for this series"
            )
        if name in seen_names:
            issues.append(
                f"  [DUP_SERIES] Slide {slide_idx+1} Shape {shape_idx}: "
                f"duplicate series name '{name}'"
            )
        seen_names.add(name)
    return issues


# ─── check 5: page number sanity ────────────────────────────────────────────

_PAGE_NUM_RE = re.compile(r'^(\d+)\s*/\s*(\d+)$')

def check_page_numbers(slide_idx: int, total_slides: int, shape_idx: int, shape, fix: bool) -> list[str]:
    issues = []
    if not shape.has_text_frame:
        return issues
    txt = _shape_text(shape)
    m = _PAGE_NUM_RE.match(txt)
    if not m:
        return issues
    shown_page = int(m.group(1))
    shown_total = int(m.group(2))
    expected_page = slide_idx + 1
    if shown_page != expected_page or shown_total != total_slides:
        issues.append(
            f"  [PAGE_NUM] Slide {slide_idx+1} Shape {shape_idx}: "
            f"shows '{txt}' but should be '{expected_page} / {total_slides}'"
        )
        if fix:
            from pptx.oxml.ns import qn as _qn
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if _PAGE_NUM_RE.match(run.text.strip()):
                        run.text = f"{expected_page} / {total_slides}"
                        break
            issues.append(f"    → FIXED: page number corrected to '{expected_page} / {total_slides}'")
    return issues


# ─── main ───────────────────────────────────────────────────────────────────

def run(input_path: Path, fix: bool, out_path: Path | None) -> None:
    print(f"Loading: {input_path}")
    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    print(f"Slides: {total_slides}\n")

    all_issues: list[str] = []

    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            all_issues.extend(check_text_overflow(si, shi, shape, fix))
            all_issues.extend(check_chart_external_data(si, shi, shape, fix))
            all_issues.extend(check_blank_shapes(si, shi, shape))
            all_issues.extend(check_series_values(si, shi, shape))
            all_issues.extend(check_page_numbers(si, total_slides, shi, shape, fix))

    if not all_issues:
        print("✓ No issues found.")
    else:
        print(f"Found {len([x for x in all_issues if not x.startswith('    →')])} issue(s):\n")
        for line in all_issues:
            print(line)

    if fix:
        dest = out_path or input_path.with_stem(input_path.stem + '_fixed')
        prs.save(dest)
        print(f"\nSaved fixed deck → {dest}")

        # validate by reloading
        p2 = Presentation(dest)
        print(f"Validation: reloaded OK ({len(p2.slides)} slides)")


def main():
    ap = argparse.ArgumentParser(description="Validate and optionally repair Honasa MT PPT decks.")
    ap.add_argument("--input", "-i", required=True, help="Input .pptx file")
    ap.add_argument("--fix", "-f", action="store_true", help="Apply automatic fixes")
    ap.add_argument("--out", "-o", help="Output path (default: <input>_fixed.pptx)")
    args = ap.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        print(f"ERROR: {inp} does not exist")
        sys.exit(1)

    out = Path(args.out) if args.out else None
    run(inp, args.fix, out)


if __name__ == "__main__":
    main()
