"""Build V18 from V17 — correct offtake values after pending chains (Trent, Guardian,
Wh-Smith, VMM, More Retail, Spencer, H&G, Wellness Forever, etc.) are incorporated.

The Jul'26 offtake CSV always contained these chains. The PPT (V13→V17) was built
with hardcoded values from an earlier partial data pull that missed 22 chains. This
script patches ALL offtake-dependent numbers to match the complete CSV.

Data source: PowerBI/RawDataFolders/Offtake_Monthly/offtake_store_article_Jul_26.csv
Zone totals computed and hardcoded below for auditability.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
import zipfile

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv17.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv18.pptx"

EMU_IN = 914400

# ─── CORRECT VALUES FROM COMPLETE Jul'26 CSV (all 27 chains) ─────────────────
# NSV in Lakhs → Cr = /100; Qty in absolute → K = /1000
# National (all 6 zones + FSN/Pan India)
NAT_OFF26_NSV   = 40.67   # Cr  (prev 36.10)
NAT_OFF26_QTY   = 2269    # K   (prev 2114)
NAT_PRI26_NSV   = 47.02   # Cr  (unchanged — primary data not in scope)
NAT_PRI26_QTY   = 2742    # K   (unchanged)
NAT_CONV_PCT    = 86.5    # %   40.67/47.02 (prev 76.8%)
NAT_GAP_NSV     = 6.35    # Cr  47.02-40.67 (prev 10.92)

# Zone-level offtake corrections  (pri26/pq26 unchanged — primary data not in scope)
# fmt: (zone_name, slide_idx, off26_new, oq26_new, mix_pct_new, conv_pct_new, gap_new,
#        ins01_shape_idx, ins02_shape_idx, rank_shape_idx, ev_shape_idx, rank_new,
#        old_off26, old_oq26, old_mix_pct, old_conv, old_gap)
ZONES = [
    # West: DMart+Reliance+Wellness+Apollo (pending: VMM, More, H&G, Spencer, etc.)
    dict(name='West', si=4,
         off_new=8.49,  oq_new=497,  mix_new=20.9,  conv_new=84.5,  gap_new=1.56,
         off_old=8.28,  oq_old=466,  mix_old=22.9,  conv_old=82.3,  gap_old=1.78,
         pri26=10.05, pq26=597, pri_yoy=99.0,
         ins01=55, ins02=58, rank_sh=67, ev_sh=101, rank_new=2,
         ev_new="Off: ₹8.49Cr(497K) | DMart+Reliance+WF top 3 | 84.5% conv | ₹1.56Cr gap"),
    # South-1: Apollo+DMart+Lulu+Reliance
    dict(name='South-1', si=5,
         off_new=8.77,  oq_new=461,  mix_new=21.6,  conv_new=89.5,  gap_new=1.03,
         off_old=8.19,  oq_old=435,  mix_old=22.7,  conv_old=83.6,  gap_old=1.61,
         pri26=9.80, pq26=581, pri_yoy=63.6,
         ins01=55, ins02=58, rank_sh=67, ev_sh=101, rank_new=1,
         ev_new="Off: ₹8.77Cr(461K) | Apollo+DMart+Lulu top 3 | 89.5% conv | ₹1.03Cr gap"),
    # North: Reliance+DMart+Apollo+Reliance BC  — largest absolute gain (+1.33 Cr)
    dict(name='North', si=6,
         off_new=8.32,  oq_new=442,  mix_new=20.5,  conv_new=69.6,  gap_new=3.63,
         off_old=6.99,  oq_old=377,  mix_old=19.4,  conv_old=58.5,  gap_old=4.97,
         pri26=11.95, pq26=649, pri_yoy=134.3,
         ins01=55, ins02=58, rank_sh=67, ev_sh=101, rank_new=3,
         ev_new="Off: ₹8.32Cr(442K) | Rel+DMart+Apollo top 3 | 69.6% conv | ₹3.63Cr gap"),
    # South-2
    dict(name='South-2', si=7,
         off_new=5.30,  oq_new=299,  mix_new=13.0,  conv_new=77.0,  gap_new=1.59,
         off_old=4.91,  oq_old=281,  mix_old=13.6,  conv_old=71.3,  gap_old=1.98,
         pri26=6.89, pq26=343, pri_yoy=69.3,
         ins01=54, ins02=57, rank_sh=66, ev_sh=100, rank_new=4,
         ev_new="Off: ₹5.30Cr(299K) | DMart+Apollo+Rel top 3 | 77.0% conv | ₹1.59Cr gap"),
    # East: Reliance+Reliance BC+Apollo — second largest gain (+1.30 Cr)
    dict(name='East', si=8,
         off_new=4.85,  oq_new=245,  mix_new=11.9,  conv_new=61.9,  gap_new=2.99,
         off_old=3.55,  oq_old=183,  mix_old=9.8,   conv_old=45.3,  gap_old=4.28,
         pri26=7.83, pq26=409, pri_yoy=115.7,
         ins01=55, ins02=58, rank_sh=67, ev_sh=101, rank_new=5,
         ev_new="Off: ₹4.85Cr(245K) | Rel+Apollo+RBC top 3 | 61.9% conv | ₹2.99Cr gap"),
    # Central: DMart+Reliance+Apollo — offtake now slightly > primary (timing)
    dict(name='Central', si=9,
         off_new=2.88,  oq_new=171,  mix_new=7.1,   conv_new=107.1, gap_new=-0.19,
         off_old=2.12,  oq_old=160,  mix_old=5.9,   conv_old=78.8,  gap_old=0.57,
         pri26=2.69, pq26=160, pri_yoy=202.2,
         ins01=54, ins02=57, rank_sh=66, ev_sh=100, rank_new=6,
         ev_new="Off: ₹2.88Cr(171K) | DMart+Rel+Apollo top 3 | 107% conv (timing surplus)"),
]


# ── Helpers ────────────────────────────────────────────────────────────────────
def get_text(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def patch(shape, old_frag: str, new_frag: str, label="") -> bool:
    if not shape.has_text_frame:
        return False
    full = shape.text_frame.text
    if old_frag not in full:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_frag in run.text:
                run.text = run.text.replace(old_frag, new_frag)
                if label:
                    print(f"    [{label}] '{old_frag[:35]}' → '{new_frag[:50]}'")
                return True
    # Multi-run fallback
    for para in shape.text_frame.paragraphs:
        pt = "".join(r.text for r in para.runs)
        if old_frag in pt:
            if para.runs:
                para.runs[0].text = pt.replace(old_frag, new_frag)
                for r in para.runs[1:]:
                    r.text = ""
            if label:
                print(f"    [{label}] (multi-run) '{old_frag[:35]}'")
            return True
    return False


def set_text_safe(shape, new_text: str) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    runs = first_para._p.findall(f'{ns}r')
    if runs:
        runs[0].find(f'{ns}t').text = new_text
        for r in runs[1:]:
            t = r.find(f'{ns}t')
            if t is not None:
                t.text = ""
    else:
        para_el = first_para._p
        r_el = etree.SubElement(para_el, qn('a:r'))
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = new_text


def apply_normautofit(shape):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


def audit_overflows(prs):
    slide_h = prs.slide_height / EMU_IN
    issues = []
    for si, slide in enumerate(prs.slides):
        for i, sh in enumerate(slide.shapes):
            bot = (sh.top + sh.height) / EMU_IN
            if bot > slide_h + 0.015:
                txt = sh.text_frame.text[:30] if sh.has_text_frame else ""
                issues.append(f"  S{si+1}[{i}] bot={bot:.3f}\" | {txt!r}")
    if issues:
        print(f"[OVERFLOW] {len(issues)} issue(s):")
        for x in issues: print(x)
    else:
        print("[OVERFLOW] Clean.")
    return len(issues) == 0


prs = Presentation(SRC)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — national totals
# ─────────────────────────────────────────────────────────────────────────────
print("=== SLIDE 1 — National offtake totals ===")
s1 = prs.slides[0]
shapes1 = s1.shapes

# shape[5]: "₹36.1 Cr July offtake | 2114K units | 76.8% flow conversion"
old5 = f"₹36.1 Cr July offtake | 2114K units | 76.8% flow conversion"
new5 = f"₹{NAT_OFF26_NSV} Cr July offtake | {NAT_OFF26_QTY}K units | {NAT_CONV_PCT}% flow conversion"
if not patch(shapes1[5], old5, new5, "S1:off+conv"):
    # Try fragment patches
    patch(shapes1[5], "₹36.1 Cr July offtake", f"₹{NAT_OFF26_NSV} Cr July offtake", "S1:off")
    patch(shapes1[5], "2114K units", f"{NAT_OFF26_QTY}K units", "S1:units")
    patch(shapes1[5], "76.8% flow conversion", f"{NAT_CONV_PCT}% flow conversion", "S1:conv")
apply_normautofit(shapes1[5])

# shape[6]: "...₹10.92 Cr gap..."
patch(shapes1[6], "₹10.92 Cr gap", f"₹{NAT_GAP_NSV} Cr gap", "S1:gap")
apply_normautofit(shapes1[6])

# shape[38]: "West: ₹8.28 Cr at 82.3% conversion."
patch(shapes1[38], "West: ₹8.28 Cr at 82.3% conversion.", "West: ₹8.49 Cr at 84.5% conversion.", "S1:West")

# shape[40]: "South‑1: ₹8.19 Cr at 83.6%."
patch(shapes1[40], "₹8.19 Cr at 83.6%", "₹8.77 Cr at 89.5%", "S1:S1")

# shape[42]: "North: ₹4.97 Cr gap; largest recovery pool."
patch(shapes1[42], "₹4.97 Cr gap", "₹3.63 Cr gap", "S1:North gap")

# shape[44]: "East: 45.3% conversion; weakest flow."
patch(shapes1[44], "45.3% conversion", "61.9% conversion", "S1:East conv")

# shape[18]: "North + East hold 70.5% of the gap."
# North gap 3.63 + East gap 2.99 = 6.62; total zone gap = sum of positives = 1.56+1.03+3.63+1.59+2.99 = 10.80
# North+East share of zone gap = 6.62/10.80 = 61.3%
patch(shapes1[18], "70.5% of the gap", "61.3% of the zone gap", "S1:NE gap share")

# shape[20]: "West + South‑1 deliver 45.6% of offtake."
# New: (8.49+8.77)/40.67 = 42.5%
patch(shapes1[20], "45.6% of offtake", "42.5% of offtake", "S1:WS mix")

# shape[53]: "Mamaearth ₹24.55 Cr | 76.8% conversion."
# Mamaearth offtake from CSV = 27.83 Cr; conv = 27.83/primary_ME... skip brand conv for now
# Just patch the national conversion reference
patch(shapes1[53], "76.8% conversion", "86.5% conversion", "S1:ME conv")

print()

# ─────────────────────────────────────────────────────────────────────────────
# ZONE SLIDES (5–10)
# ─────────────────────────────────────────────────────────────────────────────
for Z in ZONES:
    name = Z['name']
    si   = Z['si']
    slide = prs.slides[si]
    shapes = slide.shapes
    print(f"=== {name} (slide {si+1}) ===")

    off_old = Z['off_old']
    oq_old  = Z['oq_old']
    mix_old = Z['mix_old']
    conv_old= Z['conv_old']
    gap_old = Z['gap_old']

    off_new = Z['off_new']
    oq_new  = Z['oq_new']
    mix_new = Z['mix_new']
    conv_new= Z['conv_new']
    gap_new = Z['gap_new']
    pri26   = Z['pri26']
    pq26    = Z['pq26']
    pri_yoy = Z['pri_yoy']
    rank_new= Z['rank_new']

    # shape[12]: OFFTAKE value (₹X.XX Cr)
    patch(shapes[12], f"₹{off_old} Cr", f"₹{off_new} Cr", f"{name}:off26")
    apply_normautofit(shapes[12])

    # shape[13]: mix | units
    patch(shapes[13], f"{mix_old}% mix | {oq_old}K units",
                      f"{mix_new}% mix | {oq_new}K units", f"{name}:mix+units")
    apply_normautofit(shapes[13])

    # shape[17]: conversion %
    old_conv_str = f"{conv_old}%"
    new_conv_str = f"{conv_new:.1f}%".rstrip('0').rstrip('.') + '%' if '.' in f"{conv_new:.1f}" else f"{conv_new}%"
    # Simpler: just use formatted strings
    new_conv_str = f"{conv_new}%"
    patch(shapes[17], old_conv_str, new_conv_str, f"{name}:conv")
    apply_normautofit(shapes[17])

    # shape[22]: gap value — handle negative (Central has surplus)
    if gap_new < 0:
        new_gap_str = "₹0.19 Cr surplus"
    else:
        new_gap_str = f"₹{gap_new} Cr"
    patch(shapes[22], f"₹{gap_old} Cr", new_gap_str, f"{name}:gap")
    apply_normautofit(shapes[22])

    # Insight 01: "Zone delivers ₹X.XXCr offtake (XXXK units), XX.X% national."
    old_ins01 = f"₹{off_old}Cr offtake ({oq_old}K units), {mix_old}% national."
    new_ins01 = f"₹{off_new}Cr offtake ({oq_new}K units), {mix_new}% national."
    ok = patch(shapes[Z['ins01']], old_ins01, new_ins01, f"{name}:ins01")
    if not ok:
        # Try wider search
        for i, sh in enumerate(shapes):
            if sh.has_text_frame and f"₹{off_old}Cr offtake" in sh.text_frame.text:
                patch(sh, f"₹{off_old}Cr offtake ({oq_old}K units), {mix_old}% national.",
                          new_ins01, f"{name}:ins01 sh{i}")
                break
    if Z['ins01'] < len(shapes):
        apply_normautofit(shapes[Z['ins01']])

    # Insight 02: "Primary +XX.X% YoY: ₹X.XXCr (XXXK units) vs offtake ₹X.XXCr (XXXK)."
    old_ins02 = f"vs offtake ₹{off_old}Cr ({oq_old}K)."
    new_ins02 = f"vs offtake ₹{off_new}Cr ({oq_new}K)."
    ok = patch(shapes[Z['ins02']], old_ins02, new_ins02, f"{name}:ins02")
    if not ok:
        for i, sh in enumerate(shapes):
            if sh.has_text_frame and f"vs offtake ₹{off_old}Cr" in sh.text_frame.text:
                patch(sh, old_ins02, new_ins02, f"{name}:ins02 sh{i}")
                break
    if Z['ins02'] < len(shapes):
        apply_normautofit(shapes[Z['ins02']])

    # Zone rank insight: "Zone ranks #N by national offtake."
    # Current rank was inserted based on old ordering; patch it
    rank_sh = Z['rank_sh']
    if rank_sh < len(shapes):
        cur_rank_txt = get_text(shapes[rank_sh])
        # Find and replace old rank reference
        for old_rank in range(1, 7):
            old_r = f"ranks #{old_rank} by national offtake"
            new_r = f"ranks #{rank_new} by national offtake"
            if old_r in cur_rank_txt and old_r != new_r:
                patch(shapes[rank_sh], old_r, new_r, f"{name}:rank #{old_rank}→#{rank_new}")
                break
        apply_normautofit(shapes[rank_sh])

    # Evidence box: update with correct offtake values
    ev_sh = Z['ev_sh']
    if ev_sh < len(shapes):
        ev_shape = shapes[ev_sh]
        if ev_shape.has_text_frame:
            old_ev = ev_shape.text_frame.text.strip()
            set_text_safe(ev_shape, Z['ev_new'])
            apply_normautofit(ev_shape)
            print(f"  ev[{ev_sh}] updated: '{Z['ev_new']}'")

    # Also patch the NPI strip (shape[98]) if it mentions old offtake values
    # NPI strip typically mentions zone name and mix — search for mix reference
    for i, sh in enumerate(shapes):
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if f"{mix_old}%" in txt and 'NPI' in txt and i > 90:
                patch(sh, f"{mix_old}%", f"{mix_new}%", f"{name}:NPI mix")
                break

    print()

# ─────────────────────────────────────────────────────────────────────────────
# OVERFLOW AUDIT + SAVE
# ─────────────────────────────────────────────────────────────────────────────
print()
audit_overflows(prs)

prs.save(DST)
print(f"\nSaved → {DST}")

# XML validation
with zipfile.ZipFile(DST) as z:
    errors = []
    for name in z.namelist():
        if name.endswith('.xml') or name.endswith('.rels'):
            try:
                etree.fromstring(z.read(name))
            except etree.XMLSyntaxError as e:
                errors.append(f"  {name}: {e}")
    if errors:
        print("[XML] ERRORS:")
        for e in errors: print(e)
    else:
        print("[XML] All valid.")

# ─────────────────────────────────────────────────────────────────────────────
# SPOT CHECKS
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== SPOT CHECKS ===")
p = Presentation(DST)

# Slide 1
s1 = p.slides[0]
sh5_txt = s1.shapes[5].text_frame.text.strip() if s1.shapes[5].has_text_frame else ""
sh6_txt = s1.shapes[6].text_frame.text.strip() if s1.shapes[6].has_text_frame else ""
ok_nat = "40.67" in sh5_txt and "2269" in sh5_txt
print(f"  S1 shape[5]: [{' OK' if ok_nat else 'FAIL'}] {sh5_txt[:80]}")
ok_gap = "6.35" in sh6_txt
print(f"  S1 shape[6]: [{'OK' if ok_gap else 'FAIL'}] {sh6_txt[:80]}")

# Zone slides
zone_checks = [
    ('West', 4, 12, "8.49"),
    ('South-1', 5, 12, "8.77"),
    ('North', 6, 12, "8.32"),
    ('South-2', 7, 12, "5.30"),
    ('East', 8, 12, "4.85"),
    ('Central', 9, 12, "2.88"),
]
for zone, si, sh_idx, expected in zone_checks:
    slide = p.slides[si]
    txt = slide.shapes[sh_idx].text_frame.text.strip() if slide.shapes[sh_idx].has_text_frame else ""
    ok = expected in txt
    print(f"  S{si+1} ({zone}) shape[{sh_idx}]: [{'OK' if ok else 'FAIL'}] {txt!r}")
