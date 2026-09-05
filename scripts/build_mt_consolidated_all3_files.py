#!/usr/bin/env python3
"""
Build consolidated MT July'26 deck merging all 3 PPT files
Strategy: Expand File 3 (base design) + embed File 1 (Nielsen) + File 2 (MT Operations)
Final: 9 core slides + optional appendix
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# =====================================================================
# Colour scheme (File 3 aesthetic: minimalist, monochromatic)
# =====================================================================
NAVY = RGBColor(0x1F, 0x4E, 0x78)
DARK_GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT_GREY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN = RGBColor(0x70, 0xAD, 0x47)
ACCENT_RED = RGBColor(0xC5, 0x50, 0x4E)

# =====================================================================
# Helper functions
# =====================================================================

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GREY

    return slide

def add_content_slide_eiao(prs, headline, evidence, implication, action, owner):
    """Add EIAO-structured content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Evidence
    y_pos = 1.1
    ev_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    tf = ev_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"EVIDENCE: {evidence}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED

    # Implication
    y_pos += 0.9
    imp_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    tf = imp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"IMPLICATION: {implication}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GREY

    # Action
    y_pos += 0.9
    act_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    tf = act_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"ACTION: {action}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    # Owner
    y_pos += 0.9
    owner_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
    tf = owner_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Owner: {owner}"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = LIGHT_GREY

    return slide

def add_table_slide(prs, headline, table_data, column_widths, note=""):
    """Add a slide with a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = headline_box.text_frame
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Table
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for i, col_width in enumerate(column_widths):
        table_shape.columns[i].width = Inches(col_width)

    # Populate table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = str(cell_data)

            # Header row formatting
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER
            else:
                # Alternate row colours
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREY
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.color.rgb = DARK_GREY

    # Footer note
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
        tf = note_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = LIGHT_GREY

    return slide

# =====================================================================
# Main: Build consolidated deck
# =====================================================================

def build_consolidated_all3():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Honasa Modern Trade", "Q1 FY27 Performance & Q2 Strategy | July 2026")

    # Slide 2: BEST EVER Q1 (File 3 Slide 1 + File 2 zone data)
    add_content_slide_eiao(
        prs,
        "Best Ever Q1 FY27 — ₹114.39 Cr, +64% YoY",
        "Q1 offtake ₹114.39 Cr (+64% YoY); July offtake ₹36.1 Cr. West/South-1 self-sustaining (82-84% conversion); North/East = 70% of national ₹10.92 Cr gap.",
        "Gap concentration: D-Mart + Reliance = 90.7% of recovery requirement. Two chains hold the strategy.",
        "Zone scorecard: West ₹8.28 Cr (82.3% conv), South-1 ₹8.19 Cr (83.6%), North ₹6.99 Cr (58.5% - FIX), East ₹3.55 Cr (45.3% - URGENT).",
        "Sales VP + Regional Leads | Weekly tracking"
    )

    # Slide 3: TDC ENGINE (File 3 Slide 2 + File 2 supply risk)
    add_content_slide_eiao(
        prs,
        "The Derma Co. Growth Engine — +365% YoY",
        "TDC ₹29.5 Cr (+365% YoY); Face Cleanser ₹7.13 Cr (growth hero); DMart ₹13.66 Cr, Reliance ₹7.31 Cr.",
        "Supply constraint: TDC ₹4.16 Cr flow gap @ 72.6% conversion. Do NOT assume linear 365% continues into Q2.",
        "Q2 action: Validate TDC supply capacity for Aug-Sep before aggressive primary billing. If constrained, prioritise highest-velocity SKUs (Face Cleanser + Acne).",
        "TDC Category Lead + Supply Chain | 15-Aug supply audit"
    )

    # Slide 4: FACE WASH DEEP DIVE (File 3 Slide 3 + File 1 Nielsen + File 2 data)
    add_content_slide_eiao(
        prs,
        "Face Wash Uplift — ₹32.59 Cr Q1, +33% YoY, Nielsen MAT +56%",
        "ME FW ₹32.59 Cr (+33% YoY); DMart share 27% (+10 pp). Nielsen external: MAT +56% YoY (vs category +10.9%), ME 11.2% share (+2.4 pp), WD 89.2%, ND 57.8%. 150 ml = category tail-wind (+59.7% YoY, 39.1% value).",
        "Internal D-Mart 27% > Nielsen MAT 11.2% = different bases. Use Nielsen for competitive health; internal for delivery. Rice ₹15.81 Cr leading because aligned to 150 ml pack trend.",
        "Validation gate: Reconcile WD/PDO/OOS by 31-Aug. Confirm 150 ml trail in Rice/Ubtan/Turmeric lines. Expand FSDU by 30% if Nielsen validates demand.",
        "Category Lead + Analytics | 31-Aug field reconciliation"
    )

    # Slide 5: SHAMPOO & SUN CARE (File 3 Slide 4 + File 1 Nielsen pack data + File 2 pilots)
    add_content_slide_eiao(
        prs,
        "Shampoo & Sun Care — Scale Headroom, Pack Gaps",
        "Shampoo ₹22.02 Cr (+65% YoY), Reliance ₹8.98 Cr driving; Sun Care ₹21.75 Cr (all brands best-ever). Nielsen: 650 ml = 72.1% of category value, +58% YoY; ME only has 3/16 formats (gap vs available 16).",
        "Pack white-space: Test 650/1000/180 ml in DMart + Apollo (Nielsen validated as growth tiers). L3M Shampoo +80.3% shows demand pull; format gap limits penetration.",
        "Q2 pilot: Deploy top-3 formats (650/1000/180) with trial packs (₹75–150 range). Scale only after 4-week velocity data. Onion/Rosemary revival = low-risk, high-upside play.",
        "Category Lead + Sales | 15-Sep pilot launch; 30-Sep velocity review"
    )

    # Slide 6: MARKET SHARE & DISTRIBUTION (File 3 Slide 5 + File 1 Nielsen + governance)
    add_content_slide_eiao(
        prs,
        "Share & Distribution Health — Reconciliation Framework",
        "ME FW 10.5% share (+3.1 pp), Shampoo 3.7% (+1.2 pp). Q1 sell-out 83.9% (₹21.9 Cr gap). Nielsen defines: WD = weekly distribution (ME 89.2%), ND = relative numerical dist (57.8%).",
        "Measurement gap: Internal D-Mart SAH vs Nielsen WD do not map 1:1. Use Nielsen (MAT) for competitive health; internal NSV for execution targets. Flag: Shampoo brand-level Nielsen absent (File 1 note).",
        "Governance move: Obtain Nielsen Shampoo brand-level data by 15-Aug. Field-validate SAH/WD/OOS by 31-Aug. Do NOT claim share changes until bases reconciled.",
        "Analytics + Category Lead + Sales | 31-Aug reconciliation audit"
    )

    # Slide 7: ZONE PERFORMANCE SCORECARD (File 2 consolidated)
    table_data = [
        ["Zone", "Jul Offtake (₹ Cr)", "Conversion %", "Gap (₹ Cr)", "Status", "Top Chain"],
        ["West", "8.28", "82.3%", "1.78", "WATCH", "D-Mart 94.2%"],
        ["South-1", "8.19", "83.6%", "1.61", "WATCH", "Apollo 81.0%"],
        ["North", "6.99", "58.5%", "4.97", "🔴 FIX", "D-Mart 58.5%"],
        ["South-2", "4.91", "71.3%", "1.98", "🔴 FIX", "D-Mart +269%"],
        ["East", "3.55", "45.3%", "4.28", "🔴 URGENT", "Reliance 44.9%"],
        ["Central", "2.12", "78.8%", "0.57", "WATCH", "D-Mart benchmark"],
    ]
    add_table_slide(
        prs,
        "Zone Performance Scorecard — Priority Actions",
        table_data,
        [1.2, 1.3, 1.2, 1.2, 1.2, 1.8],
        "North + East = 70.5% of national gap. Priority: Validate Reliance Paisa Vasool by 5-Aug before North recovery spending."
    )

    # Slide 8: Q2 DECISIONS (File 3 Slide 6 + 5 moves + governance + phasing)
    add_content_slide_eiao(
        prs,
        "Q2 FY27: Five Strategic Moves + 90-Day Phasing",
        "Move 1: Protect FW lead (FSDU +30%, validate Nielsen WD by 31-Aug). Move 2: Scale Shampoo smartly (pilot 3 formats in DMart/Apollo, 15-Sep). Move 3: Win velocity (PDO > category by 30-Sep). Move 4: Fix specific loopholes (Reliance Paisa Vasool Aug 5; North recovery after validation; East pause new loading).",
        "Move 5 (Governance): Separate measurement bases. Nielsen MAT = competitive health (external). Internal NSV = business delivery (owned metrics). Do NOT compare as same base. Weekly steering on both.",
        "Phasing: Phase 0-30d (reconcile fields, validate data, run Paisa Vasool). Phase 31-60d (size white-space, run pilots, field-validate Nielsen). Phase 61-90d (scale proven cells only). Gate each phase: no action without data validation.",
        "Sales VP + Category Lead + Analytics | Weekly steering; gates on 5-Aug / 31-Aug / 30-Sep"
    )

    # Slide 9: EXECUTION OWNERSHIP (File 2 RKAM/NKAM)
    table_data = [
        ["Role", "Owner", "Zone/Chain", "Target KPI", "Deadline"],
        ["RKAM", "Regional Lead", "West", "Maintain 82.3% conversion", "31-Aug"],
        ["RKAM", "Regional Lead", "North", "58.5% → 65% conversion", "30-Sep"],
        ["RKAM", "Regional Lead", "East", "45.3% → 50%+ conversion (after Aug validation)", "30-Sep"],
        ["NKAM", "Account Manager", "D-Mart", "₹4.29 Cr gap closure (Phase 1: South-2)", "31-Aug"],
        ["NKAM", "Account Manager", "Reliance", "Paisa Vasool validation; hold Haircare fixes until 5-Aug", "05-Aug"],
        ["Category", "Category Lead", "All", "Pack pilot launch; Nielsen reconciliation", "15-Sep / 31-Aug"],
    ]
    add_table_slide(
        prs,
        "Execution Ownership — RKAM/NKAM Accountability",
        table_data,
        [1.0, 1.5, 1.2, 1.8, 1.2],
        "Weekly scoreboard: Offtake vs. plan, Flow conversion % vs. target, Hero-SKU OSA > 95%."
    )

    # Save
    output_path = "/home/user/mt-dashboard/MT_Jul26_Honasa_CONSOLIDATED_All3_Files_v1.pptx"
    prs.save(output_path)
    print(f"\n✅ CONSOLIDATED DECK CREATED: {output_path}")
    print(f"   Slides: 9 core | Merged: File 3 (design) + File 1 (Nielsen) + File 2 (operations)")
    print(f"   Design preserved: File 3 aesthetic (minimal, text-based, monochromatic)")
    print(f"   Data integrated: Zone performance, Nielsen benchmarks, pack gaps, governance framework")
    return output_path

if __name__ == "__main__":
    path = build_consolidated_all3()
    print(f"\n✅ Ready for distribution. File size: ", end="")
    import os
    print(f"{os.path.getsize(path) / (1024*1024):.2f} MB")
