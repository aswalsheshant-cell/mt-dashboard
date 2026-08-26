"""Build V16 from V15 — fix footer layout on all slides.

Problems fixed:
1. Slide 3: 0.95" dead gap between content end (~11.67") and footer start (12.62").
   → Move footer block up to T=11.85" (0.18" gap after content).
2. Zone slides (5-10): Evidence boxes too short (0.46") for their content.
   NPI strip at 12.29", footer at 12.62" — only 0.71" of footer height.
   → Move NPI strip to T=12.03", footer container to T=12.30",
     expand evidence content boxes to H=0.80" (bot≈13.25").
3. All EVIDENCE content boxes: normAutofit + expanded height = readable text.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu
from lxml import etree

SRC = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv15.pptx"
DST = "/home/user/mt-dashboard/MT_Jul26_Honasa_Finalv16.pptx"

EMU_IN = 914400


def emu(inches: float) -> int:
    return int(inches * EMU_IN)


def get_shape_top(shape) -> int:
    return shape.top

def get_shape_left(shape) -> int:
    return shape.left

def get_shape_h(shape) -> int:
    return shape.height

def get_shape_w(shape) -> int:
    return shape.width


def set_top(shape, top_emu: int):
    sp = shape._element
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    xfrm = sp.find(f'.//{ns_a}xfrm')
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    if off is not None:
        off.set('y', str(top_emu))


def set_height(shape, h_emu: int):
    sp = shape._element
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    xfrm = sp.find(f'.//{ns_a}xfrm')
    if xfrm is None:
        return
    ext = xfrm.find(qn('a:ext'))
    if ext is not None:
        ext.set('cy', str(h_emu))


def apply_normautofit(shape):
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


def shapes_in_band(slide, top_min_in: float, top_max_in: float):
    """Return list of (idx, shape) whose top is within [top_min_in, top_max_in]."""
    result = []
    for i, sh in enumerate(slide.shapes):
        t = sh.top / EMU_IN
        if top_min_in <= t <= top_max_in:
            result.append((i, sh))
    return result


def audit_overflows(prs):
    slide_h = prs.slide_height / EMU_IN
    issues = []
    for si, slide in enumerate(prs.slides):
        for i, sh in enumerate(slide.shapes):
            bot = (sh.top + sh.height) / EMU_IN
            if bot > slide_h + 0.015:
                txt = sh.text_frame.text[:30] if sh.has_text_frame else ""
                issues.append(f"  S{si+1} shape[{i}] bot={bot:.3f}\" > {slide_h:.2f}\" | {txt!r}")
    if issues:
        print(f"[OVERFLOW] {len(issues)} issue(s):")
        for x in issues:
            print(x)
    else:
        print("[OVERFLOW] Clean — no overflows.")
    return len(issues) == 0


prs = Presentation(SRC)
slide_h_emu = int(prs.slide_height)
slide_h_in  = slide_h_emu / EMU_IN

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 (index 2) — Market share slide
# Gap: content ends ~11.67", footer starts at 12.62". Move footer to 11.87".
# ─────────────────────────────────────────────────────────────────────────────
print("=== Slide 3 — close footer gap ===")
slide3 = prs.slides[2]

FOOTER_START_OLD_S3 = 12.55   # all shapes at T > this are footer shapes
NEW_FOOTER_TOP_S3   = 11.87   # move footer container to here
shift_s3 = emu(NEW_FOOTER_TOP_S3 - 12.62)  # delta EMU (negative = up)

# Container (T≈12.62)
container_top = emu(NEW_FOOTER_TOP_S3)
container_h   = emu(1.38)   # extends to 11.87+1.38=13.25" (fits within 13.33")

# Inside footer: label row at +0.05" above content
label_top   = emu(NEW_FOOTER_TOP_S3 + 0.05)
divider_top = emu(NEW_FOOTER_TOP_S3 + 0.07)
divider_h   = emu(1.15)
content_top = emu(NEW_FOOTER_TOP_S3 + 0.24)
content_h   = emu(1.05)    # was 0.38" — expand to 1.05"

moved = 0
for i, sh in enumerate(slide3.shapes):
    t_in = sh.top / EMU_IN
    if t_in < FOOTER_START_OLD_S3:
        continue
    w = sh.width / EMU_IN
    h_in = sh.height / EMU_IN
    txt = sh.text_frame.text.strip() if sh.has_text_frame else ""

    if h_in < 0.05:          # zero-width divider lines
        set_top(sh, divider_top)
        set_height(sh, divider_h)
    elif any(lbl in txt for lbl in ['EVIDENCE', 'IMPLICATION', 'ACTION', 'OWNER']) and h_in < 0.25:
        # label row
        set_top(sh, label_top)
    elif w > 5.0:             # wide container box
        set_top(sh, container_top)
        set_height(sh, container_h)
    else:                     # content boxes
        set_top(sh, content_top)
        set_height(sh, content_h)
        apply_normautofit(sh)
    moved += 1

print(f"  Moved/resized {moved} footer shapes; gap closed to 0.20\"")

# ─────────────────────────────────────────────────────────────────────────────
# ZONE SLIDES (indices 4-9) — Move NPI strip up, expand evidence section
# ─────────────────────────────────────────────────────────────────────────────
# Layout targets:
#   NPI strip [98]/[97]: T=12.29" → 12.03"
#   Footer container:    T=12.62" → 12.30"  H=0.62" → 1.00"
#   Label row:                       12.67" → 12.35"
#   Divider lines (W=0):             12.70" → 12.37"  H=0.33" → 0.80"
#   Content boxes:                   12.83" → 12.55"  H=0.46" → 0.73"
#   (bot = 12.55+0.73=13.28" — within 13.33")
#
NPI_NEW_TOP   = emu(12.03)
CONT_NEW_TOP  = emu(12.30)
CONT_NEW_H    = emu(1.00)
LABEL_NEW_TOP = emu(12.35)
DIV_NEW_TOP   = emu(12.37)
DIV_NEW_H     = emu(0.80)
EVID_NEW_TOP  = emu(12.55)
EVID_NEW_H    = emu(0.73)

for si in [4, 5, 6, 7, 8, 9]:
    slide = prs.slides[si]
    n_shapes = len(slide.shapes)
    print(f"\n=== Slide {si+1} — footer lift + evidence expand ({n_shapes} shapes) ===")
    moved = 0

    for i, sh in enumerate(slide.shapes):
        t_in = sh.top / EMU_IN
        if t_in < 12.00:
            continue                       # don't touch anything above 12"

        w_in = sh.width  / EMU_IN
        h_in = sh.height / EMU_IN
        txt  = sh.text_frame.text.strip() if sh.has_text_frame else ""

        # NPI mix summary strip (full-width, short H, T≈12.29")
        if 12.20 <= t_in <= 12.40 and w_in > 5.0 and h_in < 0.40:
            set_top(sh, NPI_NEW_TOP)
            moved += 1
            continue

        # Container box (full-width, T≈12.62")
        if t_in > 12.50 and w_in > 5.0:
            set_top(sh, CONT_NEW_TOP)
            set_height(sh, CONT_NEW_H)
            moved += 1
            continue

        # Label headers (EVIDENCE / IMPLICATION / ACTION / OWNER), T≈12.67"
        if t_in > 12.55 and any(lbl in txt for lbl in ['EVIDENCE', 'IMPLICATION', 'ACTION', 'OWNER']):
            set_top(sh, LABEL_NEW_TOP)
            moved += 1
            continue

        # Vertical divider lines (W=0), T≈12.70"
        if t_in > 12.60 and w_in < 0.05:
            set_top(sh, DIV_NEW_TOP)
            set_height(sh, DIV_NEW_H)
            moved += 1
            continue

        # Evidence content boxes (four columns), T≈12.83"
        if t_in > 12.70 and w_in < 2.0 and h_in > 0.10:
            set_top(sh, EVID_NEW_TOP)
            set_height(sh, EVID_NEW_H)
            apply_normautofit(sh)
            moved += 1
            continue

    print(f"  Repositioned {moved} shapes; evidence boxes now H=0.73\" at T=12.55\"")

# ─────────────────────────────────────────────────────────────────────────────
# Also lift & expand evidence boxes on slides not yet touched (slides 1-2, 11-18)
# that have EVIDENCE footer blocks
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Other slides — expand evidence boxes where present ===")
done_slides = {0, 1, 2, 4, 5, 6, 7, 8, 9}
for si, slide in enumerate(prs.slides):
    if si in done_slides:
        continue
    modified = 0
    for i, sh in enumerate(slide.shapes):
        t_in = sh.top / EMU_IN
        if t_in < 12.50:
            continue
        h_in = sh.height / EMU_IN
        w_in = sh.width  / EMU_IN
        txt  = sh.text_frame.text.strip() if sh.has_text_frame else ""
        # Content boxes in the evidence band
        if t_in > 12.70 and w_in < 2.0 and h_in > 0.10:
            if h_in < 0.50:
                set_height(sh, emu(0.50))
                apply_normautofit(sh)
                modified += 1
    if modified:
        print(f"  Slide {si+1}: expanded {modified} evidence content box(es)")

# ─────────────────────────────────────────────────────────────────────────────
# Overflow audit + save
# ─────────────────────────────────────────────────────────────────────────────
print()
audit_overflows(prs)

prs.save(DST)
print(f"\nSaved → {DST}")


# ─────────────────────────────────────────────────────────────────────────────
# Validate key positions
# ─────────────────────────────────────────────────────────────────────────────
def validate(path):
    import zipfile
    from lxml import etree as _et

    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            data = z.read(name)
            if name.endswith('.xml') or name.endswith('.rels'):
                try:
                    _et.fromstring(data)
                except _et.XMLSyntaxError as e:
                    print(f"  XML ERROR in {name}: {e}")
                    return False

    p = Presentation(path)
    print(f"\n=== VALIDATE {path} ===")
    checks = [
        # (slide_idx, shape_idx, expected_text)
        (0,  5,  '2114K units'),
        (0,  6,  '+98.9% YoY'),
        (4,  8,  '+99.0% YoY'),
        (4, 13,  '466K units'),
        (5,  8,  '+63.6% YoY'),
        (6,  8,  '+134.3% YoY'),
        (7,  8,  '+69.3% YoY'),
        (8,  8,  '+115.7% YoY'),
        (9,  8,  '+202.2% YoY'),
    ]
    ok = True
    for si, shi, exp in checks:
        sh = p.slides[si].shapes[shi]
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        status = "OK" if exp in txt else "FAIL"
        if status == "FAIL":
            ok = False
        print(f"  S{si+1}:{shi} [{status}] '{exp}' → '{txt[:60]}'")

    # Position checks for zone slide 5 evidence box
    sl5 = p.slides[4]
    ev_box = sl5.shapes[101]
    ev_top = ev_box.top / EMU_IN
    ev_h   = ev_box.height / EMU_IN
    ev_bot = ev_top + ev_h
    pos_ok = ev_top < 12.70 and ev_bot < 13.33
    print(f"  S5 ev_box[101]: top={ev_top:.2f}\" H={ev_h:.2f}\" bot={ev_bot:.2f}\" "
          f"[{'OK' if pos_ok else 'FAIL'}]")
    if not pos_ok:
        ok = False

    # Slide 3 footer gap check — footer container should be at T < 12.10"
    sl3 = p.slides[2]
    footer_tops = [sh.top / EMU_IN for sh in sl3.shapes
                   if sh.width / EMU_IN > 5.0 and sh.top / EMU_IN > 11.5]
    if footer_tops:
        fmin = min(footer_tops)
        gap_ok = fmin < 12.10
        print(f"  S3 footer top={fmin:.2f}\" [{'OK' if gap_ok else 'FAIL — gap still large'}]")
        if not gap_ok:
            ok = False

    return ok


valid = validate(DST)
print(f"\n{'ALL CHECKS PASSED' if valid else 'SOME CHECKS FAILED'}")
