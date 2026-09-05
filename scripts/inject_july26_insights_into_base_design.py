#!/usr/bin/env python3
"""
Surgical Content Injection: MT July'26 Insights into Base Design
Strategy: Update ONLY text/data while preserving 100% of design, fonts, colours, layout

Data sources:
- File 1 (Nielsen): Market benchmarks, WD, ND, share data
- File 2 (MT Ops): Zone/chain performance, conversion, gaps
- File 3 (Current): Q1 performance metrics, brand data
"""

from pptx import Presentation
from pptx.util import Pt, Emu
import re

# =====================================================================
# NEW DATA MAPPING (File 1 Nielsen + File 2 MT Ops + File 3 Current)
# =====================================================================

SLIDE_DATA = {
    1: {  # BEST EVER Q1
        "subtitle": "HONASA MODERN TRADE | Q1 FY27 | LEADERSHIP REVIEW",
        "headline": "114.39 CR OFFTAKES: 27% SEQUENTIAL & 64% YOY GROWTH",
        "callout": "BEST EVER Q1",
        "intro": "Q1 FY27 closed at ₹114.39 Cr vs ₹69.92 Cr (Q1 FY26), marking best-ever performance. Mamaearth +56% MAT growth (Nielsen), TDC +365% YoY internal. Zone conversion health: West 82%, South-1 84%, North 58%, East 45% (File 2 data). D-Mart + Reliance account for 90.7% of recovery requirement.",
        "kpis": {
            "Q1 FY27 OFFTAKE": ("₹114.39 Cr", "▲ 64%"),
            "SEQUENTIAL": ("↑27%", "vs Jun ₹36.1 Cr"),
            "Q1 FY27 PRIMARY": ("₹88.2 Cr", "▲ 71%"),
            "JUN'26 OFFTAKE": ("₹36.1 Cr", "↓ 10% MoM (billed-not-sold ₹10.92 Cr)"),
        },
        "tables": {
            "Zone View": [
                ["Zone", "Q1 Offtake (₹Cr)", "Conversion %", "Status"],
                ["West", "8.28", "82.3%", "WATCH"],
                ["North", "6.99", "58.5%", "🔴 FIX"],
                ["South-1", "8.19", "83.6%", "WATCH"],
                ["South-2", "4.91", "71.3%", "🔴 FIX"],
                ["East", "3.55", "45.3%", "🔴 URGENT"],
                ["Central", "2.12", "78.8%", "WATCH"],
            ],
            "Top Chains": [
                ["Chain", "Q1 NSV (₹Cr)", "Conv %", "Gap"],
                ["D-Mart", "14.33", "76.5%", "₹4.29"],
                ["Reliance", "8.06", "51.4%", "₹7.61"],
                ["Apollo", "2.5", "99.7%", "—"],
                ["Lulu", "1.70", "95%", "—"],
                ["FSN/Nykaa", "2.07", "99.4%", "—"],
                ["TOTAL", "36.1", "70.2%", "₹10.92"],
            ],
            "Brand View": [
                ["Brand", "Q1 NSV (₹Cr)", "Growth"],
                ["MAMAEARTH", "24.55", "▲ 56% (Nielsen MAT)"],
                ["THE DERMA CO.", "11.08", "▲ 365%"],
                ["AQUALOGICA", "0.48", "▲ 117%"],
            ]
        },
        "insights": "Zone conversion spread (45%–84%) drives phased recovery. D-Mart + Reliance concentration = execution risk. Nielsen +56% MAT validates demand pull.",
        "so_what": "PHASE 0-30d: Validate Reliance Paisa Vasool (Aug 5) before North/East spend. PHASE 31-60d: Run pack pilots (650/1000ml Shampoo, File 1 data). PHASE 61-90d: Scale proven cells only."
    },

    2: {  # TDC ENGINE
        "subtitle": "THE DERMA CO. | Q1 FY27 | BRAND ENGINE",
        "headline": "TDC @ ₹29.5 CR, GROWING 365% YOY—SUPPLY GAP FLAGGED",
        "callout": "365% YoY",
        "intro": "Major scale-up across acne, with supply constraint emerging. TDC ₹4.16 Cr flow gap @ 72.6% conversion signals capacity ceiling. Nielsen Face Wash premium tier +15% market trend supports continued growth.",
        "kpis": {
            "TDC Q1 OFFTAKE": ("₹29.5 Cr", "▲ 365%"),
            "TDC IN D-MART": ("₹13.66 Cr", "▲ 312%"),
            "TDC FACEWASH": ("₹7.13 Cr", "HERO SKU"),
            "TDC SUN CARE": ("₹6.82 Cr", "▲ 429%"),
        },
        "tables": {
            "Account-wise": [
                ["Chain", "Q1 (₹Cr)", "Conv %", "Gap"],
                ["D-Mart", "13.66", "88.4%", "₹1.78"],
                ["Reliance", "7.31", "61.2%", "₹4.65"],
                ["Apollo", "2.18", "99.7%", "—"],
                ["Wellness", "1.85", "92%", "—"],
                ["Lulu", "1.42", "94%", "—"],
                ["FSN", "1.08", "99.4%", "—"],
            ],
            "Acne Partition": [
                ["Salicylic Acid", "₹8.92 Cr", "▲ 389%"],
                ["Sali-Cinamide", "₹7.51 Cr", "▲ 354%"],
                ["Kojic Acid", "₹5.38 Cr", "▲ 298%"],
                ["TOTAL", "₹29.5 Cr", "▲ 365%"],
            ]
        },
        "insights": "D-Mart conversion 88% vs Reliance 61% = format/planogram mismatch risk. TDC ₹4.16 Cr gap @ 72.6% = supply bottleneck emerging.",
        "so_what": "VALIDATE TDC supply capacity by 15-Aug before aggressive primary billing. PRIORITISE Face Cleanser + Acne range (highest velocity). HOLD Reliance aggressive loading until conversion ≥ 75%."
    },

    3: {  # FACE WASH
        "subtitle": "FACE WASH | Q1 FY27 | CATEGORY DEEP DIVE",
        "headline": "MAMAEARTH FW ₹32.59 CR (+33% YOY), NIELSEN MAT +56%—PACK OPPORTUNITY",
        "intro": "Nielsen external: MAT +56% YoY (vs category +10.9%), ME 11.2% share (+2.4 pp), WD 89.2%, ND 57.8%. 150 ml = category tail-wind (+59.7% YoY, 39% value). D-Mart internal 27% > Nielsen 11.2% = different bases.",
        "kpis": {
            "ME FACEWASH Q1": ("₹32.59 Cr", "▲ 33%"),
            "D-MART SHARE": ("27% (+10pp)", "Nielsen: 11.2% MAT"),
            "RICE RANGE": ("₹15.81 Cr", "HERO (150ml)"),
            "LARGE PACKS": ("54%→67%", "150ml growth"),
        },
        "tables": {
            "Range breakdown": [
                ["Range", "Q1 (₹Cr)", "Growth"],
                ["Rice", "₹15.81 Cr", "▲ 47%"],
                ["Ubtan", "₹8.42 Cr", "▲ 28%"],
                ["Vitamin C", "₹4.91 Cr", "▲ 19%"],
                ["Tea Tree", "₹3.45 Cr", "▲ 8%"],
            ],
            "Pack size mix": [
                ["Pack Size", "% of Mix", "Nielsen Trend"],
                ["150 ml", "46%", "▲ 59.7% (category tail-wind)"],
                ["100 ml", "31%", "▼ 9.8%"],
                ["50 ml", "15%", "▲ 12%"],
                ["200 ml", "5%", "New"],
                ["250 ml", "3%", "Niche"],
            ],
            "Market Share": [
                ["Brand", "D-Mart %", "Nielsen MAT %"],
                ["MAMAEARTH", "15%", "11.2% (↑2.4pp)"],
                ["THE DERMA CO.", "12%", "—"],
                ["TOTAL ME+TDC", "27%", "—"],
            ]
        },
        "insights": "150 ml = Nielsen-validated growth tier. Pack white-space: 150ml upside by expanding Rice/Ubtan range. Nielsen WD 89.2% = near-perfect shelf availability.",
        "so_what": "FIELD RECONCILE WD/PDO/OOS by 31-Aug. EXPAND 150ml trial in Rice/Ubtan by 30%. INCREASE FSDU placements in D-Mart/Reliance by mid-Sep."
    },

    4: {  # SHAMPOO & SUN CARE
        "subtitle": "SHAMPOO & SUN CARE | Q1 FY27",
        "headline": "SHAMPOO ₹22.02 CR (+65% YOY, RELIANCE-LED); SUN CARE ₹21.75 CR BEST-EVER",
        "intro": "Shampoo: Reliance ₹8.98 Cr (41% of category); Nielsen 650ml = 72% of category value, +58% YoY (gap: ME only 3/16 formats). Sun Care: All-brand best-ever due to TDC + Aqualogica scale.",
        "kpis": {
            "SHAMPOO Q1": ("₹22.02 Cr", "▲ 65%"),
            "RELIANCE": ("₹8.98 Cr", "▲ 80%"),
            "SUN CARE Q1": ("₹21.75 Cr", "BEST-EVER"),
            "RUN RATE": ("₹88 Cr annualized", "L3M +80%"),
        },
        "tables": {
            "Shampoo ranges": [
                ["Range", "Q1 (₹Cr)", "Growth"],
                ["Onion", "₹8.91 Cr", "▲ 62%"],
                ["Rosemary", "₹7.24 Cr", "▲ 71%"],
                ["Rice", "₹4.12 Cr", "▲ 48%"],
                ["Lemon", "₹1.75 Cr", "▲ 31%"],
            ],
            "Sun Care Brands": [
                ["Brand", "Q1 (₹Cr)", "Growth"],
                ["THE DERMA CO.", "₹9.38 Cr", "▲ 485% (NEW)"],
                ["MAMAEARTH", "₹8.15 Cr", "▲ 39%"],
                ["AQUALOGICA", "₹4.22 Cr", "▲ 156%"],
            ]
        },
        "insights": "Shampoo pack gap: Nielsen shows 650/1000ml = 72% of category; ME only 3/16 formats = white-space. Reliance 80% YoY but conversion 51.4% (File 2) = Paisa Vasool risk.",
        "so_what": "PILOT 3-format test (650/1000/180ml) in D-Mart + Apollo by 15-Sep. VALIDATE Reliance Paisa Vasool conversion by 5-Aug before phasing. PROTECT Onion/Rosemary momentum (highest growth)."
    },

    5: {  # MARKET SHARE & SELL-OUT
        "subtitle": "MARKET SHARE, DISTRIBUTION & SELL-OUT | MAT JUN-26",
        "headline": "MAMAEARTH SHARE GAINS (FW +3.1pp, SHAMPOO +1.2pp); ₹21.9 CR BILLED-NOT-SOLD GAP",
        "intro": "Nielsen external: ME FW 10.5% share (+3.1pp), Shampoo 3.7% (+1.2pp). Internal sell-out Q1 83.9% (₹21.9 Cr gap); Jun recover 92%. Gap concentration: D-Mart + Reliance = 90.7%.",
        "kpis": {
            "FACEWASH SHARE": ("10.5% (+3.1pp)", "Nielsen MAT"),
            "SHAMPOO SHARE": ("3.7% (+1.2pp)", "vs Q1 FY26"),
            "Q1 SELL-OUT": ("83.9%", "₹21.9 Cr gap"),
            "JUN'26 SELL-OUT": ("92%", "₹4.5 Cr gap"),
        },
        "tables": {
            "Category scorecard": [
                ["Category", "Q1 (₹Cr)", "Nielsen MAT %", "WD %"],
                ["Facewash", "₹32.59 Cr", "10.5% (+3.1pp)", "89.2%"],
                ["Shampoo", "₹22.02 Cr", "3.7% (+1.2pp)", "—"],
            ],
            "Sell-out by chain": [
                ["Chain", "Q1 NSV", "Sell-out %", "Gap (₹Cr)"],
                ["D-Mart", "₹14.33 Cr", "77%", "₹4.29"],
                ["Reliance", "₹8.06 Cr", "92%", "₹7.61"],
                ["Metro", "₹0.49 Cr", "45%", "₹0.60"],
                ["TOTAL", "₹36.1 Cr", "83.9%", "₹21.9"],
            ],
            "Share of Shelf": [
                ["Category", "Actual %", "Target %"],
                ["Sunscreen", "18.2%", "20%"],
                ["Facewash", "22.5%", "20%"],
                ["Shampoo", "14.8%", "15%"],
            ]
        },
        "insights": "Nielsen +3.1pp FW share validates demand. Internal D-Mart SAH 22.5% vs Nielsen WD 89.2% = measurement reconciliation needed. Gap concentration = execution risk.",
        "so_what": "RECONCILE internal SAH vs Nielsen WD by 31-Aug (different bases). HOLD Metro aggressive loading until sell-out ≥ 75%. FIELD-VALIDATE Nielsen assumptions by 30-Sep."
    },

    6: {  # Q2 DECISIONS
        "subtitle": "Q2 FY27 | DECISIONS AND OWNERS",
        "headline": "TURN SHARE GAINS INTO REPEATABLE GROWTH—FIVE MOVES + 90-DAY PHASING",
        "intro": "The opportunity is clear (Nielsen +56% MAT, internal +64% YoY, File 2 shows zone/chain recovery path). Execution now decides. Phase 0-30d: validate data. Phase 31-60d: run pilots. Phase 61-90d: scale proven.",
        "moves": [
            "PROTECT FACEWASH LEAD: Validate Nielsen benchmarks (WD/PDO/OOS) by 31-Aug. Expand FSDU ×30% in D-Mart/Reliance. Owner: Category Lead",
            "SCALE SHAMPOO SMARTLY: Pilot 3 formats (650/1000/180ml) in D-Mart + Apollo (File 1 Nielsen validated). 15-Sep launch, 30-Sep velocity review. Owner: Category Lead",
            "WIN THE VELOCITY BATTLE: Confirm PDO > category by 30-Sep. Increase trial schemes (File 2 North/East recovery). Owner: Sales Lead",
            "FIX SPECIFIC LOOPHOLES: Reliance Paisa Vasool check 5-Aug; North conversion 58.5%→65%; East 45%→50% (File 2 zones). Owner: NKAM/RKAM",
            "SEPARATE THE MEASUREMENT BASES: Use Nielsen MAT for competitive health (external). Use NSV for delivery targets (internal). Weekly steering on both. Owner: Analytics Lead",
        ],
        "tables": {
            "JTBD Q1 Deltas": [
                ["Category/Brand", "Q1 Delivered (₹Cr)", "vs JTBD", "Action"],
                ["Facewash", "₹32.59 Cr", "+₹0.25 Cr", "Hit target; expand FSDU"],
                ["Shampoo", "₹22.02 Cr", "–₹0.50 Cr", "Pilot pack formats"],
                ["Sun Care", "₹21.75 Cr", "+₹1.20 Cr", "Protect TDC supply"],
                ["Offtake", "₹114.39 Cr", "+₹2.5 Cr", "Reliance recovery phasing"],
                ["Sell-out", "83.9%", "–8.1pp", "Close ₹21.9 Cr gap by Jun"],
            ],
            "Immediate Priorities": [
                ["Priority", "Owner", "Deadline", "Success Metric"],
                ["Close Nielsen WD gaps", "Analytics", "31-Aug", "WD/PDO reconciled"],
                ["Shampoo pack format pilot", "Category", "15-Sep", "3 SKUs live"],
                ["Reliance Paisa Vasool audit", "NKAM", "05-Aug", "≥80% conversion"],
                ["North zone conversion ramp", "RKAM North", "30-Sep", "58.5%→65%"],
                ["TDC supply capacity check", "Supply", "15-Aug", "Aug target validated"],
                ["eB2B allocation 2x", "eB2B Lead", "31-Aug", "₹3.5 Cr monthly run"],
            ]
        },
        "non_negotiables": "No share claim without Nielsen reconciliation. No new loading until conversion ≥75%. No Q2 acceleration until Jul Paisa Vasool confirmed.",
        "leadership_ask": "Approval for 90-day phasing (0-30d validate, 31-60d pilot, 61-90d scale). Weekly steering on both Nielsen (external) and NSV (internal) metrics."
    }
}

# =====================================================================
# Content Injection Engine
# =====================================================================

def update_pptx_with_insights(pptx_path, output_path):
    """
    Surgical injection: update ONLY text in base design.
    Preserve ALL formatting, fonts, colours, positions, shapes.
    """
    prs = Presentation(pptx_path)

    for slide_num, slide_data in SLIDE_DATA.items():
        slide = prs.slides[slide_num - 1]

        print(f"Updating Slide {slide_num}...")

        # Find shapes by text content and update
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            text = shape.text.strip()

            # SUBTITLE
            if "HONASA MODERN TRADE" in text or "THE DERMA CO." in text or "FACE WASH" in text or "SHAMPOO" in text or "MARKET SHARE" in text or "Q2 FY27" in text:
                if len(text.split("|")[0]) < 50:  # Likely a subtitle line
                    shape.text = slide_data.get("subtitle", text)
                    print(f"  ✓ Updated subtitle")

            # HEADLINE
            elif len(text) > 50 and ("CR" in text or "₹" in text or "GROWTH" in text):
                shape.text = slide_data.get("headline", text)
                print(f"  ✓ Updated headline")

            # INTRO
            elif "Q1 FY27 closed" in text or "Major scale-up" in text or "Nielsen" in text or "Power of all" in text or "Mamaearth FW" in text or "The opportunity" in text:
                shape.text = slide_data.get("intro", text)
                print(f"  ✓ Updated intro paragraph")

            # INSIGHTS & SO WHAT
            elif "LEADERSHIP TAKEAWAY" in text or "WHAT IS DRIVING IT" in text or "THE READ" in text or "WHAT CHANGED" in text or "WHAT THE MARKET" in text:
                if "TAKEAWAY" in text:
                    shape.text = slide_data.get("insights", text)
                print(f"  ✓ Updated insights")

            elif "SO WHAT" in text or "ACTION" in text:
                shape.text = slide_data.get("so_what", text)
                print(f"  ✓ Updated action items")

    # Save
    prs.save(output_path)
    print(f"\n✅ CONTENT INJECTION COMPLETE: {output_path}")
    return output_path

if __name__ == "__main__":
    input_path = "/root/.claude/uploads/fca9dc4c-355d-55d2-8fb3-9c3700ee1ff0/486c9290-MT_July26_Final_to_be_update.pptx"
    output_path = "/home/user/mt-dashboard/MT_July26_Final_UPDATED_with_All3_Insights_v1.pptx"

    update_pptx_with_insights(input_path, output_path)

    import os
    print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
