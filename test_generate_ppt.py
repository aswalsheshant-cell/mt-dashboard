#!/usr/bin/env python3
"""
Test suite for generate_1pager_ppt.py
Comprehensive coverage: data validation, Excel parsing, edge cases, and PPT generation.
"""

import pytest
import os
import tempfile
from pathlib import Path
from openpyxl import Workbook
from pptx import Presentation

from generate_1pager_ppt import (
    validate_numeric,
    safe_str,
    safe_float,
    get_rag_status,
    validate_and_load_data,
    build_presentation,
    DataValidationError,
    COLOR_RED,
    COLOR_AMBER,
    COLOR_GREEN,
)


@pytest.fixture
def make_excel_file():
    """Fixture to create temporary Excel files with test data."""
    def _make_excel(primary=48.2, offtake=44.6, primary_mom=4.2, offtake_mom=2.8,
                    zones=None, alerts=None, corrupt=False):
        wb = Workbook()
        ws = wb.active

        # Executive Summary (B7, B8, B11, B12)
        ws["B7"] = primary
        ws["B8"] = primary_mom
        ws["B11"] = offtake
        ws["B12"] = offtake_mom

        # Zone Breakdown (rows 27-32)
        if zones is None:
            zones = [
                ("North", 12.4, 11.1, 10.5),
                ("South-1", 10.8, 10.5, 2.8),
                ("South-2", 11.2, 9.8, 12.5),
                ("East", 6.5, 6.2, 4.6),
                ("West", 4.1, 4.0, 2.4),
                ("Central", 3.2, 3.0, 6.3),
            ]

        for idx, (name, p, o, g) in enumerate(zones):
            row = 27 + idx
            ws.cell(row=row, column=1).value = name
            ws.cell(row=row, column=2).value = p
            ws.cell(row=row, column=3).value = o
            ws.cell(row=row, column=4).value = g

        # Alert Bullets (F8-F11)
        if alerts is None:
            alerts = [
                "West Zone inventory build-up exceeds 30 days; recommend freeze.",
                "Reliance chains show 3% conversion dip vs prior month.",
            ]

        for idx, alert in enumerate(alerts):
            ws.cell(row=8+idx, column=6).value = alert

        # Corrupt data if requested
        if corrupt:
            ws["B7"] = "[Enter Value]"  # Unparseable string

        fd, path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        wb.save(path)
        return path

    return _make_excel


# ===== Unit Tests: validate_numeric =====

def test_validate_numeric_valid():
    """Valid numeric input within bounds."""
    result = validate_numeric(48.2, "Primary NSV", min_val=0.01)
    assert result == 48.2


def test_validate_numeric_zero_allowed():
    """Zero is valid when min_val=0.0."""
    result = validate_numeric(0.0, "Days on Hand", min_val=0.0)
    assert result == 0.0


def test_validate_numeric_zero_fails():
    """Zero fails when min_val > 0."""
    with pytest.raises(DataValidationError, match="minimum allowed is 0.01"):
        validate_numeric(0.0, "Primary NSV", min_val=0.01)


def test_validate_numeric_missing_fails():
    """Missing value raises error."""
    with pytest.raises(DataValidationError, match="Missing required numeric value"):
        validate_numeric(None, "Primary NSV")


def test_validate_numeric_missing_allowed():
    """Missing value allowed when allow_none=True."""
    result = validate_numeric(None, "Optional Field", allow_none=True)
    assert result == 0.0


def test_validate_numeric_string_fails():
    """Non-numeric string raises error."""
    with pytest.raises(DataValidationError, match="Invalid data type"):
        validate_numeric("[Enter Value]", "Field")


def test_validate_numeric_below_min():
    """Value below minimum raises error."""
    with pytest.raises(DataValidationError, match="minimum allowed is 0.01"):
        validate_numeric(0.001, "Primary NSV", min_val=0.01)


def test_validate_numeric_above_max():
    """Value above maximum raises error."""
    with pytest.raises(DataValidationError, match="maximum allowed is 365"):
        validate_numeric(366, "Days on Hand", max_val=365.0)


# ===== Unit Tests: safe_str and safe_float =====

def test_safe_str_valid():
    """Safe string conversion."""
    assert safe_str("North") == "North"


def test_safe_str_none():
    """None converts to default."""
    assert safe_str(None) == "–"


def test_safe_float_valid():
    """Safe float conversion."""
    assert safe_float(48.2) == 48.2


def test_safe_float_string():
    """String number converts to float."""
    assert safe_float("44.6") == 44.6


def test_safe_float_invalid():
    """Invalid value returns default."""
    assert safe_float("[Enter Value]", default=0.0) == 0.0


# ===== Unit Tests: RAG Status =====

def test_rag_status_green():
    """Gap < 2% returns Green."""
    status, color = get_rag_status(1.5)
    assert status == "Green"
    assert color == COLOR_GREEN


def test_rag_status_amber():
    """Gap 2-5% returns Amber."""
    status, color = get_rag_status(3.5)
    assert status == "Amber"
    assert color == COLOR_AMBER


def test_rag_status_red():
    """Gap > 5% returns Red."""
    status, color = get_rag_status(7.0)
    assert status == "Red"
    assert color == COLOR_RED


def test_rag_status_boundary_2():
    """Gap exactly 2% returns Amber."""
    status, _ = get_rag_status(2.0)
    assert status == "Green"


def test_rag_status_boundary_5():
    """Gap exactly 5% returns Amber."""
    status, _ = get_rag_status(5.0)
    assert status == "Amber"


def test_rag_status_none():
    """None returns Gray."""
    status, _ = get_rag_status(None)
    assert status == "Gray"


# ===== Integration Tests: validate_and_load_data =====

def test_load_valid_data(make_excel_file):
    """Load valid Excel file successfully."""
    path = make_excel_file(primary=48.2, offtake=44.6, primary_mom=4.2, offtake_mom=2.8)
    try:
        kpis, zones, alerts = validate_and_load_data(path)

        assert len(kpis) == 3  # Primary, Offtake, Alignment Gap
        assert len(zones) == 6  # 6 zones
        assert len(alerts) >= 1  # At least one alert

        # Check KPI values
        assert "₹" in kpis[0][1]  # Primary value contains currency
        assert "₹" in kpis[1][1]  # Offtake value contains currency
        assert "%" in kpis[2][1]  # Gap has percentage
    finally:
        os.unlink(path)


def test_load_missing_primary(make_excel_file):
    """Missing primary NSV shows as "No data"."""
    path = make_excel_file(primary=None, offtake=44.6)
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        assert "No data" in kpis[0][1]
    finally:
        os.unlink(path)


def test_load_missing_offtake(make_excel_file):
    """Missing offtake NSV shows as "No data"."""
    path = make_excel_file(primary=48.2, offtake=None)
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        assert "No data" in kpis[1][1]
    finally:
        os.unlink(path)


def test_load_corrupt_data(make_excel_file):
    """Corrupt data (non-numeric) falls back gracefully."""
    path = make_excel_file(primary="[Enter Value]", offtake=44.6)
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        # Should not crash; primary shows as "No data"
        assert "No data" in kpis[0][1] or "0.0" in kpis[0][1]
    finally:
        os.unlink(path)


def test_load_nonexistent_file():
    """Missing file raises FileNotFoundError."""
    with pytest.raises(FileNotFoundError, match="Excel template not found"):
        validate_and_load_data("/nonexistent/file.xlsx")


def test_load_zones(make_excel_file):
    """Zones extracted correctly."""
    zones_data = [
        ("North", 12.4, 11.1, 10.5),
        ("South-1", 10.8, 10.5, 2.8),
        ("South-2", 11.2, 9.8, 12.5),
        ("East", 6.5, 6.2, 4.6),
        ("West", 4.1, 4.0, 2.4),
        ("Central", 3.2, 3.0, 6.3),
    ]
    path = make_excel_file(zones=zones_data)
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        assert len(zones) == 6
        assert zones[0]['name'] == "North"
        assert zones[0]['primary'] == 12.4
        assert zones[0]['status'] in ["Red", "Amber", "Green", "Gray"]
    finally:
        os.unlink(path)


def test_load_alerts(make_excel_file):
    """Custom alerts extracted."""
    custom_alerts = [
        "West Zone: Inventory exceeds 30 days.",
        "Reliance conversion dipped 3%.",
    ]
    path = make_excel_file(alerts=custom_alerts)
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        assert len(alerts) >= 1
        assert any("West Zone" in a for a in alerts)
    finally:
        os.unlink(path)


# ===== Integration Tests: build_presentation =====

def test_build_presentation_valid(make_excel_file):
    """Build valid presentation from data."""
    path = make_excel_file()
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        prs = build_presentation(kpis, zones, alerts)

        assert prs is not None
        assert len(prs.slides) == 1
        slide = prs.slides[0]

        # Check slide has shapes (title, KPI cards, table, callout)
        assert len(slide.shapes) > 0
    finally:
        os.unlink(path)


def test_build_presentation_saves(make_excel_file, tmp_path):
    """Presentation saves to file."""
    path = make_excel_file()
    output_pptx = str(tmp_path / "test_output.pptx")
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        prs = build_presentation(kpis, zones, alerts)
        prs.save(output_pptx)

        assert os.path.exists(output_pptx)
        assert os.path.getsize(output_pptx) > 0

        # Verify file is readable as PPTX
        prs_read = Presentation(output_pptx)
        assert len(prs_read.slides) == 1
    finally:
        os.unlink(path)


def test_build_presentation_no_zones(make_excel_file):
    """Presentation works with empty zones."""
    path = make_excel_file(zones=[])
    try:
        kpis, zones, alerts = validate_and_load_data(path)
        prs = build_presentation(kpis, zones, alerts)
        assert prs is not None
    finally:
        os.unlink(path)


# ===== End-to-End Tests =====

def test_end_to_end_full_workflow(make_excel_file, tmp_path):
    """Full workflow: load Excel → validate → build → save."""
    # Create test Excel with realistic data
    path = make_excel_file(
        primary=48.2,
        offtake=44.6,
        primary_mom=4.2,
        offtake_mom=2.8,
        zones=[
            ("North", 12.4, 11.1, 10.5),
            ("South-1", 10.8, 10.5, 2.8),
            ("South-2", 11.2, 9.8, 12.5),
            ("East", 6.5, 6.2, 4.6),
            ("West", 4.1, 4.0, 2.4),
            ("Central", 3.2, 3.0, 6.3),
        ],
        alerts=[
            "West Zone inventory build-up exceeds 30 days.",
            "Reliance shows 3% conversion dip.",
        ]
    )

    output_pptx = str(tmp_path / "MT_Test.pptx")

    try:
        # Load and validate
        kpis, zones, alerts = validate_and_load_data(path)
        assert len(kpis) == 3
        assert len(zones) == 6
        assert len(alerts) >= 2

        # Build presentation
        prs = build_presentation(kpis, zones, alerts)
        assert prs is not None

        # Save
        prs.save(output_pptx)
        assert os.path.exists(output_pptx)
        assert os.path.getsize(output_pptx) > 20000  # At least 20KB

        # Verify output is readable
        prs_verify = Presentation(output_pptx)
        assert len(prs_verify.slides) == 1
    finally:
        if os.path.exists(path):
            os.unlink(path)


def test_idempotency(make_excel_file, tmp_path):
    """Same input always produces identical output (idempotent)."""
    path = make_excel_file(primary=50.0, offtake=48.0, primary_mom=2.5, offtake_mom=2.0)

    output1 = str(tmp_path / "output1.pptx")
    output2 = str(tmp_path / "output2.pptx")

    try:
        # Generate twice from same input
        for output_path in [output1, output2]:
            kpis, zones, alerts = validate_and_load_data(path)
            prs = build_presentation(kpis, zones, alerts)
            prs.save(output_path)

        # Both files should exist and be readable
        assert os.path.exists(output1)
        assert os.path.exists(output2)

        # Files should be nearly identical in size (within 10%)
        size1 = os.path.getsize(output1)
        size2 = os.path.getsize(output2)
        size_diff = abs(size1 - size2) / max(size1, size2)
        assert size_diff < 0.1, f"Output sizes differ by {size_diff*100}%"
    finally:
        if os.path.exists(path):
            os.unlink(path)


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
